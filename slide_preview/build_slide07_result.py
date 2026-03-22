"""
슬라이드 7: RESULT — 대시보드·자동화 체계
HTML → 고해상도 PNG 변환 후 브라우저 프레임에 삽입
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 경로 ──
PREVIEW_DIR = Path(__file__).parent
OUTPUT = PREVIEW_DIR / "slide07_result.pptx"
FREE_HTML = PREVIEW_DIR / "dashboard_html_free.html"
PAID_HTML = PREVIEW_DIR / "dashboard_html_paid.html"
FREE_PNG = PREVIEW_DIR / "dashboard_free_3x.png"
PAID_PNG = PREVIEW_DIR / "dashboard_paid_3x.png"

# ── 색상 ──
CORAL = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CORAL_SUBTLE = RGBColor(0xFF, 0xF8, 0xF9)
TEAL = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT = RGBColor(0xE0, 0xF7, 0xF5)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_100 = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BROWSER_BAR = RGBColor(0x3C, 0x3C, 0x3C)
BROWSER_BTN_R = RGBColor(0xFF, 0x5F, 0x57)
BROWSER_BTN_Y = RGBColor(0xFF, 0xBD, 0x2E)
BROWSER_BTN_G = RGBColor(0x28, 0xCA, 0x41)
BADGE_BLUE = RGBColor(0x38, 0x7A, 0xFF)
BADGE_BLUE_BG = RGBColor(0xEE, 0xF3, 0xFF)

FONT = "Pretendard"


# ══════════════════════════════════════════
# HTML → PNG 렌더링 (Playwright, 3x DPI)
# ══════════════════════════════════════════
def render_html_to_png(html_path: Path, png_path: Path,
                       width: int = 720, height: int = 560):
    """HTML 파일을 3x DPI로 렌더링하여 PNG 저장 (뷰포트 고정)"""
    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": width, "height": height},
            device_scale_factor=3,  # 3x 밀도 → 선명한 텍스트
        )
        page.goto(f"file://{html_path.resolve()}")
        # 폰트 로딩 대기
        page.wait_for_timeout(1500)
        # 뷰포트 크기 고정 캡처 (프레임 비율에 맞춤)
        page.screenshot(path=str(png_path), full_page=False)
        browser.close()
    print(f"  ✓ PNG 렌더링: {png_path.name} ({width}×{height}px × 3x)")


# ══════════════════════════════════════════
# PPTX 헬퍼
# ══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return txBox


def add_rich_text(slide, left, top, width, height, runs,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 14))
        run.font.color.rgb = r.get("color", CHARCOAL)
        run.font.bold = r.get("bold", False)
        run.font.name = r.get("font", FONT)
    return txBox


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size_inch, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size_inch), Inches(size_inch)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_left_bar(slide, left, top, height, color, width_pt=4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Pt(width_pt), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════
# 빌드
# ══════════════════════════════════════════
def build():
    print("=" * 55)
    print("슬라이드 7: RESULT — 대시보드·자동화 체계 (HTML PNG)")
    print("=" * 55)

    # 1) HTML → PNG 렌더링
    print("\n[1/3] HTML → 고해상도 PNG 렌더링...")
    render_html_to_png(FREE_HTML, FREE_PNG, width=720, height=500)
    render_html_to_png(PAID_HTML, PAID_PNG, width=780, height=480)

    # 2) PPTX 빌드
    print("\n[2/3] PPTX 빌드...")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 헤더 ──
    add_text(s, 0.7, 0.35, 10, 0.6,
             "RESULT — 대시보드·자동화 체계", size=43, bold=True)
    add_rect(s, 0.7, 0.95, 11.9, 0.03, CORAL)

    # ── 인사이트 박스 ──
    add_rect(s, 0.7, 1.15, 11.9, 0.55, CORAL_LIGHT, radius=0.04)
    add_rich_text(s, 0.95, 1.25, 11.4, 0.35, [
        {"text": "💡 인사이트:  ", "size": 15, "color": CORAL, "bold": True},
        {"text": "호스트가 ", "size": 15},
        {"text": "통제할 수 있는 변수만", "size": 15, "bold": True},
        {"text": "으로 객실당 수익을 극대화하는 ", "size": 15},
        {"text": "데이터 기반 의사결정 체계", "size": 15, "bold": True},
    ], anchor=MSO_ANCHOR.MIDDLE)

    # ══════════════════════════════════════════
    # 메인 영역: 브라우저 프레임 2개 (무료 + 유료)
    # ══════════════════════════════════════════
    MAIN_Y = 1.9
    FRAME_H = 3.85  # 프레임 높이 축소 (하단 설명 공간 확보)

    # ── 무료 버전 (좌측) ──
    FREE_X = 0.7
    FREE_W = 5.45

    # 레이블 배지
    add_rect(s, FREE_X, MAIN_Y - 0.03, 1.6, 0.28,
             RGBColor(0xE8, 0xF5, 0xE9), radius=0.5)
    add_text(s, FREE_X + 0.08, MAIN_Y - 0.03, 1.44, 0.28,
             "🆓 무료 이메일 진단", size=9, color=RGBColor(0x2E, 0x7D, 0x32),
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 브라우저 프레임 (무료)
    _build_browser_frame(s, FREE_X, MAIN_Y + 0.32, FREE_W, FRAME_H,
                         str(FREE_PNG), "revpar-check.streamlit.app")

    # ── 유료 버전 (우측) ──
    PAID_X = 6.4
    PAID_W = 6.1

    # 레이블 배지
    add_rect(s, PAID_X, MAIN_Y - 0.03, 1.6, 0.28,
             RGBColor(0xFF, 0xF0, 0xF3), radius=0.5)
    add_text(s, PAID_X + 0.08, MAIN_Y - 0.03, 1.44, 0.28,
             "💎 PRO 대시보드", size=9, color=CORAL,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 브라우저 프레임 (유료)
    _build_browser_frame(s, PAID_X, MAIN_Y + 0.32, PAID_W, FRAME_H,
                         str(PAID_PNG), "revpar-dashboard.streamlit.app")

    # 연결 화살표
    arrow_x = FREE_X + FREE_W + 0.12
    arrow_y = MAIN_Y + 0.32 + FRAME_H / 2
    add_text(s, arrow_x, arrow_y, 0.5, 0.3,
             "→", size=28, color=CORAL, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ══════════════════════════════════════════
    # 하단 설명 카드 3개: 문제 · 장점 · 플랫폼 이익
    # ══════════════════════════════════════════
    CARD_Y = MAIN_Y + 0.32 + FRAME_H + 0.2
    CARD_H = 1.15
    CARD_W = 3.73
    CARD_GAP = 0.2

    desc_cards = [
        {
            "num": "문제",
            "title": "호스트의 현실",
            "body": "여러 채널에 분산된 운영 정보 →\n숙소 현황 파악이 어렵고\n수익 개선 포인트를 놓침",
            "accent": CORAL,
            "bg": RGBColor(0xFF, 0xF8, 0xF9),
        },
        {
            "num": "해결",
            "title": "대시보드의 가치",
            "body": "분산된 지표를 단일 화면에 통합\nML 기반 시장 진단 + 맞춤형 체크리스트\n→ 데이터 기반 의사결정 지원",
            "accent": TEAL,
            "bg": RGBColor(0xE0, 0xF7, 0xF5),
        },
        {
            "num": "확장",
            "title": "플랫폼이 얻는 이익",
            "body": "무료 이메일 → 대시보드 유입 →\n유료 전환 → 구독 수익 창출\n호스트 성장 = 플랫폼 수수료 증가",
            "accent": BADGE_BLUE,
            "bg": BADGE_BLUE_BG,
        },
    ]

    for j, card in enumerate(desc_cards):
        cx = 0.7 + j * (CARD_W + CARD_GAP)

        # 카드 배경
        add_rect(s, cx, CARD_Y, CARD_W, CARD_H, card["bg"], radius=0.04)
        # 좌측 악센트 바
        add_left_bar(s, cx, CARD_Y + 0.06, CARD_H - 0.12, card["accent"], width_pt=5)

        # 번호 배지
        badge_w = len(card["num"]) * 0.14 + 0.2
        badge = add_rect(s, cx + 0.2, CARD_Y + 0.1, badge_w, 0.24,
                         card["accent"], radius=0.5)
        add_text(s, cx + 0.2, CARD_Y + 0.1, badge_w, 0.24,
                 card["num"], size=10, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 타이틀
        add_text(s, cx + 0.2 + badge_w + 0.1, CARD_Y + 0.1, CARD_W - badge_w - 0.5, 0.24,
                 card["title"], size=14, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        # 본문
        add_text(s, cx + 0.2, CARD_Y + 0.42, CARD_W - 0.4, 0.7,
                 card["body"], size=11, color=GRAY_500)

    # ── 페이지 번호 ──
    add_text(s, 12.5, 7.05, 0.5, 0.3, "7",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # 3) 저장
    print("\n[3/3] 저장...")
    prs.save(str(OUTPUT))
    print(f"\n✅ 완료: {OUTPUT}")


def _build_browser_frame(slide, x, y, w, h, image_path, url_text):
    """브라우저 프레임 + 이미지 삽입"""
    BAR_H = 0.32

    # 외곽 그림자
    add_rect(slide, x + 0.04, y + 0.04, w, h, GRAY_200, radius=0.03)

    # 프레임 배경
    add_rect(slide, x, y, w, h, WHITE, border_color=GRAY_200, radius=0.03)

    # 브라우저 바
    add_rect(slide, x, y, w, BAR_H, BROWSER_BAR, radius=0.03)
    add_rect(slide, x, y + BAR_H - 0.06, w, 0.06, BROWSER_BAR)

    # 트래픽 라이트
    btn_y = y + 0.1
    add_circle(slide, x + 0.15, btn_y, 0.1, BROWSER_BTN_R)
    add_circle(slide, x + 0.3, btn_y, 0.1, BROWSER_BTN_Y)
    add_circle(slide, x + 0.45, btn_y, 0.1, BROWSER_BTN_G)

    # URL 바
    url_w = min(w - 1.3, 3.5)
    add_rect(slide, x + 0.7, y + 0.06, url_w, 0.18,
             RGBColor(0x55, 0x55, 0x55), radius=0.5)
    add_text(slide, x + 0.85, y + 0.06, url_w - 0.3, 0.18,
             url_text, size=8, color=GRAY_400,
             anchor=MSO_ANCHOR.MIDDLE)

    # LIVE 배지
    add_rect(slide, x + w - 0.7, y + 0.06, 0.45, 0.18,
             TEAL, radius=0.5)
    add_text(slide, x + w - 0.7, y + 0.06, 0.45, 0.18,
             "● LIVE", size=7, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 이미지 삽입 (콘텐츠 영역)
    img_x = x + 0.05
    img_y = y + BAR_H + 0.05
    img_w = w - 0.1
    img_h = h - BAR_H - 0.1

    try:
        from PIL import Image
        with Image.open(image_path) as img:
            orig_w, orig_h = img.size
            # 이미지 비율 유지
            img_ratio = orig_w / orig_h
            frame_ratio = img_w / img_h
            if img_ratio > frame_ratio:
                # 이미지가 더 넓음 → 너비 맞춤
                final_w = img_w
                final_h = img_w / img_ratio
            else:
                # 이미지가 더 높음 → 높이 맞춤
                final_h = img_h
                final_w = img_h * img_ratio

            # 센터 정렬
            offset_x = (img_w - final_w) / 2
            offset_y = 0  # 상단 정렬

            slide.shapes.add_picture(
                image_path,
                Inches(img_x + offset_x),
                Inches(img_y + offset_y),
                Inches(final_w),
                Inches(final_h),
            )
            print(f"  ✓ 이미지 삽입: {Path(image_path).name} → {final_w:.1f}×{final_h:.1f}in")
    except Exception as e:
        print(f"  ⚠ 이미지 삽입 실패: {e}")
        # 플레이스홀더 표시
        add_rect(slide, img_x, img_y, img_w, img_h, GRAY_100, radius=0.02)
        add_text(slide, img_x, img_y + img_h / 2 - 0.2, img_w, 0.4,
                 "📸 스크린샷 위치", size=16, color=GRAY_400,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    build()
