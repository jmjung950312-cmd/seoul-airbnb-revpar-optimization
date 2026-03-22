"""
EDA 슬라이드 4장 — 리디자인 v2
왼쪽 히어로 텍스트 + 오른쪽 차트 (H1 패턴 통일)
"""
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 경로 ──
BASE = Path(__file__).parent.parent
DATA_PATH = BASE / "data" / "processed" / "seoul_airbnb_features.csv"
CHART_DIR = Path(__file__).parent / "charts"
OUTPUT = Path(__file__).parent / "eda_slides_native.pptx"
CHART_DIR.mkdir(exist_ok=True)

# ── 색상 ──
CORAL = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Pretendard"

# matplotlib
plt.rcParams['font.family'] = 'AppleGothic'
plt.rcParams['axes.unicode_minus'] = False

C_CORAL = '#FF385C'
C_CHARCOAL = '#484848'
C_GRAY500 = '#767676'
C_GRAY300 = '#DBDBDB'
C_BAR_GRAY = '#D9D9D9'
C_BLUE = '#4A90D9'
C_HIGHLIGHT = '#FF385C'

TOURIST_DISTRICTS = ['Jung-gu', 'Jongno-gu', 'Yongsan-gu', 'Gangnam-gu', 'Mapo-gu']


# ═══════════════════════════════════════════
# PPTX 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return txBox


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRAY_200, width=1, dash=False):
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dash:
        conn.line.dash_style = 4
    return conn


def add_header(slide, title):
    """헤더: 제목 43pt + coral 밑줄 (breadcrumb 제거)"""
    add_text(slide, 0.7, 0.35, 10, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    """인사이트 박스: 핑크 배경 + 20.5pt 텍스트"""
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL, bold=False)


def add_stat_card(slide, left, top, width, height, value, desc,
                  accent=False, value_color=CHARCOAL):
    """보조 수치 카드"""
    bg = CORAL_LIGHT if accent else GRAY_BG
    add_rect(slide, left, top, width, height, bg, radius=0.04)
    if accent:
        add_rect(slide, left, top, 0.03, height, CORAL)
    add_text(slide, left + 0.18, top + 0.12, 1.5, 0.35, value,
             size=22, color=value_color, bold=True)
    add_text(slide, left + 1.75, top + 0.1, width - 2.0, height - 0.2, desc,
             size=13, color=GRAY_500)


# ═══════════════════════════════════════════
# 데이터
# ═══════════════════════════════════════════
def load_data():
    df = pd.read_csv(DATA_PATH)
    ao = df[(df['refined_status'] == 'Active') &
            (df['operation_status'] == 'Operating')].copy()
    print(f"Active+Operating: {len(ao)}건")
    return ao


# ═══════════════════════════════════════════
# 차트 1: H3 사진 — 오른쪽 배치용 (세로형)
# ═══════════════════════════════════════════
def chart_h3_photos(ao):
    bins = list(range(0, 101, 10))
    labels = [f"{bins[i]+1}-{bins[i+1]}" for i in range(len(bins)-1)]
    labels[0] = "1-10"

    ao_f = ao[ao['photos_count'].between(1, 100)].copy()
    ao_f['photo_bin'] = pd.cut(ao_f['photos_count'], bins=bins, labels=labels,
                               include_lowest=True)
    grp = ao_f.groupby('photo_bin', observed=False)
    medians = grp['ttm_revpar'].median()
    overall_median = ao['ttm_revpar'].median()

    fig, ax = plt.subplots(figsize=(8, 6.5), dpi=200)
    x = np.arange(len(labels))

    # 최적 구간 하이라이트 (21-40)
    ax.axvspan(1.5, 3.5, color='#FFE0E6', alpha=0.5, zorder=0,
               label='효율 최적 구간 (21~40장)')

    # 바
    bar_colors = [C_CORAL if 2 <= i <= 3 else C_BAR_GRAY for i in range(len(labels))]
    ax.bar(x, medians.values, color=bar_colors, width=0.6,
           edgecolor='white', linewidth=0.5, zorder=2)

    # 트렌드 라인
    ax.plot(x, medians.values, color='#FF385C', marker='o', markersize=4,
            linewidth=2, zorder=3, alpha=0.6)

    # 서울 중위값
    ax.axhline(y=overall_median, color=C_GRAY500, linestyle='--',
               linewidth=1.2, alpha=0.6, zorder=1)
    ax.text(len(labels) - 0.3, overall_median + 2000,
            f"중위값 ₩{overall_median/1000:.0f}K",
            ha='right', fontsize=9, color=C_GRAY500)

    # 값 라벨 (바 위에, 클린하게)
    for i, v in enumerate(medians.values):
        ax.text(i, v + 1500, f"₩{v/1000:.0f}K",
                ha='center', va='bottom', fontsize=8, color=C_CHARCOAL,
                fontweight='bold' if 2 <= i <= 3 else 'normal')

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, rotation=0)
    ax.set_xlabel("사진 수 구간", fontsize=10, color=C_CHARCOAL, labelpad=8)
    ax.set_ylabel("")
    ax.set_ylim(0, medians.values.max() * 1.18)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.tick_params(left=False, labelleft=False)
    ax.grid(axis='y', alpha=0.15, linestyle='-', color=C_GRAY300)

    plt.tight_layout()
    path = CHART_DIR / "h3_photos_revpar.png"
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ {path.name}")
    return path, medians


# ═══════════════════════════════════════════
# 차트 2: H4 최소숙박 — 2-panel
# ═══════════════════════════════════════════
def chart_h4_min_nights(ao):
    ao_f = ao.copy()

    def mn_bin(v):
        if v <= 1: return '1박'
        elif v <= 3: return '2~3박'
        elif v <= 7: return '4~7박'
        else: return '8박+'

    ao_f['mn_group'] = ao_f['min_nights'].apply(mn_bin)
    order = ['1박', '2~3박', '4~7박', '8박+']
    ao_f['mn_group'] = pd.Categorical(ao_f['mn_group'], categories=order, ordered=True)

    grp = ao_f.groupby('mn_group', observed=False)
    rev_med = grp['ttm_revpar'].median()
    occ_med = grp['ttm_occupancy'].median()

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8, 5.5), dpi=200)

    # ── 왼쪽: RevPAR ──
    colors_rev = [C_BAR_GRAY if i != 1 else C_CORAL for i in range(4)]
    ax1.bar(order, rev_med.values, color=colors_rev, width=0.5,
            edgecolor='white', linewidth=0.5)

    for i, v in enumerate(rev_med.values):
        ax1.text(i, v + 1200, f"₩{v/1000:.0f}K",
                 ha='center', va='bottom', fontsize=10,
                 color=C_CORAL if i == 1 else C_GRAY500,
                 fontweight='bold' if i == 1 else 'normal')

    ax1.set_title("RevPAR (수익)", fontsize=12, fontweight='bold',
                  color=C_CHARCOAL, pad=15)
    ax1.set_ylim(0, rev_med.values.max() * 1.25)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.spines['left'].set_visible(False)
    ax1.tick_params(left=False, labelleft=False)
    ax1.grid(axis='y', alpha=0.15, linestyle='-')

    # ── 오른쪽: 점유율 ──
    colors_occ = [C_BAR_GRAY if i != 2 else C_BLUE for i in range(4)]
    ax2.bar(order, occ_med.values * 100, color=colors_occ, width=0.5,
            edgecolor='white', linewidth=0.5)

    for i, v in enumerate(occ_med.values * 100):
        ax2.text(i, v + 1, f"{v:.0f}%",
                 ha='center', va='bottom', fontsize=10,
                 color=C_BLUE if i == 2 else C_GRAY500,
                 fontweight='bold' if i == 2 else 'normal')

    ax2.set_title("Occupancy (예약률)", fontsize=12, fontweight='bold',
                  color=C_CHARCOAL, pad=15)
    ax2.set_ylim(0, occ_med.values.max() * 100 * 1.25)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.spines['left'].set_visible(False)
    ax2.tick_params(left=False, labelleft=False)
    ax2.grid(axis='y', alpha=0.15, linestyle='-')

    plt.tight_layout(w_pad=2)
    path = CHART_DIR / "h4_min_nights.png"
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ {path.name}")
    return path, rev_med, occ_med


# ═══════════════════════════════════════════
# 차트 3: H5 추가요금 — 스택 바
# ═══════════════════════════════════════════
def chart_h5_extra_fee(ao):
    ao_f = ao.copy()
    ao_f['is_tourist'] = ao_f['district'].isin(TOURIST_DISTRICTS)

    segments = {
        '소형\n(1-6인)': ao_f[ao_f['guests_group'] == '1_6'],
        '대형\n(7인+)': ao_f[ao_f['guests_group'] == '7+'],
        '관광\n핵심자치구': ao_f[ao_f['is_tourist']],
        '주거\n중심자치구': ao_f[~ao_f['is_tourist']],
    }

    no_fee, with_fee, labels = [], [], []
    for name, sub in segments.items():
        nf = sub[sub['extra_guest_fee_policy'] == 0]['ttm_revpar'].median()
        wf = sub[sub['extra_guest_fee_policy'] == 1]['ttm_revpar'].median()
        no_fee.append(nf)
        with_fee.append(wf)
        labels.append(name)

    no_fee = np.array(no_fee)
    with_fee = np.array(with_fee)
    diffs = no_fee - with_fee
    pcts = diffs / no_fee * 100

    fig, ax = plt.subplots(figsize=(8, 6), dpi=200)
    x = np.arange(len(labels))

    ax.bar(x, with_fee / 1000, color=C_BAR_GRAY, width=0.45,
           edgecolor='white', linewidth=0.5, label='추가요금 부과 RevPAR')
    ax.bar(x, diffs / 1000, bottom=with_fee / 1000, color=C_CORAL,
           width=0.45, label='손실분', alpha=0.85, edgecolor='white', linewidth=0.5)

    for i in range(len(labels)):
        total = no_fee[i] / 1000
        d = diffs[i] / 1000
        pct = pcts[i]
        ax.text(i, total + 1.5, f"▼₩{d:.0f}K\n(-{pct:.0f}%)",
                ha='center', va='bottom', fontsize=10, fontweight='bold',
                color=C_CORAL, linespacing=1.3)
        # 베이스 값
        ax.text(i, with_fee[i] / 1000 / 2, f"₩{with_fee[i]/1000:.0f}K",
                ha='center', va='center', fontsize=9, color=C_GRAY500)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=10)
    ax.set_ylim(0, no_fee.max() / 1000 * 1.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.tick_params(left=False, labelleft=False)
    ax.grid(axis='y', alpha=0.15, linestyle='-')
    ax.legend(loc='upper right', framealpha=0.9, fontsize=9,
              edgecolor=C_GRAY300)

    plt.tight_layout()
    path = CHART_DIR / "h5_extra_fee.png"
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ {path.name}")
    return path, diffs, pcts


# ═══════════════════════════════════════════
# 차트 4: 히트맵
# ═══════════════════════════════════════════
def chart_h3_heatmap(ao):
    ao_f = ao[(ao['photos_count'] > 0) & (ao['rating_overall'] > 0)].copy()

    rating_bins = [0, 4.5, 4.70, 4.85, 4.95, 5.01]
    rating_labels = ['~4.5', '4.50-4.70', '4.70-4.85', '4.85-4.95', '4.95-5.0']
    photo_bins = [0, 10, 20, 35, 60, 999]
    photo_labels = ['0-10장', '11-20장', '21-35장', '36-60장', '60+장']

    ao_f['rating_grp'] = pd.cut(ao_f['rating_overall'], bins=rating_bins,
                                labels=rating_labels, include_lowest=True)
    ao_f['photo_grp'] = pd.cut(ao_f['photos_count'], bins=photo_bins,
                               labels=photo_labels, include_lowest=True)
    pivot = ao_f.groupby(['rating_grp', 'photo_grp'],
                         observed=False)['ttm_revpar'].median().unstack()

    fig, ax = plt.subplots(figsize=(8, 6), dpi=200)

    cmap = LinearSegmentedColormap.from_list('custom',
        ['#FFF8E1', '#FFD54F', '#FF8A65', '#5C6BC0', '#303F9F'], N=256)

    data = pivot.values
    im = ax.imshow(data, cmap=cmap, aspect='auto')

    for i in range(data.shape[0]):
        for j in range(data.shape[1]):
            val = data[i, j]
            if np.isnan(val):
                continue
            norm_val = (val - np.nanmin(data)) / (np.nanmax(data) - np.nanmin(data) + 1e-6)
            clr = 'white' if norm_val > 0.55 else C_CHARCOAL
            ax.text(j, i, f"₩{val:,.0f}",
                    ha='center', va='center', fontsize=9.5, fontweight='bold', color=clr)

    # 피크 행 강조
    ax.add_patch(plt.Rectangle((-0.5, 3 - 0.5), data.shape[1], 1,
                               fill=False, edgecolor=C_CORAL, linewidth=3))

    ax.set_xticks(range(len(photo_labels)))
    ax.set_xticklabels(photo_labels, fontsize=10)
    ax.set_yticks(range(len(rating_labels)))
    ax.set_yticklabels(rating_labels, fontsize=10)
    ax.set_xlabel("사진 수 구간", fontsize=11, color=C_CHARCOAL, labelpad=10)
    ax.set_ylabel("평점 구간", fontsize=11, color=C_CHARCOAL, labelpad=10)

    # 하단 주석
    for j, (txt, clr_) in enumerate([
        ("← 준비 부족", C_GRAY500), ("", ""), ("최적 효율", C_CORAL),
        ("", ""), ("효율 체감 →", C_GRAY500)]):
        if txt:
            ax.text(j, len(rating_labels) + 0.15, txt,
                    ha='center', va='top', fontsize=9, color=clr_,
                    fontweight='bold' if clr_ == C_CORAL else 'normal')

    cbar = plt.colorbar(im, ax=ax, shrink=0.8, pad=0.02)
    cbar.ax.yaxis.set_major_formatter(mticker.FuncFormatter(
        lambda v, _: f"₩{v/1000:.0f}K"))
    cbar.set_label("")

    plt.tight_layout()
    path = CHART_DIR / "h3_heatmap.png"
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ {path.name}")
    return path, pivot


# ═══════════════════════════════════════════
# PPTX 빌드 — 왼쪽 히어로 + 오른쪽 차트
# ═══════════════════════════════════════════
def build_pptx(charts, data):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    photo_medians = data['photo_medians']
    rev_med = data['rev_med']
    occ_med = data['occ_med']
    h5_diffs, h5_pcts = data['h5']
    heatmap_pivot = data['heatmap']

    # 공통 레이아웃 상수
    HERO_W = 4.8   # 왼쪽 히어로 폭
    DIV_X = 5.2    # 세로 구분선 X
    CHART_X = 5.5  # 차트 시작 X
    CHART_W = 7.3  # 차트 폭

    # ════════════════════════════════════════
    # 슬라이드 1: H3 사진 수
    # ════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "H3 — 사진 수와 수익의 관계")
    add_insight_box(s, "인사이트 :  사진 수 21~35장 구간에서 RevPAR 효율이 가장 높음. 그 이후 추가 효과는 절반으로 감소.")

    # 세로 구분선
    add_line(s, DIV_X, 2.1, DIV_X, 6.5, GRAY_200, 1)

    # ── 왼쪽 히어로 ──
    L = 0.7
    add_text(s, L, 2.5, 4, 0.9, "21~35장",
             size=83, color=CORAL, bold=True)
    add_text(s, L, 3.45, 4, 0.3, "RevPAR 효율 최적 구간",
             size=20, color=CHARCOAL, bold=True)
    add_text(s, L, 3.85, 4.2, 0.7,
             "사진 30장 추가 시 수익 +132%\n그 이후 60장을 더 찍어도 추가 효과는 절반으로 감소",
             size=14, color=GRAY_500)

    # 보조 카드
    opt_val = photo_medians.values[2]  # 21-30
    low_val = photo_medians.values[0]  # 1-10
    gain_pct = (opt_val - low_val) / low_val * 100

    add_stat_card(s, L, 4.8, 4.2, 0.65,
                  f"+{gain_pct:.0f}%",
                  f"1-10장 → 21-30장\n수익 증가율",
                  accent=True, value_color=CORAL)

    peak_val = photo_medians.values.max()
    add_stat_card(s, L, 5.6, 4.2, 0.65,
                  f"₩{peak_val:,.0f}",
                  f"91-100장 피크값\n(n=122, 표본 소수)",
                  accent=False)

    add_text(s, L, 6.45, 4.2, 0.4,
             "→ 21~35장이 비용 대비 효과가 가장 높은 구간",
             size=13, color=CORAL, bold=True)

    # ── 오른쪽 차트 ──
    s.shapes.add_picture(str(charts['h3_photos']),
                         Inches(CHART_X), Inches(2.0), Inches(CHART_W), Inches(5.0))

    add_text(s, 12.5, 7.0, 0.5, 0.3, "14",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════
    # 슬라이드 2: H4 최소숙박
    # ════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "H4 — 최소숙박일, 수익과 예약률의 딜레마")
    add_insight_box(s, "인사이트 :  짧으면 예약↑ 수익↓, 길면 수익↑ 예약↓ — 두 지표의 균형점이 RevPAR을 결정한다.")

    add_line(s, DIV_X, 2.1, DIV_X, 6.5, GRAY_200, 1)

    # ── 왼쪽 히어로 ──
    L = 0.7
    add_text(s, L, 2.5, 4.2, 0.9, "2~3박",
             size=83, color=CORAL, bold=True)
    add_text(s, L, 3.45, 4, 0.3, "수익 최적 구간",
             size=20, color=CHARCOAL, bold=True)

    peak_rev = rev_med.values[1]
    peak_occ_idx = occ_med.values.argmax()
    peak_occ = occ_med.values[peak_occ_idx] * 100
    peak_occ_label = ['1박', '2~3박', '4~7박', '8박+'][peak_occ_idx]

    add_text(s, L, 3.85, 4.2, 0.7,
             f"수익은 2~3박에서 최대 (₩{peak_rev/1000:.0f}K)\n"
             f"예약률은 {peak_occ_label}에서 최대 ({peak_occ:.0f}%)",
             size=14, color=GRAY_500)

    add_stat_card(s, L, 4.8, 4.2, 0.65,
                  f"₩{peak_rev/1000:.0f}K",
                  "2~3박 RevPAR 중위값\n수익 기준 최고점",
                  accent=True, value_color=CORAL)

    BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)
    BLUE_RGB = RGBColor(0x4A, 0x90, 0xD9)
    add_rect(s, L, 5.6, 4.2, 0.65, BLUE_LIGHT, radius=0.04)
    add_rect(s, L, 5.6, 0.03, 0.65, BLUE_RGB)
    add_text(s, L + 0.18, 5.72, 1.5, 0.35, f"{peak_occ:.0f}%",
             size=22, color=BLUE_RGB, bold=True)
    add_text(s, L + 1.75, 5.7, 2.3, 0.45,
             f"{peak_occ_label} 점유율 중위값\n예약률 기준 최고점",
             size=13, color=GRAY_500)

    add_text(s, L, 6.45, 4.2, 0.4,
             "→ 2~3박 설정이 수익·예약 균형의 최적점",
             size=13, color=CORAL, bold=True)

    # ── 오른쪽 차트 ──
    s.shapes.add_picture(str(charts['h4']),
                         Inches(CHART_X), Inches(2.0), Inches(CHART_W), Inches(5.0))

    add_text(s, 12.5, 7.0, 0.5, 0.3, "15",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════
    # 슬라이드 3: H5 추가요금
    # ════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "H5 — 추가요금, 수익을 깎는 숨은 변수")

    max_idx = h5_diffs.argmax()
    max_drop_pct = h5_pcts[max_idx]

    add_insight_box(s, f"인사이트 :  추가요금을 부과하는 숙소는 부과하지 않는 숙소 대비 수익이 최대 {max_drop_pct:.0f}% 낮음. 특히 대형 숙소에서 손실 극대화.")

    add_line(s, DIV_X, 2.1, DIV_X, 6.5, GRAY_200, 1)

    # ── 왼쪽 히어로 ──
    L = 0.7

    add_text(s, L, 2.5, 4, 0.9, f"-{max_drop_pct:.0f}%",
             size=83, color=CORAL, bold=True)
    add_text(s, L, 3.45, 4, 0.3, "대형 숙소 RevPAR 손실",
             size=20, color=CHARCOAL, bold=True)
    add_text(s, L, 3.85, 4.2, 0.7,
             "추가요금을 부과하는 숙소는 부과하지 않는\n"
             "숙소 대비 수익이 최대 56% 낮음\n"
             "특히 대형 숙소(7인+)에서 손실이 극대화",
             size=14, color=GRAY_500)

    add_stat_card(s, L, 4.8, 4.2, 0.65,
                  f"₩{h5_diffs[max_idx]/1000:.0f}K",
                  "대형 숙소 객실당\n수익 손실 금액",
                  accent=True, value_color=CORAL)

    add_stat_card(s, L, 5.6, 4.2, 0.65,
                  f"-{h5_pcts[0]:.0f}%",
                  f"소형 숙소도 -{h5_pcts[0]:.0f}% 손실\n추가요금 전략 재검토 필요",
                  accent=False)

    add_text(s, L, 6.45, 4.2, 0.4,
             "→ 대형+관광지: 추가요금 즉시 제거 권장",
             size=13, color=CORAL, bold=True)

    # ── 오른쪽 차트 ──
    s.shapes.add_picture(str(charts['h5']),
                         Inches(CHART_X), Inches(2.0), Inches(CHART_W), Inches(5.0))

    add_text(s, 12.5, 7.0, 0.5, 0.3, "16",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════
    # 슬라이드 4: 사진·평점 히트맵
    # ════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "The Paradox of Perfection")
    add_insight_box(s, "인사이트 :  완벽한 5.0점이 아닌 4.85~4.95점 구간에서 RevPAR이 가장 높음. 만점 추구가 오히려 수익을 낮춘다.")

    add_line(s, DIV_X, 2.1, DIV_X, 6.5, GRAY_200, 1)

    # ── 왼쪽 히어로 ──
    L = 0.7
    add_text(s, L, 2.5, 4, 0.9, "4.85~4.95",
             size=70, color=CORAL, bold=True)
    add_text(s, L, 3.45, 4, 0.3, "최적 평점 구간",
             size=20, color=CHARCOAL, bold=True)
    add_text(s, L, 3.85, 4.2, 0.7,
             "완벽한 5.0점이 아닌 4.85~4.95점 구간에서\n"
             "RevPAR 중위값이 가장 높음\n"
             "최상의 성과는 '신뢰할 수 있는 구간'에서 나옴",
             size=14, color=GRAY_500)

    peak_row = heatmap_pivot.loc['4.85-4.95']
    peak_val = peak_row.max()

    if '4.95-5.0' in heatmap_pivot.index:
        val_50 = heatmap_pivot.loc['4.95-5.0'].iloc[2]  # 21-35장 기준
        drop = (1 - val_50 / heatmap_pivot.loc['4.85-4.95'].iloc[2]) * 100
    else:
        val_50, drop = 0, 0

    add_stat_card(s, L, 4.8, 4.2, 0.65,
                  f"₩{peak_val:,.0f}",
                  "4.85~4.95 × 60+장\n최고 RevPAR 조합",
                  accent=True, value_color=CORAL)
    add_stat_card(s, L, 5.6, 4.2, 0.65,
                  f"-{drop:.0f}%",
                  f"만점 5.0의 함정\n₩{val_50:,.0f} (같은 사진 구간)")

    add_text(s, L, 6.45, 4.2, 0.4,
             "→ 만점 추구보다 4.85~4.95 유지가 수익에 유리",
             size=13, color=CORAL, bold=True)

    # ── 오른쪽 차트 ──
    s.shapes.add_picture(str(charts['heatmap']),
                         Inches(CHART_X), Inches(2.0), Inches(CHART_W), Inches(5.0))

    add_text(s, 12.5, 7.0, 0.5, 0.3, "17",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # 저장
    prs.save(str(OUTPUT))
    print(f"\n✓ PPTX 저장: {OUTPUT}")
    return OUTPUT


# ═══════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════
def main():
    print("=" * 50)
    print("EDA 슬라이드 v2 빌드")
    print("=" * 50)

    ao = load_data()

    print("\n차트 생성:")
    p1, photo_med = chart_h3_photos(ao)
    p2, rev_med, occ_med = chart_h4_min_nights(ao)
    p3, h5_diffs, h5_pcts = chart_h5_extra_fee(ao)
    p4, heatmap_pivot = chart_h3_heatmap(ao)

    charts = {
        'h3_photos': p1, 'h4': p2,
        'h5': p3, 'heatmap': p4,
    }
    data = {
        'photo_medians': photo_med, 'rev_med': rev_med,
        'occ_med': occ_med, 'h5': (h5_diffs, h5_pcts),
        'heatmap': heatmap_pivot,
    }

    print("\nPPTX 빌드:")
    build_pptx(charts, data)
    print("완료!")


if __name__ == "__main__":
    main()
