"""
슬라이드 12 — 가설 리스트 (2가지 버전)
Version A: 분석 기반 채택 가설 7개 (호스트 통제 vs 시장 요인 분류)
Version B: 사용자 공유 리스트 (7개 가설 테이블)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 경로 ──
BASE = Path(__file__).parent.parent
OUTPUT_A = Path(__file__).parent / "slide12_hypotheses_A.pptx"
OUTPUT_B = Path(__file__).parent / "slide12_hypotheses_B.pptx"

# ── 색상 ──
CORAL = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xA6, 0x99)
BLUE = RGBColor(0x4A, 0x90, 0xD9)

FONT = "Pretendard"


# ═══════════════════════════════════════════
# 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, line_spacing=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRAY_200, width=1):
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_header(slide, title):
    """헤더: 제목 43pt + coral 밑줄"""
    add_text(slide, 0.7, 0.35, 10, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    """인사이트 박스: 핑크 배경"""
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL, bold=False)


def add_badge(slide, left, top, num, color=CORAL):
    """번호 뱃지 (원형)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.32), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def new_prs():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs


# ═══════════════════════════════════════════
# Version A: 분석 기반 채택 가설 7개 (좌우 분할)
# ═══════════════════════════════════════════
def build_version_a():
    prs = new_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])

    add_header(s, "수익에 영향을 미치는 주요 가설 7가지")
    add_insight_box(s, "인사이트 :  9개 가설 중 7개 채택 — 호스트가 통제 가능한 5개 + 시장 환경 2개로 분류")

    # ── 좌우 분류 소제목 ──
    # 왼쪽: 호스트 통제 가능
    host_x = 0.7
    market_x = 7.2

    add_rect(s, host_x, 2.1, 5.8, 0.45, CORAL_LIGHT, radius=0.04)
    add_text(s, host_x + 0.2, 2.15, 5.4, 0.35,
             "호스트가 통제 가능한 변수", size=18, color=CORAL, bold=True)

    add_rect(s, market_x, 2.1, 5.4, 0.45, RGBColor(0xE8, 0xF0, 0xFE), radius=0.04)
    add_text(s, market_x + 0.2, 2.15, 5.0, 0.35,
             "시장 · 환경 변수", size=18, color=BLUE, bold=True)

    # ── 왼쪽 가설 5개 ──
    host_hypotheses = [
        ("H1", "슈퍼호스트 프리미엄", "슈퍼호스트는 비슈퍼호스트 대비 +83.1% 높은 수익", "채택"),
        ("H3", "사진의 힘", "사진 21~35장 구간에서 RevPAR 효율이 가장 높음", "채택"),
        ("H4", "최소숙박 최적점", "최소 2~3박 설정 시 수익·예약 균형 최적", "채택"),
        ("H5", "추가요금 역효과", "7인+ 대형 숙소에서 추가요금 시 수익 -55.9%", "채택"),
        ("+", "완벽주의의 역설", "만점 5.0이 아닌 4.85~4.95 구간이 수익 최고", "추가 발견"),
    ]

    y_start = 2.8
    card_h = 0.82
    card_gap = 0.08

    for i, (hid, title, desc, result) in enumerate(host_hypotheses):
        y = y_start + i * (card_h + card_gap)

        # 카드 배경
        add_rect(s, host_x, y, 5.8, card_h, WHITE, border_color=GRAY_200, radius=0.03)
        # 좌측 coral 악센트 라인
        add_rect(s, host_x, y, 0.04, card_h, CORAL)

        # 번호 뱃지
        if hid == "+":
            add_badge(s, host_x + 0.2, y + 0.1, "+", GREEN)
        else:
            add_badge(s, host_x + 0.2, y + 0.1, hid[1:])

        # 제목
        add_text(s, host_x + 0.65, y + 0.07, 4.8, 0.3,
                 title, size=16, color=CHARCOAL, bold=True)

        # 설명
        add_text(s, host_x + 0.65, y + 0.38, 4.8, 0.35,
                 desc, size=12, color=GRAY_500)

        # 결과 뱃지
        badge_color = GREEN if result == "추가 발견" else CORAL
        badge_w = 1.0 if result == "추가 발견" else 0.6
        add_rect(s, host_x + 5.8 - badge_w - 0.15, y + 0.12,
                 badge_w, 0.28, badge_color, radius=0.04)
        add_text(s, host_x + 5.8 - badge_w - 0.15, y + 0.14,
                 badge_w, 0.24, result,
                 size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # ── 오른쪽 가설 2개 ──
    market_hypotheses = [
        ("H2", "객실 유형별 RevPAR", "entire_home이 private_room 대비 2.7배 높은 RevPAR", "채택"),
        ("H6", "공급 집중 → RevPAR 압박", "마포구(21% 공급) 등 공급 과잉 자치구에서\nRevPAR 하락 압력 확인", "채택"),
    ]

    for i, (hid, title, desc, result) in enumerate(market_hypotheses):
        y = y_start + i * (card_h + card_gap)

        add_rect(s, market_x, y, 5.4, card_h, WHITE, border_color=GRAY_200, radius=0.03)
        add_rect(s, market_x, y, 0.04, card_h, BLUE)
        add_badge(s, market_x + 0.2, y + 0.1, hid[1:], BLUE)
        add_text(s, market_x + 0.65, y + 0.07, 4.4, 0.3,
                 title, size=16, color=CHARCOAL, bold=True)
        add_text(s, market_x + 0.65, y + 0.38, 4.4, 0.35,
                 desc, size=12, color=GRAY_500)
        add_rect(s, market_x + 5.4 - 0.75, y + 0.12, 0.6, 0.28, BLUE, radius=0.04)
        add_text(s, market_x + 5.4 - 0.75, y + 0.14, 0.6, 0.24, result,
                 size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # ── 오른쪽 하단: 제외 가설 안내 ──
    excl_y = y_start + 2 * (card_h + card_gap) + 0.15
    add_rect(s, market_x, excl_y, 5.4, 1.6, GRAY_BG, radius=0.04)
    add_text(s, market_x + 0.3, excl_y + 0.12, 5.0, 0.3,
             "분석 대상에서 제외", size=14, color=GRAY_400, bold=True)
    excluded = [
        "H7  POI 근접성 — 부분 채택 (유형 > 거리)",
        "H8  자치구 인구 — 약한 증거 (rho = -0.037)",
        "H9  Dormant 비율 — 미검증 (현황 파악만)",
    ]
    for j, text in enumerate(excluded):
        add_text(s, market_x + 0.3, excl_y + 0.5 + j * 0.32, 4.8, 0.3,
                 text, size=12, color=GRAY_400)

    # ── 하단 CTA ──
    add_rect(s, 0.7, 6.75, 11.8, 0.45, CORAL_LIGHT, radius=0.04)
    add_text(s, 1.0, 6.82, 11.2, 0.3,
             "이 중 호스트가 오늘 당장 바꿀 수 있는 5가지를 깊이 봅니다",
             size=18, color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # 페이지 번호
    add_text(s, 12.5, 7.0, 0.5, 0.3, "12",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    prs.save(str(OUTPUT_A))
    print(f"Version A 저장: {OUTPUT_A}")


# ═══════════════════════════════════════════
# Version B: 나열 + 디자인 폴리시
# ═══════════════════════════════════════════
def build_version_b():
    prs = new_prs()
    s = prs.slides.add_slide(prs.slide_layouts[6])

    add_header(s, "수익에 영향을 미치는 주요 가설 7가지")
    add_insight_box(s, "인사이트 :  호스트가 직접 통제할 수 있는 변수들이 RevPAR에 어떤 영향을 미치는지 7가지 관점에서 검증합니다.")

    # ── 가설 리스트 ──
    hypotheses = [
        ("1", "슈퍼호스트는 더 높은 수익을 만들 것이다", "슈퍼호스트 지위"),
        ("2", "숙소 사진 수는 예약 성과에 영향을 줄 것이다", "사진 수"),
        ("3", "최소 숙박일 설정은 수익에 영향을 줄 것이다", "최소숙박"),
        ("4", "추가 게스트 요금 정책은 수익에 영향을 줄 것이다", "추가요금"),
        ("5", "숙소 평점은 예약 성과에 영향을 줄 것이다", "평점"),
        ("6", "완벽한 평점(5.0)은 항상 최적이 아닐 것이다", "만점의 역설"),
        ("7", "시장에는 비활성 숙소가 많을 것이다", "시장 건전성"),
    ]

    card_x = 0.7
    card_w = 11.8
    card_h = 0.58
    card_gap = 0.06
    y_start = 2.1

    for i, (num, text, keyword) in enumerate(hypotheses):
        y = y_start + i * (card_h + card_gap)

        # 카드 배경
        add_rect(s, card_x, y, card_w, card_h, WHITE, border_color=GRAY_200, radius=0.03)

        # 좌측 coral 악센트 바
        add_rect(s, card_x, y, 0.04, card_h, CORAL)

        # 번호 뱃지
        add_badge(s, card_x + 0.25, y + 0.13, num)

        # 키워드 태그 (coral 배경)
        kw_w = len(keyword) * 0.18 + 0.3
        add_rect(s, card_x + 0.7, y + 0.13, kw_w, 0.32, CORAL_LIGHT, radius=0.04)
        add_text(s, card_x + 0.7 + 0.1, y + 0.15, kw_w - 0.2, 0.26,
                 keyword, size=12, color=CORAL, bold=True)

        # 가설 텍스트
        text_x = card_x + 0.7 + kw_w + 0.2
        add_text(s, text_x, y + 0.14, card_w - (text_x - card_x) - 0.3, 0.3,
                 text, size=15, color=CHARCOAL)

    # ── 하단 안내 ──
    bottom_y = y_start + 7 * (card_h + card_gap) + 0.1
    add_line(s, card_x, bottom_y, card_x + card_w, bottom_y, GRAY_200, 1)
    add_text(s, card_x, bottom_y + 0.12, card_w, 0.3,
             "이 중 호스트가 오늘 바꿀 수 있는 5가지를 다음 장에서 깊이 봅니다  →",
             size=14, color=GRAY_500, align=PP_ALIGN.RIGHT)

    # 페이지 번호
    add_text(s, 12.5, 7.0, 0.5, 0.3, "12",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    prs.save(str(OUTPUT_B))
    print(f"Version B 저장: {OUTPUT_B}")


# ═══════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════
def main():
    print("=" * 50)
    print("슬라이드 12 — 가설 리스트 빌드 (2버전)")
    print("=" * 50)

    print("\n[Version A] 분석 기반 채택 가설 7개 (좌우 분할)")
    build_version_a()

    print("\n[Version B] 사용자 공유 리스트 (심플 테이블)")
    build_version_b()

    print("\n완료!")


if __name__ == "__main__":
    main()
