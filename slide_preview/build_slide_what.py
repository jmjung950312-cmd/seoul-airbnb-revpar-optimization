"""
슬라이드 5: WHAT — 객실당 수익의 구조
Layout: 히어로 센터 + 극적 대비 (차트 없음)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 경로 ──
PREVIEW_DIR = Path(__file__).parent
OUTPUT = PREVIEW_DIR / "slide_what_native.pptx"

# ── 색상 ──
CORAL = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CORAL_SUBTLE = RGBColor(0xFF, 0xF8, 0xF9)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_300 = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Pretendard"
SLIDE_W = 13.333  # 인치 (16:9)


# ── 헬퍼 ──
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title):
    """헤더: 제목 43pt + coral 밑줄"""
    add_text(slide, 0.7, 0.35, 10, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    """인사이트 박스: 핑크 배경 + 20pt"""
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL, bold=False)


# ── PPTX 빌드 ──
def build():
    print("=" * 50)
    print("슬라이드 5: WHAT — 객실당 수익의 구조")
    print("=" * 50)

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    s = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 1) 헤더 ──
    add_header(s, "WHAT — 객실당 수익의 구조")
    add_insight_box(s,
        "객실 요금과 예약률의 균형점을 찾아, 호스트가 직접 바꿀 수 있는 변수로 수익을 최적화하는 방법을 데이터로 검증")

    # ── 2) 수식 한 줄 — 슬라이드 문맥 설정 (작고 절제) ──
    FORMULA_Y = 2.2
    add_text(s, 0.7, FORMULA_Y, 11.8, 0.35,
             "RevPAR  =  ADR  ×  Occupancy Rate",
             size=18, color=GRAY_400, bold=False, align=PP_ALIGN.CENTER)

    # ── 3) 히어로 존 — 58.2%가 슬라이드를 지배 ──
    HERO_Y = 2.85

    # 히어로 배경 패널 (아주 연한 코럴)
    add_rect(s, 0.7, HERO_Y, 11.8, 2.5, CORAL_SUBTLE, radius=0.05)

    # 58.2% — 슬라이드의 주인공
    add_text(s, 0.7, HERO_Y + 0.15, 11.8, 1.2,
             "58.2%",
             size=96, color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # 서브타이틀
    add_text(s, 0.7, HERO_Y + 1.35, 11.8, 0.4,
             "RevPAR 변동의 58%를 점유율 하나가 설명한다",
             size=22, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

    # 보조 설명
    add_text(s, 0.7, HERO_Y + 1.85, 11.8, 0.4,
             "가격(ADR)은 37.5%  ·  교호작용 4.3%  — 점유율이 가격의 1.6배",
             size=15, color=GRAY_500, bold=False, align=PP_ALIGN.CENTER)

    # ── 4) 메시지 카드 2개 — 통계 없이 핵심만 ──
    EVIDENCE_Y = 5.7
    CARD_W = 5.6
    CARD_H = 0.65
    TOTAL_W = 11.8
    GAP = TOTAL_W - CARD_W * 2
    X1 = 0.7
    X2 = X1 + CARD_W + GAP

    # 카드 1: 핵심 메시지
    add_rect(s, X1, EVIDENCE_Y, CARD_W, CARD_H, WHITE,
             border_color=GRAY_200, radius=0.03)
    add_text(s, X1 + 0.25, EVIDENCE_Y + 0.1, CARD_W - 0.5, 0.3,
             "가격 인하 ≠ 점유율 상승", size=17, color=CHARCOAL, bold=True)
    add_text(s, X1 + 0.25, EVIDENCE_Y + 0.38, CARD_W - 0.5, 0.25,
             "가격과 예약률은 독립적으로 움직인다", size=11, color=GRAY_500)

    # 카드 2: 결론
    add_rect(s, X2, EVIDENCE_Y, CARD_W, CARD_H, CORAL_LIGHT,
             border_color=CORAL, radius=0.03)
    add_text(s, X2 + 0.25, EVIDENCE_Y + 0.1, CARD_W - 0.5, 0.3,
             "점유율을 높이는 변수를 찾아라", size=17, color=CORAL, bold=True)
    add_text(s, X2 + 0.25, EVIDENCE_Y + 0.38, CARD_W - 0.5, 0.25,
             "이 프로젝트의 핵심 질문", size=11, color=GRAY_500)

    # 구분선 (히어로와 카드 사이 여백 강조)
    add_rect(s, 3.5, EVIDENCE_Y - 0.25, 6.3, 0.008, GRAY_200)

    # 페이지 번호
    add_text(s, 12.5, 7.0, 0.5, 0.3, "5",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # 저장
    prs.save(str(OUTPUT))
    print(f"\n✓ PPTX 저장: {OUTPUT}")
    print("완료!")


if __name__ == "__main__":
    build()
