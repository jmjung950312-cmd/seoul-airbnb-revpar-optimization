"""
슬라이드 6: HOW — 분석 프레임워크 & AI 도입 효과
- 상단: 6단계 파이프라인 플로우 (간결)
- 중단: AI 도입 효과 3 카드 (객관성 · 품질 · 재현성)
- 하단: 인사이트 배너
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 디자인 토큰 ──
DOMINANT   = RGBColor(0x48, 0x48, 0x48)
SUPPORTING = RGBColor(0x76, 0x76, 0x76)
ACCENT     = RGBColor(0xFF, 0x5A, 0x5F)
ACCENT_SEC = RGBColor(0x00, 0xA6, 0x99)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG    = RGBColor(0xF5, 0xF5, 0xF5)
INSIGHT_BG = RGBColor(0xFF, 0xF0, 0xF3)
ORANGE     = RGBColor(0xFC, 0x64, 0x2D)
BORDER     = RGBColor(0xE8, 0xE8, 0xE8)

FONT = "Noto Sans KR"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = Inches(11.733)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, radius=Inches(0.15)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # 둥근 정도 조절 (0~1, 작을수록 덜 둥글)
    shape.adjustments[0] = 0.08
    return shape


def set_text(shape, text, size=14, bold=False, color=DOMINANT, align=PP_ALIGN.LEFT, font_name=FONT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tf


def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=DOMINANT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return txBox


slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ═══════════════════════════════════════════
# 1. 타이틀 + 악센트 언더라인
# ═══════════════════════════════════════════
title_box = add_text_box(
    slide, MARGIN_L, Inches(0.5), Inches(10), Inches(0.8),
    "HOW — 분석 프레임워크 & AI 도입 효과",
    size=28, bold=True, color=DOMINANT
)

# 악센트 언더라인
slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_L, Inches(1.25), CONTENT_W, Pt(3)
).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = ACCENT
slide.shapes[-1].line.fill.background()

# ═══════════════════════════════════════════
# 2. 6단계 파이프라인 플로우 (상단)
# ═══════════════════════════════════════════
steps = ["데이터 수집", "전처리", "EDA", "모델링", "대시보드", "자동화"]
step_colors = [SUPPORTING, SUPPORTING, ACCENT, ACCENT, ACCENT_SEC, ACCENT_SEC]

flow_y = Inches(1.65)
step_w = Inches(1.45)
step_h = Inches(0.42)
gap = Inches(0.28)
total_flow_w = len(steps) * step_w + (len(steps) - 1) * gap
flow_start_x = MARGIN_L + (CONTENT_W - total_flow_w) / 2

for i, (step, color) in enumerate(zip(steps, step_colors)):
    x = flow_start_x + i * (step_w + gap)
    pill = add_rounded_rect(slide, x, flow_y, step_w, step_h, color, radius=Inches(0.1))
    pill.adjustments[0] = 0.5  # 완전 둥근 pill
    set_text(pill, step, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    pill.text_frame.paragraphs[0].space_before = Pt(0)
    pill.text_frame.paragraphs[0].space_after = Pt(0)
    pill.text_frame.auto_size = None

    # 화살표 (마지막 제외)
    if i < len(steps) - 1:
        arrow_x = x + step_w + Inches(0.04)
        add_text_box(
            slide, arrow_x, flow_y, gap - Inches(0.08), step_h,
            "→", size=16, bold=True, color=SUPPORTING, align=PP_ALIGN.CENTER
        )

# "AI가 자동 수행" 브래킷 텍스트
bracket_x = flow_start_x + 2 * (step_w + gap)
bracket_w = 2 * step_w + gap
add_text_box(
    slide, bracket_x, flow_y + step_h + Inches(0.02), bracket_w, Inches(0.3),
    "▲ AI 에이전트가 자동 수행", size=10, color=ACCENT, align=PP_ALIGN.CENTER
)

# ═══════════════════════════════════════════
# 3. AI 도입 효과 3 카드 (메인 영역)
# ═══════════════════════════════════════════
cards = [
    {
        "icon": "🎯",
        "title": "객관적 가설 수립",
        "result": "9개 가설 자동 도출",
        "detail": "분석자 편향 제거\n놓칠 수 있는 가설까지 커버",
        "accent": ACCENT,
    },
    {
        "icon": "🛡️",
        "title": "일관된 품질 검증",
        "result": "누수 변수 6개 자동 검출",
        "detail": "잘못된 예측이 호스트에게\n전달되는 것 방지",
        "accent": ORANGE,
    },
    {
        "icon": "🔄",
        "title": "반복 가능한 분석 체계",
        "result": "동일 데이터 → 동일 결과",
        "detail": "월별 데이터 갱신 시\n재분석 즉시 가능",
        "accent": ACCENT_SEC,
    },
]

card_w = Inches(3.5)
card_h = Inches(3.0)
card_gap = Inches(0.37)
total_cards_w = 3 * card_w + 2 * card_gap
card_start_x = MARGIN_L + (CONTENT_W - total_cards_w) / 2
card_y = Inches(2.65)

for i, card in enumerate(cards):
    x = card_start_x + i * (card_w + card_gap)

    # 카드 배경
    bg = add_rounded_rect(slide, x, card_y, card_w, card_h, CARD_BG, radius=Inches(0.12))
    bg.adjustments[0] = 0.06

    # 상단 악센트 바
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, card_y, card_w, Pt(5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = card["accent"]
    accent_bar.line.fill.background()
    accent_bar.adjustments[0] = 1.0

    # 아이콘
    add_text_box(
        slide, x + Inches(0.3), card_y + Inches(0.3), Inches(0.6), Inches(0.6),
        card["icon"], size=32, align=PP_ALIGN.LEFT
    )

    # 타이틀
    add_text_box(
        slide, x + Inches(0.3), card_y + Inches(0.9), card_w - Inches(0.6), Inches(0.4),
        card["title"], size=18, bold=True, color=DOMINANT
    )

    # 결과 숫자/텍스트 (강조)
    add_text_box(
        slide, x + Inches(0.3), card_y + Inches(1.35), card_w - Inches(0.6), Inches(0.35),
        card["result"], size=14, bold=True, color=card["accent"]
    )

    # 상세 설명
    add_text_box(
        slide, x + Inches(0.3), card_y + Inches(1.8), card_w - Inches(0.6), Inches(0.9),
        card["detail"], size=12, color=SUPPORTING
    )

# ═══════════════════════════════════════════
# 4. 인사이트 배너 (하단)
# ═══════════════════════════════════════════
insight_y = Inches(6.0)
insight_h = Inches(0.55)

# 배경
insight_bg = add_rounded_rect(
    slide, MARGIN_L, insight_y, CONTENT_W, insight_h,
    hex_to_rgb("FFF0F3"), radius=Inches(0.08)
)
insight_bg.adjustments[0] = 0.15

# 좌측 악센트 바
left_bar = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    MARGIN_L, insight_y, Pt(5), insight_h
)
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = ACCENT
left_bar.line.fill.background()

# 인사이트 텍스트
insight_txBox = slide.shapes.add_textbox(
    MARGIN_L + Inches(0.3), insight_y, CONTENT_W - Inches(0.5), insight_h
)
tf = insight_txBox.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

# "💡 인사이트:" 볼드
run1 = p.add_run()
run1.text = "💡 인사이트: "
run1.font.size = Pt(13)
run1.font.bold = True
run1.font.color.rgb = ACCENT
run1.font.name = FONT

# 본문
run2 = p.add_run()
run2.text = "AI를 도입한 이유는 기술이 아니라, 분석의 "
run2.font.size = Pt(13)
run2.font.color.rgb = DOMINANT
run2.font.name = FONT

# 볼드 강조
run3 = p.add_run()
run3.text = "객관성·품질·재현성"
run3.font.size = Pt(13)
run3.font.bold = True
run3.font.color.rgb = DOMINANT
run3.font.name = FONT

run4 = p.add_run()
run4.text = "을 담보하기 위해서입니다"
run4.font.size = Pt(13)
run4.font.color.rgb = DOMINANT
run4.font.name = FONT

# ═══════════════════════════════════════════
# 5. 슬라이드 번호
# ═══════════════════════════════════════════
add_text_box(
    slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.5), Inches(0.5), Inches(0.3),
    "6", size=11, color=SUPPORTING, align=PP_ALIGN.RIGHT
)

# ═══════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════
out = "/Users/jungmo/Documents/data-projects/seoul_airbnb_cleaned_B/slide_preview/slide06_HOW_v2.pptx"
prs.save(out)
print(f"저장 완료: {out}")
