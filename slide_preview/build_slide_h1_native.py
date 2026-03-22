"""
H1 슈퍼호스트 프리미엄 — 네이티브 python-pptx 빌드
PowerPoint에서 직접 편집 가능한 버전
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUTPUT = Path(__file__).parent / "slide_h1_native.pptx"

# ── 색상 ──
CORAL = RGBColor(0xFF, 0x38, 0x5C)      # 에어비앤비 공식
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3) # 연한 핑크 배경
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF7, 0xF7, 0xF7)
BAR_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Pretendard"


def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    # 첫 번째 문단은 이미 존재
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)

    return txBox


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    """사각형 도형 추가"""
    if radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        # python-pptx 라운드 코너 조정
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRAY_400, width=1.5, dash=True):
    """직선 추가"""
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    if dash:
        connector.line.dash_style = 4  # dash
    return connector


def build():
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333"
    prs.slide_height = Emu(6858000)   # 7.5"

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # ════════════════════════════════════════
    # 헤더
    # ════════════════════════════════════════
    add_text(slide, 0.7, 0.45, 3, 0.25, "EDA · 가설검증",
             size=13, color=GRAY_500, bold=False)
    add_text(slide, 0.7, 0.7, 6, 0.5, "H1 — 슈퍼호스트 프리미엄 효과",
             size=28, color=CHARCOAL, bold=True)

    # 제목 밑줄 (coral)
    add_rect(slide, 0.7, 1.18, 1.5, 0.03, CORAL)

    # ════════════════════════════════════════
    # 왼쪽: 히어로 패널
    # ════════════════════════════════════════
    L = 0.7   # 왼쪽 시작

    # +83.1% 히어로 숫자
    add_text(slide, L, 2.5, 4.5, 1.0, "+83.1%",
             size=80, color=CORAL, bold=True)

    # 라벨
    add_text(slide, L, 3.55, 4.5, 0.35, "슈퍼호스트 RevPAR 프리미엄",
             size=20, color=CHARCOAL, bold=True)

    # 서브 설명
    add_text(slide, L, 3.95, 4.5, 0.55,
             "일반 호스트 중위값 ₩31,825 대비\n슈퍼호스트 중위값 ₩61,205",
             size=15, color=GRAY_500)

    # ── 보조 카드 1: +148.2% ──
    card1_y = 4.75
    # 배경 (연한 핑크)
    add_rect(slide, L, card1_y, 4.2, 0.7, CORAL_LIGHT, radius=0.05)
    # 좌측 액센트 바
    add_rect(slide, L, card1_y, 0.03, 0.7, CORAL)
    # 숫자
    add_text(slide, L + 0.2, card1_y + 0.12, 1.3, 0.4, "+148.2%",
             size=24, color=CORAL, bold=True)
    # 설명
    add_text(slide, L + 1.6, card1_y + 0.08, 2.4, 0.55,
             "슈퍼호스트 + 즉시예약\n두 조건 결합 시 프리미엄",
             size=14, color=GRAY_500)

    # ── 보조 카드 2: ₩47,850 ──
    card2_y = 5.6
    add_rect(slide, L, card2_y, 4.2, 0.7, GRAY_BG, radius=0.05)
    add_text(slide, L + 0.2, card2_y + 0.12, 1.4, 0.4, "₩47,850",
             size=24, color=CHARCOAL, bold=True)
    add_text(slide, L + 1.7, card2_y + 0.08, 2.3, 0.55,
             "Active+Operating\n전체 중위값 (기준선)",
             size=14, color=GRAY_500)

    # ════════════════════════════════════════
    # 세로 구분선
    # ════════════════════════════════════════
    add_line(slide, 5.3, 1.8, 5.3, 6.5, color=GRAY_200, width=1, dash=False)

    # ════════════════════════════════════════
    # 오른쪽: 바 차트
    # ════════════════════════════════════════
    CX = 5.8    # 차트 영역 시작 X
    CW = 7.0    # 차트 영역 폭
    BW = 1.2    # 바 너비
    GAP = 0.55  # 바 간격
    BAR_BOTTOM = 5.95  # 바 하단 Y
    BAR_MAX_H = 3.6    # 최대 바 높이 (63708 기준)

    # 차트 제목
    add_text(slide, CX, 1.6, CW, 0.3, "슈퍼호스트 × 예약방식별 RevPAR 중위값",
             size=14, color=GRAY_500, bold=True)

    # 바 데이터: (라벨, 값, 색상여부)
    bars = [
        ("슈퍼호스트\n즉시예약", 63708, True),
        ("슈퍼호스트\n수동승인", 59497, True),
        ("일반 호스트\n즉시예약", 32863, False),
        ("일반 호스트\n수동승인", 30846, False),
    ]

    max_val = 63708

    for i, (label, val, is_coral) in enumerate(bars):
        bx = CX + i * (BW + GAP) + 0.3
        bar_h = (val / max_val) * BAR_MAX_H
        bar_top = BAR_BOTTOM - bar_h
        bar_color = CORAL if is_coral else BAR_GRAY
        val_color = CORAL if is_coral else CHARCOAL

        # 바
        add_rect(slide, bx, bar_top, BW, bar_h, bar_color, radius=0.02)

        # 값 라벨 (바 위)
        add_text(slide, bx - 0.15, bar_top - 0.35, BW + 0.3, 0.3,
                 f"₩{val:,}",
                 size=16, color=val_color, bold=True, align=PP_ALIGN.CENTER)

        # 카테고리 라벨 (바 아래)
        add_text(slide, bx - 0.15, BAR_BOTTOM + 0.08, BW + 0.3, 0.5,
                 label,
                 size=13, color=GRAY_500, align=PP_ALIGN.CENTER)

    # 중위값 점선: 47850/63708 × BAR_MAX_H = 2.703"
    median_h = (47850 / max_val) * BAR_MAX_H
    median_y = BAR_BOTTOM - median_h
    line_x1 = CX + 0.1
    line_x2 = CX + CW - 0.3
    add_line(slide, line_x1, median_y, line_x2, median_y,
             color=GRAY_400, width=1.2, dash=True)

    # 중위값 라벨
    add_text(slide, line_x2 - 1.8, median_y - 0.28, 2.0, 0.25,
             "전체 중위값 ₩47,850",
             size=12, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # ── 범례 ──
    legend_y = 6.55
    # 구분선
    add_line(slide, CX + 0.3, legend_y, CX + CW - 0.5, legend_y,
             color=GRAY_200, width=0.8, dash=False)

    # 슈퍼호스트 범례
    lx_center = CX + CW / 2
    add_rect(slide, lx_center - 1.4, legend_y + 0.12, 0.15, 0.15, CORAL, radius=0.02)
    add_text(slide, lx_center - 1.2, legend_y + 0.1, 1.0, 0.2, "슈퍼호스트",
             size=14, color=GRAY_500, bold=False)

    # 일반 호스트 범례
    add_rect(slide, lx_center + 0.3, legend_y + 0.12, 0.15, 0.15, BAR_GRAY, radius=0.02)
    add_text(slide, lx_center + 0.5, legend_y + 0.1, 1.0, 0.2, "일반 호스트",
             size=14, color=GRAY_500, bold=False)

    # ── 페이지 번호 ──
    add_text(slide, 12.5, 7.0, 0.5, 0.3, "13",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    # 저장
    prs.save(str(OUTPUT))
    print(f"✓ PPTX 저장: {OUTPUT}")
    print("  → PowerPoint에서 열어 직접 편집 가능합니다")


if __name__ == "__main__":
    build()
