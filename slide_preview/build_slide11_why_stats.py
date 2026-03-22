"""
슬라이드 11 — 왜 통계적 검증이 필요한가?
가설 검증 슬라이드(12) 앞에 삽입하는 소개 슬라이드
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 경로 ──
OUTPUT = Path(__file__).parent / "slide11_why_stats.pptx"

# ── 색상 (디자인 시스템 기준) ──
CORAL = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x00, 0xA6, 0x99)
DARK_BG = RGBColor(0x2B, 0x2B, 0x2B)

FONT = "Pretendard"


# ═══════════════════════════════════════════
# 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, line_spacing=None, anchor=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    if anchor:
        tf.paragraphs[0].alignment = align
        tf.word_wrap = True
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRAY_200, width=1):
    conn = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_header(slide, title):
    """헤더: 제목 43pt + coral 밑줄"""
    add_text(slide, 0.7, 0.35, 10, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_icon_circle(slide, left, top, text, bg_color=CORAL):
    """아이콘 원형 (숫자 또는 기호)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# 빌드
# ═══════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 헤더 ──
    add_header(s, "왜 통계적 검증이 필요한가?")

    # ══════════════════════════════════════
    # 좌측: 목적 + 검증 방법 선택 기준
    # ══════════════════════════════════════
    left_x = 0.7
    left_w = 5.6

    # ── 카드 1: 목적 ──
    card1_y = 1.35
    card1_h = 2.1
    add_rect(s, left_x, card1_y, left_w, card1_h, WHITE, border_color=GRAY_200, radius=0.03)
    # 좌측 coral 악센트 라인
    add_rect(s, left_x, card1_y, 0.05, card1_h, CORAL)

    add_icon_circle(s, left_x + 0.3, card1_y + 0.25, "?", CORAL)
    add_text(s, left_x + 0.95, card1_y + 0.28, 4.3, 0.35,
             "목적", size=22, color=CHARCOAL, bold=True)

    add_text(s, left_x + 0.95, card1_y + 0.75, 4.3, 1.1,
             "EDA에서 발견된 패턴이\n단순한 관찰이 아니라\n실제 데이터에서 유의한 차이인지\n확인하기 위해 통계 검증을 수행",
             size=15, color=GRAY_500, line_spacing=6)

    # ── 카드 2: 검증 방법 선택 기준 ──
    card2_y = 3.7
    card2_h = 2.85
    add_rect(s, left_x, card2_y, left_w, card2_h, WHITE, border_color=GRAY_200, radius=0.03)
    # 좌측 teal 악센트 라인
    add_rect(s, left_x, card2_y, 0.05, card2_h, TEAL)

    add_icon_circle(s, left_x + 0.3, card2_y + 0.25, "!", TEAL)
    add_text(s, left_x + 0.95, card2_y + 0.28, 4.3, 0.35,
             "비모수 검정을 선택한 이유", size=22, color=CHARCOAL, bold=True)

    # 이유 항목들
    reasons = [
        ("RevPAR 분포가 정규분포를 따르지 않음", "skewness = 3.76, 극도의 우편향"),
        ("중위값과 평균의 괴리가 큼", "중위값 ₩8,868 vs 평균 ₩31,310"),
        ("순위·분포 기반 검정이 더 신뢰성 높음", "평균 기반 검정은 이상치에 취약"),
    ]

    for i, (main_text, sub_text) in enumerate(reasons):
        ry = card2_y + 0.85 + i * 0.6
        # 번호
        add_rect(s, left_x + 0.4, ry + 0.02, 0.28, 0.28, GRAY_BG, radius=0.15)
        add_text(s, left_x + 0.4, ry + 0.03, 0.28, 0.26,
                 str(i + 1), size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        # 주 텍스트
        add_text(s, left_x + 0.85, ry, 4.3, 0.25,
                 main_text, size=14, color=CHARCOAL, bold=True)
        # 보조 텍스트
        add_text(s, left_x + 0.85, ry + 0.26, 4.3, 0.25,
                 sub_text, size=12, color=GRAY_400)

    # ══════════════════════════════════════
    # 우측: 검증 방법 역할 테이블
    # ══════════════════════════════════════
    right_x = 6.6
    right_w = 5.9

    add_text(s, right_x, 1.35, right_w, 0.4,
             "사용된 검증 방법", size=22, color=CHARCOAL, bold=True)

    # ── 테이블 헤더 ──
    table_y = 1.85
    header_h = 0.45
    add_rect(s, right_x, table_y, right_w, header_h, CHARCOAL, radius=0.02)
    add_text(s, right_x + 0.2, table_y + 0.1, 1.6, 0.3,
             "방법", size=14, color=WHITE, bold=True)
    add_text(s, right_x + 2.0, table_y + 0.1, 1.5, 0.3,
             "비교 대상", size=14, color=WHITE, bold=True)
    add_text(s, right_x + 3.5, table_y + 0.1, 2.2, 0.3,
             "역할", size=14, color=WHITE, bold=True)

    # ── 테이블 행 ──
    methods = [
        {
            "name": "Spearman",
            "target": "연속 변수 쌍",
            "role": "변수 간 관계 확인\n→ 순위 기반 상관 분석",
            "used_in": "H3, H6",
            "color": TEAL,
        },
        {
            "name": "Mann-Whitney U",
            "target": "두 그룹",
            "role": "두 집단 간 수익 차이 검증\n→ 비모수 독립표본 검정",
            "used_in": "H1, H5",
            "color": CORAL,
        },
        {
            "name": "Kruskal-Wallis",
            "target": "세 그룹 이상",
            "role": "다중 집단 간 수익 차이 검증\n→ 비모수 분산분석(ANOVA 대안)",
            "used_in": "H2, H4",
            "color": RGBColor(0x4A, 0x90, 0xD9),
        },
    ]

    row_h = 1.15
    row_gap = 0.06

    for i, m in enumerate(methods):
        ry = table_y + header_h + 0.06 + i * (row_h + row_gap)
        bg = WHITE if i % 2 == 0 else GRAY_BG

        # 행 배경
        add_rect(s, right_x, ry, right_w, row_h, bg, border_color=GRAY_200, radius=0.02)

        # 좌측 색상 바
        add_rect(s, right_x, ry, 0.05, row_h, m["color"])

        # 방법명
        add_text(s, right_x + 0.2, ry + 0.15, 1.6, 0.3,
                 m["name"], size=15, color=CHARCOAL, bold=True)

        # 적용 가설 뱃지
        add_rect(s, right_x + 0.2, ry + 0.6, 0.8, 0.3, m["color"], radius=0.04)
        add_text(s, right_x + 0.2, ry + 0.63, 0.8, 0.24,
                 m["used_in"], size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # 비교 대상
        add_text(s, right_x + 2.0, ry + 0.15, 1.3, 0.3,
                 m["target"], size=13, color=GRAY_500, bold=True)

        # 역할
        add_text(s, right_x + 3.5, ry + 0.12, 2.2, 0.9,
                 m["role"], size=13, color=GRAY_500, line_spacing=4)

    # ── 하단 인사이트 박스 ──
    add_rect(s, 0.7, 6.75, 11.8, 0.45, CORAL_LIGHT, radius=0.04)
    add_text(s, 1.0, 6.82, 11.2, 0.3,
             "모든 가설이 p < 0.001에서 통계적으로 유의미 → EDA 패턴이 우연이 아님을 확인",
             size=18, color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # ── 페이지 번호 ──
    add_text(s, 12.5, 7.0, 0.5, 0.3, "11",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    prs.save(str(OUTPUT))
    print(f"저장 완료: {OUTPUT}")


if __name__ == "__main__":
    build()
