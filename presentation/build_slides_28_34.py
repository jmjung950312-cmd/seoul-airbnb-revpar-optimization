"""
슬라이드 28~34 — 섹션 6 + 마무리 (7장)
PDF 기준 페이지 번호 사용 (발표용_.pdf 기준)

디자인 시스템: build_eda_slides.py 패턴 통일
인지 과학 원칙:
  - Miller's Law: 슬라이드당 3~5 청크
  - Peak-End Rule: 32(임팩트)=Peak, 33(한계점)=End
  - Progressive Disclosure: 핵심만 보여주고 디테일은 대시보드로
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 디자인 토큰 (build_eda_slides.py 기준) ──
CORAL       = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL    = RGBColor(0x48, 0x48, 0x48)
GRAY_BG     = RGBColor(0xF7, 0xF7, 0xF7)
GRAY_400    = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500    = RGBColor(0x76, 0x76, 0x76)
GRAY_200    = RGBColor(0xE8, 0xE8, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT  = RGBColor(0xE0, 0xF7, 0xF5)
ORANGE      = RGBColor(0xFC, 0x64, 0x2D)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xED)
BLUE        = RGBColor(0x4A, 0x90, 0xD9)
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)

FONT = "Pretendard"
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


# ═══════════════════════════════════════════
# 공통 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title):
    add_text(slide, 0.7, 0.35, 10.5, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL)


def add_page_num(slide, num):
    add_text(slide, 12.5, 7.0, 0.5, 0.3, str(num),
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_kpi_card(slide, left, top, width, height, value, label,
                 value_color=CHARCOAL, bar_color=None, accent=False,
                 value_size=36, label_size=13):
    """KPI 카드: 큰 숫자 + 설명 + 좌측 바"""
    bg = CORAL_LIGHT if accent else GRAY_BG
    add_rect(slide, left, top, width, height, bg, radius=0.04)
    if bar_color:
        add_rect(slide, left, top + 0.06, 0.035, height - 0.12, bar_color)
    add_text(slide, left + 0.2, top + 0.1, width - 0.4, 0.5,
             value, size=value_size, color=value_color, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.2, top + height - 0.4, width - 0.4, 0.3,
             label, size=label_size, color=GRAY_500, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# p28: 비즈니스 모델 — 무료에서 시작, 단계별 확장
# ═══════════════════════════════════════════
def build_slide_28(prs):
    s = new_slide(prs)
    add_header(s, "비즈니스 모델 — 무료에서 시작, 단계별 확장")
    add_insight_box(s,
        "💡 인사이트:  무료 위험 감지로 진입장벽 제거 → "
        "가치를 경험한 호스트가 유료 전환하는 Freemium 모델")

    # ── 3-tier 카드 (가로 배치) ──
    tiers = [
        {
            "name": "Free",
            "price": "무료",
            "features": ["숙소 사전 진단", "위험 신호 알림", "이메일 리포트"],
            "color": GRAY_BG,
            "name_color": CHARCOAL,
            "price_color": TEAL,
            "accent": False,
        },
        {
            "name": "Basic",
            "price": "W9,900/월",
            "features": ["수익 예측", "기본 최적화 가이드", "행동 추천"],
            "color": GRAY_BG,
            "name_color": CHARCOAL,
            "price_color": CORAL,
            "accent": False,
        },
        {
            "name": "Enterprise",
            "price": "문의",
            "features": ["AI 에이전트 자동분석", "맞춤 가설 검증", "자동 리포팅"],
            "color": CORAL,
            "name_color": WHITE,
            "price_color": WHITE,
            "accent": True,
        },
    ]

    card_w = 3.6
    card_gap = 0.3
    card_start = 0.85
    card_y = 2.15
    card_h = 4.2

    for i, t in enumerate(tiers):
        cx = card_start + i * (card_w + card_gap)
        bg = t["color"]

        # 카드 배경
        add_rect(s, cx, card_y, card_w, card_h, bg, radius=0.06)

        # 플랜명
        add_text(s, cx, card_y + 0.4, card_w, 0.4,
                 t["name"], size=28, color=t["name_color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 가격
        add_text(s, cx, card_y + 0.9, card_w, 0.35,
                 t["price"], size=22, color=t["price_color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 구분선
        line_color = WHITE if t["accent"] else GRAY_200
        add_rect(s, cx + 0.5, card_y + 1.4, card_w - 1.0, 0.015, line_color)

        # 기능 리스트
        feat_color = WHITE if t["accent"] else CHARCOAL
        for j, feat in enumerate(t["features"]):
            fy = card_y + 1.7 + j * 0.55
            add_text(s, cx + 0.4, fy, card_w - 0.8, 0.4,
                     f"· {feat}", size=16, color=feat_color)

        # 화살표 (마지막 카드 제외)
        if i < 2:
            ax = cx + card_w + 0.02
            add_text(s, ax, card_y + card_h / 2 - 0.2, 0.25, 0.4,
                     "→", size=24, color=GRAY_400, align=PP_ALIGN.CENTER)

    add_page_num(s, 28)
    return s


# ═══════════════════════════════════════════
# p29: 위험 감지 — R1~R5 + 자동 플로우
# ═══════════════════════════════════════════
def build_slide_29(prs):
    s = new_slide(prs)
    add_header(s, "24/7 위험 감지 — 5개 규칙이 쉬지 않습니다")
    add_insight_box(s,
        "💡 인사이트:  데이터 변경 시 자동 실행 → 리스크 분류 → "
        "HIGH 발견 시 즉시 이메일 알림. 호스트 개입 없이 24시간 동작",
        height=0.6)

    # ── R1~R5 미니카드 5개 (가로 1행) ──
    rules = [
        ("R1", "유령수익",       "리뷰 0 + 수익 > 1천만", "높음",  CORAL),
        ("R2", "점유율-수익\n불일치", "점유율 0% & 수익 > 0", "높음", CORAL),
        ("R3", "가격 이상치",    "ADR > 구평균+3σ",       "중간",  ORANGE),
        ("R4", "유령 액티브",    "Active & 수익 < 5백만",  "중간",  ORANGE),
        ("R5", "평점 없는\n고수익", "수익 > 1천만 & 평점 무", "높음", CORAL),
    ]

    card_w = 2.2
    card_gap = 0.15
    card_start_x = 0.7
    card_y = 1.95
    card_h = 1.55

    for i, (rid, name, cond, level, color) in enumerate(rules):
        cx = card_start_x + i * (card_w + card_gap)

        add_rect(s, cx, card_y, card_w, card_h, GRAY_BG, radius=0.04)
        # 상단 accent 바
        add_rect(s, cx + 0.1, card_y, card_w - 0.2, 0.035, color)

        # R# 배지 + LEVEL 배지
        add_rect(s, cx + 0.15, card_y + 0.15, 0.42, 0.24, color, radius=0.10)
        add_text(s, cx + 0.15, card_y + 0.16, 0.42, 0.24,
                 rid, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        level_bg = CORAL_LIGHT if level == "높음" else ORANGE_LIGHT
        add_rect(s, cx + card_w - 0.85, card_y + 0.15, 0.7, 0.24,
                 level_bg, radius=0.10)
        add_text(s, cx + card_w - 0.85, card_y + 0.16, 0.7, 0.24,
                 level, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)

        # 규칙명
        add_text(s, cx + 0.15, card_y + 0.5, card_w - 0.3, 0.45,
                 name, size=14, color=CHARCOAL, bold=True)

        # 조건
        add_rect(s, cx + 0.1, card_y + 1.05, card_w - 0.2, 0.38,
                 RGBColor(0xEE, 0xEE, 0xEE), radius=0.03)
        add_text(s, cx + 0.2, card_y + 1.1, card_w - 0.4, 0.3,
                 cond, size=11, color=GRAY_500)

    # ── 자동화 플로우 ──
    flow_y = 3.7
    add_text(s, 0.7, flow_y, 3, 0.25,
             "자동 실행 플로우", size=15, color=CHARCOAL, bold=True)

    steps = [
        ("데이터 변경", "hooks.py", GRAY_BG, CHARCOAL),
        ("R1~R5 검사", "detector.py", GRAY_BG, CHARCOAL),
        ("IQR / Z-score", "scorer.py", GRAY_BG, CHARCOAL),
        ("리스크 분류", "HIGH / MEDIUM", GRAY_BG, CHARCOAL),
        ("이메일 알림", "자동 발송", CORAL_LIGHT, CORAL),
    ]

    step_w = 2.1
    step_gap = 0.22
    step_y = 4.05

    for i, (title, sub, bg, tc) in enumerate(steps):
        sx = 0.7 + i * (step_w + step_gap)
        add_rect(s, sx, step_y, step_w, 0.75, bg, radius=0.04)
        if bg == CORAL_LIGHT:
            add_rect(s, sx, step_y + 0.06, 0.035, 0.63, CORAL)
        add_text(s, sx + 0.15, step_y + 0.12, step_w - 0.3, 0.3,
                 title, size=13, color=tc, bold=True)
        add_text(s, sx + 0.15, step_y + 0.42, step_w - 0.3, 0.25,
                 sub, size=11, color=GRAY_500)
        if i < len(steps) - 1:
            ax = sx + step_w + 0.02
            add_text(s, ax, step_y + 0.2, 0.2, 0.3,
                     "→", size=18, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 하단 요약 ──
    add_rect(s, 0.7, 5.15, 11.8, 0.55, TEAL_LIGHT, radius=0.04)
    add_rect(s, 0.7, 5.21, 0.035, 0.43, TEAL)
    add_text(s, 1.0, 5.25, 11.0, 0.4,
             "Claude Code Hooks 기반 — 데이터 파일 수정 시 자동 트리거, 중복 알림 7일 필터, "
             "HIGH 리스크(2개+ 규칙) 시에만 이메일 발송",
             size=14, color=TEAL, bold=True)

    # 탐지 실적 카드
    add_kpi_card(s, 0.7, 5.9, 5.7, 0.9,
                 "14건", "전체 플래그 탐지 (R1: 14건, R3: 3건)",
                 CORAL, CORAL)
    add_kpi_card(s, 6.6, 5.9, 5.9, 0.9,
                 "HIGH 1건", "listing_idx=5996 · 송파구 · W20.8M",
                 CORAL, CORAL, accent=True)

    add_page_num(s, 29)
    return s


# ═══════════════════════════════════════════
# p30: 대시보드 — 핵심 4기능
# ═══════════════════════════════════════════
def build_slide_30(prs):
    s = new_slide(prs)
    add_header(s, "호스트 컨설팅 대시보드 — 핵심 기능")
    add_insight_box(s,
        "💡 인사이트:  호스트가 직접 조건을 입력하면 즉시 시장 위치 · 건강점수 · "
        "수익 예측 · 최적 요금을 확인할 수 있는 올인원 도구")

    # ── 좌측: 대시보드 특징 서브헤더 ──
    add_text(s, 0.7, 2.05, 5.5, 0.3,
             "Streamlit 기반 · 실시간 분석", size=15, color=GRAY_500)

    # ── 4개 핵심 기능 카드 (좌측 세로 배치) ──
    features = [
        {
            "num": "1", "title": "시장 위치 4분면",
            "desc": "내 숙소가 서울에서 어디에 있는지\n한눈에 확인 (전체 · 자치구 · 유형별)",
            "color": CORAL,
        },
        {
            "num": "2", "title": "건강점수 게이지",
            "desc": "0~100점 + A~F 등급으로\n숙소 상태 종합 진단",
            "color": TEAL,
        },
        {
            "num": "3", "title": "수익 예측 KPI",
            "desc": "예측 수익 / 시장 대비 위치 /\n개선 잠재력 시각화",
            "color": BLUE,
        },
        {
            "num": "4", "title": "AI 시장 진단",
            "desc": "적정 요금 vs 현재 요금 비교\n군집별 최적 가격 제안",
            "color": PURPLE,
        },
    ]

    LX = 0.7
    LW = 5.6
    card_h = 1.0

    for i, f in enumerate(features):
        fy = 2.5 + i * (card_h + 0.2)
        color = f["color"]

        add_rect(s, LX, fy, LW, card_h, GRAY_BG, radius=0.04)
        add_rect(s, LX, fy + 0.08, 0.04, card_h - 0.16, color)

        # 번호 배지
        add_rect(s, LX + 0.25, fy + 0.15, 0.45, 0.35, color, radius=0.10)
        add_text(s, LX + 0.25, fy + 0.17, 0.45, 0.35,
                 f["num"], size=18, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)

        # 제목 + 설명
        add_text(s, LX + 0.85, fy + 0.12, LW - 1.2, 0.3,
                 f["title"], size=16, color=CHARCOAL, bold=True)
        add_text(s, LX + 0.85, fy + 0.45, LW - 1.2, 0.5,
                 f["desc"], size=12, color=GRAY_500)

    # ── 우측: 대시보드 미리보기 프레임 ──
    RX = 6.65
    RW = 5.85
    RH = 4.55

    # 외곽 프레임 (스크린샷 자리 표시)
    add_rect(s, RX, 2.2, RW, RH, GRAY_BG, border_color=GRAY_200, radius=0.06)

    # 내부 헤더
    add_rect(s, RX + 0.15, 2.35, RW - 0.3, 0.4, CORAL, radius=0.04)
    add_text(s, RX + 0.15, 2.38, RW - 0.3, 0.35,
             "분석 결과", size=16, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)

    # 시뮬레이션 KPI 예시
    mini_cards = [
        ("전체 1박 요금", "W141,517", CHARCOAL),
        ("하루 실수익", "W84,910", CORAL),
        ("월 예상 순이익", "W1,930,886", CORAL),
    ]
    for j, (label, val, vc) in enumerate(mini_cards):
        mx = RX + 0.3 + j * 1.75
        add_rect(s, mx, 2.95, 1.6, 0.75, WHITE, radius=0.03)
        add_text(s, mx + 0.1, 3.0, 1.4, 0.2,
                 label, size=9, color=GRAY_500, align=PP_ALIGN.CENTER)
        add_text(s, mx + 0.1, 3.25, 1.4, 0.3,
                 val, size=13, color=vc, bold=True, align=PP_ALIGN.CENTER)

    # 예약률 예시
    rates = [("전체 예약률", "60%"), ("평일", "0%"), ("주말", "0%")]
    for j, (label, val) in enumerate(rates):
        mx = RX + 0.3 + j * 1.75
        add_rect(s, mx, 3.85, 1.6, 0.55, WHITE, radius=0.03)
        add_text(s, mx + 0.1, 3.88, 1.4, 0.15,
                 label, size=8, color=GRAY_500, align=PP_ALIGN.CENTER)
        add_text(s, mx + 0.1, 4.05, 1.4, 0.3,
                 val, size=16, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

    # 하단 안내
    add_text(s, RX + 0.3, 4.6, RW - 0.6, 0.3,
             "* 실제 대시보드에서 조건별 실시간 시뮬레이션 가능",
             size=10, color=GRAY_400, align=PP_ALIGN.CENTER)

    # 추가 기능 탭 표시
    tabs = ["수익 요약", "요금 전략", "주변 관광지", "운영 개선", "지역 진단"]
    tab_y = 5.0
    for j, tab in enumerate(tabs):
        tx = RX + 0.2 + j * 1.1
        add_rect(s, tx, tab_y, 1.0, 0.25, WHITE, border_color=GRAY_200, radius=0.03)
        add_text(s, tx, tab_y + 0.03, 1.0, 0.2,
                 tab, size=8, color=GRAY_500, align=PP_ALIGN.CENTER)

    # 하단 CTA
    add_rect(s, 0.7, 6.5, 11.8, 0.45, GRAY_BG, radius=0.04)
    add_rect(s, 0.7, 6.56, 0.035, 0.33, TEAL)
    add_text(s, 1.0, 6.58, 11.0, 0.35,
             "같은 지역 · 같은 유형 숙소 대비 상대 점수 → "
             "호스트가 즉시 약점 파악",
             size=15, color=CHARCOAL, bold=True)

    add_page_num(s, 30)
    return s


# ═══════════════════════════════════════════
# p31: 자동화 시스템 + AI 에이전트 아키텍처
# ═══════════════════════════════════════════
def build_slide_31(prs):
    s = new_slide(prs)
    add_header(s, "자동화 시스템 + AI 에이전트 아키텍처")
    add_insight_box(s,
        "💡 인사이트:  Orchestrator가 6개 에이전트를 순차 호출, "
        "도메인 리서치부터 모델 검증까지 End-to-End 자동화",
        height=0.6)

    # ── 섹션 1: 위험 감지 자동화 ──
    add_text(s, 0.7, 1.95, 5, 0.25,
             "① 위험 감지 자동화", size=15, color=CORAL, bold=True)

    auto_steps = [
        ("데이터 변경\n감지", "hooks.py", CORAL_LIGHT),
        ("5규칙\n검사", "detector.py", GRAY_BG),
        ("통계 보조\n검사", "scorer.py", GRAY_BG),
        ("위험도\n산출", "7일 중복 필터", GRAY_BG),
        ("이메일\n발송", "email_alert.py", TEAL_LIGHT),
    ]

    step_w = 2.1
    step_gap = 0.22
    step_y = 2.3

    for i, (title, sub, bg) in enumerate(auto_steps):
        sx = 0.7 + i * (step_w + step_gap)
        add_rect(s, sx, step_y, step_w, 0.8, bg, radius=0.04)
        add_text(s, sx + 0.15, step_y + 0.08, step_w - 0.3, 0.4,
                 title, size=12, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, sx + 0.15, step_y + 0.55, step_w - 0.3, 0.2,
                 sub, size=10, color=GRAY_500, align=PP_ALIGN.CENTER)
        if i < len(auto_steps) - 1:
            add_text(s, sx + step_w + 0.02, step_y + 0.25, 0.2, 0.3,
                     "→", size=16, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 섹션 2: AI 에이전트 구성 ──
    add_text(s, 0.7, 3.35, 5, 0.25,
             "② AI 에이전트 구성 (Enterprise)", size=15, color=TEAL, bold=True)

    # Orchestrator 바
    orch_y = 3.7
    add_rect(s, 3.5, orch_y, 6.0, 0.4, GRAY_BG, radius=0.04)
    add_text(s, 3.5, orch_y + 0.05, 6.0, 0.3,
             "오케스트레이터 (순차 호출)", size=14, color=CHARCOAL, bold=True,
             align=PP_ALIGN.CENTER)

    # 6개 에이전트 (1행)
    agents = [
        ("도메인 연구", "→ 도메인 지식", CORAL),
        ("가설 생성",   "→ H1~H9",      ORANGE),
        ("탐색적 분석", "→ 38개 차트",   TEAL),
        ("변수 설계",   "→ 14개 변수",   BLUE),
        ("모델링",      "→ ML+AI분석",   PURPLE),
        ("모델 검수",   "→ 품질 평가",   RGBColor(0x2C, 0x3E, 0x50)),
    ]

    agent_w = 1.85
    agent_gap = 0.13
    agent_y = 4.25

    for i, (name, output, color) in enumerate(agents):
        ax = 0.7 + i * (agent_w + agent_gap)
        add_rect(s, ax, agent_y, agent_w, 0.85, GRAY_BG, radius=0.04)
        add_rect(s, ax + 0.05, agent_y, agent_w - 0.1, 0.03, color)
        add_text(s, ax + 0.1, agent_y + 0.15, agent_w - 0.2, 0.25,
                 name, size=12, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, ax + 0.1, agent_y + 0.5, agent_w - 0.2, 0.25,
                 output, size=11, color=GRAY_500, align=PP_ALIGN.CENTER)

    # ── 하단 결과 요약 바 ──
    add_rect(s, 0.7, 5.35, 11.8, 0.5, CORAL_LIGHT, radius=0.04)
    add_text(s, 0.7, 5.42, 11.8, 0.35,
             "결과: 9개 가설 검증 / 14개 변수 설계 / 정확도 R²=0.85",
             size=17, color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    # ── 하단 KPI 4개 ──
    results = [
        ("9",       "가설 검증",  CORAL),
        ("14",      "파생변수",   TEAL),
        ("38",      "차트",       BLUE),
        ("R² 0.85", "모델 성능",  PURPLE),
    ]

    result_w = 2.75
    result_gap = 0.2
    result_y = 6.0

    for i, (val, label, color) in enumerate(results):
        rx = 0.7 + i * (result_w + result_gap)
        add_rect(s, rx, result_y, result_w, 0.8, GRAY_BG, radius=0.04)
        add_rect(s, rx, result_y + 0.06, 0.035, 0.68, color)
        add_text(s, rx + 0.15, result_y + 0.08, result_w - 0.3, 0.4,
                 val, size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + 0.15, result_y + 0.5, result_w - 0.3, 0.2,
                 label, size=12, color=GRAY_500, align=PP_ALIGN.CENTER)

    add_page_num(s, 31)
    return s


# ═══════════════════════════════════════════
# p32: 프로젝트 임팩트 요약 (PEAK 슬라이드)
#
# 디자인 전략 (v2):
#   - Hero 앵커링: 85% 대형 숫자로 시선 고정
#   - 3 Discovery: 데이터 증거 (강력한/충격적/반직관적)
#   - 4 Delivery: 구축 결과물 compact badges
#   - Tagline: 감정의 클로징
# ═══════════════════════════════════════════
def build_slide_32(prs):
    s = new_slide(prs)
    add_header(s, "데이터가 증명한 것")

    # 서브타이틀: 스케일과 맥락을 먼저 제시
    add_text(s, 0.7, 1.15, 11.8, 0.35,
             "32,061개 리스팅 · 9개 가설 · 12개월의 데이터가 도달한 결론",
             size=18, color=GRAY_500, align=PP_ALIGN.CENTER)

    # ═══ HERO CARD: 핵심 성취 ═══
    hero_x = 1.2
    hero_w = 10.9
    hero_y = 1.7
    hero_h = 1.45

    add_rect(s, hero_x, hero_y, hero_w, hero_h, CORAL_LIGHT, radius=0.06)
    add_rect(s, hero_x, hero_y + 0.1, 0.05, hero_h - 0.2, CORAL)

    # 85% 대형 숫자 (앵커링 효과)
    add_text(s, hero_x + 0.35, hero_y + 0.1, 2.8, 1.2,
             "85%", size=72, color=CORAL, bold=True)

    # 맥락 설명 — "왜 이것이 대단한가"
    add_text(s, hero_x + 3.3, hero_y + 0.25, hero_w - 3.8, 0.4,
             "호스트가 바꿀 수 있는 것만으로, 수익의 85%가 결정됩니다",
             size=20, color=CHARCOAL, bold=True)
    add_text(s, hero_x + 3.3, hero_y + 0.75, hero_w - 3.8, 0.5,
             "LightGBM · 16개 호스트 통제 변수 · 5-Fold CV R²=0.847",
             size=14, color=GRAY_500)

    # ═══ 3 DISCOVERY CARDS: 가장 강력한 / 가장 충격적인 / 가장 반직관적인 ═══
    discoveries = [
        {
            "number": "+83.1%",
            "label": "슈퍼호스트 프리미엄",
            "detail": "₩79,087 vs ₩43,201\n가장 강력한 수익 레버",
            "color": CORAL,
        },
        {
            "number": "54.3%",
            "label": "시장 절반이 멈춰있다",
            "detail": "17,399개 Dormant 리스팅\n위기 속에 기회가 있다",
            "color": TEAL,
        },
        {
            "number": "4.9 > 5.0",
            "label": "완벽주의의 역설",
            "detail": "₩66,586 vs ₩58,206\n만점이 최선은 아니다",
            "color": BLUE,
        },
    ]

    card_w = 3.73
    card_gap = 0.25
    total_cards_w = 3 * card_w + 2 * card_gap
    card_start = 0.7 + (11.8 - total_cards_w) / 2
    card_y = 3.4
    card_h = 1.75

    for i, d in enumerate(discoveries):
        cx = card_start + i * (card_w + card_gap)

        # 그림자 (깊이감)
        add_rect(s, cx + 0.03, card_y + 0.03,
                 card_w, card_h, GRAY_200, radius=0.05)
        # 카드 본체
        add_rect(s, cx, card_y, card_w, card_h, WHITE,
                 border_color=GRAY_200, radius=0.05)
        # 상단 accent 바
        add_rect(s, cx, card_y, card_w, 0.04, d["color"])

        # 핵심 숫자
        num_size = 36 if len(d["number"]) > 5 else 42
        add_text(s, cx + 0.15, card_y + 0.2, card_w - 0.3, 0.6,
                 d["number"], size=num_size, color=d["color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 라벨
        add_text(s, cx + 0.15, card_y + 0.85, card_w - 0.3, 0.25,
                 d["label"], size=15, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)

        # 상세 근거
        add_text(s, cx + 0.15, card_y + 1.15, card_w - 0.3, 0.5,
                 d["detail"], size=12, color=GRAY_500,
                 align=PP_ALIGN.CENTER)

    # ═══ 4 DELIVERY BADGES: 구축한 시스템 ═══
    delivery_y = 5.4
    items = [
        ("실시간 대시보드", "호스트 의사결정 도구", CORAL),
        ("24/7 위험 감시", "5규칙 자동 탐지", TEAL),
        ("무료 이메일 진단", "숙소 사전 분석 리포트", BLUE),
        ("AI 에이전트", "6-Agent E2E 자동화", PURPLE),
    ]

    item_w = 2.75
    item_gap = 0.17
    total_items_w = 4 * item_w + 3 * item_gap
    item_start = 0.7 + (11.8 - total_items_w) / 2

    for i, (title, desc, color) in enumerate(items):
        ix = item_start + i * (item_w + item_gap)

        add_rect(s, ix, delivery_y, item_w, 0.55, GRAY_BG, radius=0.04)
        add_rect(s, ix, delivery_y + 0.05, 0.035, 0.45, color)
        add_text(s, ix + 0.15, delivery_y + 0.07, item_w - 0.3, 0.2,
                 title, size=13, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, ix + 0.15, delivery_y + 0.3, item_w - 0.3, 0.2,
                 desc, size=11, color=GRAY_500, align=PP_ALIGN.CENTER)

    # ═══ TAGLINE: 마지막 인상 ═══
    add_rect(s, 0.7, 6.2, 11.8, 0.55, CORAL, radius=0.04)
    add_text(s, 0.7, 6.28, 11.8, 0.4,
             "이 프로젝트는 분석에서 끝나지 않습니다.  "
             "호스트의 다음 행동을 설계합니다.",
             size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_page_num(s, 32)
    return s


# ═══════════════════════════════════════════
# p33: 한계점 및 향후 계획 (END 슬라이드)
# ═══════════════════════════════════════════
def build_slide_33(prs):
    s = new_slide(prs)
    add_header(s, "한계점 및 향후 계획")

    # ── 좌측: 현재 한계점 ──
    LX = 0.7
    LW = 5.6

    add_text(s, LX, 1.25, 3, 0.25,
             "현재 한계점", size=17, color=CHARCOAL, bold=True)

    # 한계점 카드 배경
    add_rect(s, LX, 1.6, LW, 4.7, GRAY_BG, radius=0.04)

    limitations = [
        ("계절성 및 이벤트 미반영",
         "명절, 대형 행사, 계절 변동성 등 외부 요인이 모델에 포함되지 않음"),
        ("규제 영향 미반영",
         "2025-10-16 시행된 규제 이후 데이터가 분석에 반영되지 않음"),
        ("정적 배치 분석",
         "실시간 데이터 처리가 아닌 배치 방식으로 즉각적인 대응 제한"),
        ("단일 시점 예측",
         "다기간 예측 및 시나리오 분석 기능 부재로 전략 수립 제약"),
        ("웹 배포 미완",
         "호스트가 직접 접근 가능한 웹 서비스로 배포되지 않음"),
    ]

    for i, (title, desc) in enumerate(limitations):
        ly = 1.8 + i * 0.88

        # 아이콘 대체: 번호
        add_rect(s, LX + 0.2, ly, 0.3, 0.25, ORANGE, radius=0.08)
        add_text(s, LX + 0.2, ly + 0.02, 0.3, 0.25,
                 str(i + 1), size=11, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)

        add_text(s, LX + 0.65, ly, LW - 1.0, 0.25,
                 title, size=14, color=CHARCOAL, bold=True)
        add_text(s, LX + 0.65, ly + 0.28, LW - 1.0, 0.45,
                 desc, size=11, color=GRAY_500)

    # ── 우측: 향후 로드맵 ──
    RX = 6.65
    RW = 5.85

    add_text(s, RX, 1.25, 3, 0.25,
             "향후 로드맵", size=17, color=CHARCOAL, bold=True)

    roadmap = [
        {
            "period": "단기",
            "title": "규제 이후 데이터 재수집 및 모델 재학습",
            "desc": "2025-10-16 규제 시행 이후 시장 변화를 반영한\n데이터 재수집 및 LightGBM 모델 재학습 진행",
            "color": CORAL,
        },
        {
            "period": "2026 상반기",
            "title": "실시간 스트리밍 파이프라인 구축",
            "desc": "배치 방식에서 실시간 처리로 전환하여\n즉각적인 이상 탐지 및 예측 지원",
            "color": TEAL,
        },
        {
            "period": "2026 하반기",
            "title": "다기간 예측 및 시나리오 분석 확장",
            "desc": "1개월/3개월/6개월 예측 및\nWhat-if 시나리오 분석 기능 추가",
            "color": BLUE,
        },
        {
            "period": "장기",
            "title": "웹 배포 및 접근성(WCAG) 고도화",
            "desc": "호스트 직접 접근 가능한 웹 서비스 배포\n및 WCAG 2.1 AA 수준 접근성 준수",
            "color": PURPLE,
        },
    ]

    for i, r in enumerate(roadmap):
        ry = 1.65 + i * 1.15
        color = r["color"]

        # 타임라인 좌측 바
        add_rect(s, RX, ry, 0.04, 1.0, color)

        # 기간 배지
        add_rect(s, RX + 0.2, ry, 1.1, 0.25, color, radius=0.10)
        add_text(s, RX + 0.2, ry + 0.02, 1.1, 0.25,
                 r["period"], size=11, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)

        # 제목
        add_text(s, RX + 0.2, ry + 0.35, RW - 0.5, 0.25,
                 r["title"], size=13, color=CHARCOAL, bold=True)

        # 설명
        add_text(s, RX + 0.2, ry + 0.6, RW - 0.5, 0.45,
                 r["desc"], size=11, color=GRAY_500)

    # ── 하단 마무리 ──
    add_rect(s, 0.7, 6.5, 11.8, 0.45, TEAL_LIGHT, radius=0.04)
    add_rect(s, 0.7, 6.56, 0.035, 0.33, TEAL)
    add_text(s, 1.0, 6.58, 11.0, 0.35,
             "지속적인 개선과 확장을 통해 호스트와 플랫폼 모두에게 "
             "실질적인 가치를 제공하는 시스템으로 발전시켜 나가겠습니다",
             size=14, color=TEAL, bold=True)

    add_page_num(s, 33)
    return s


# ═══════════════════════════════════════════
# p34: Q&A
# ═══════════════════════════════════════════
def build_slide_34(prs):
    s = new_slide(prs)

    # 중앙 coral 카드
    card_w = 10.5
    card_h = 4.2
    card_x = (13.333 - card_w) / 2
    card_y = 0.8

    add_rect(s, card_x, card_y, card_w, card_h, CORAL, radius=0.15)

    # Q & A 텍스트
    add_text(s, card_x, card_y + 1.0, card_w, 1.0,
             "Q & A", size=64, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, card_x, card_y + 2.3, card_w, 0.5,
             "질문과 의견을 환영합니다", size=22, color=WHITE,
             align=PP_ALIGN.CENTER)

    # 구분선
    add_rect(s, card_x + 1.5, card_y + 3.1, card_w - 3.0, 0.01, WHITE)

    # 팀 정보
    add_text(s, 0.7, 5.5, 11.8, 0.4,
             "팀 구매확정 · 김수빈 · 이영은 · 이정모",
             size=18, color=CORAL, bold=True, align=PP_ALIGN.CENTER)

    return s


# ═══════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_28(prs)   # 비즈니스 모델
    build_slide_29(prs)   # 위험 감지
    build_slide_30(prs)   # 대시보드
    build_slide_31(prs)   # 자동화 + AI 에이전트
    build_slide_32(prs)   # 프로젝트 임팩트 (PEAK)
    build_slide_33(prs)   # 한계점 + 로드맵 (END)
    build_slide_34(prs)   # Q&A

    out = "presentation/slides_28_34.pptx"
    prs.save(out)
    print(f"✓ 저장 완료: {out} (7장: 슬라이드 28~34)")
