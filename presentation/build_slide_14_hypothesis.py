"""
슬라이드 14 — 가설 검증 결과 요약 (채택 7개 집중)
슬라이드 A  — 어펜딕스: 부분 채택 가설 (H8, H9)

변경:
  - 본 슬라이드: 채택 7개(H1~H7)만 표시, 메시지 단순화
  - 어펜딕스: H8, H9 부분채택 + 소표본 한계 설명
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 디자인 토큰 (build_eda_slides.py 기준 통일) ──
CORAL       = RGBColor(0xFF, 0x38, 0x5C)   # #FF385C 에어비앤비 공식
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)   # 인사이트 박스
CHARCOAL    = RGBColor(0x48, 0x48, 0x48)
GRAY_BG     = RGBColor(0xF7, 0xF7, 0xF7)   # 카드 배경
GRAY_400    = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500    = RGBColor(0x76, 0x76, 0x76)
GRAY_200    = RGBColor(0xE8, 0xE8, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xFC, 0x64, 0x2D)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xED)
TEAL        = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT  = RGBColor(0xE0, 0xF7, 0xF5)

FONT = "Pretendard"

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


# ── 헬퍼 (build_eda_slides.py 패턴 재현) ──
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT):
    """복수줄 지원 텍스트박스"""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    """사각형/둥근사각형"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title):
    """제목 43pt + coral 밑줄"""
    add_text(slide, 0.7, 0.35, 10, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    """인사이트 박스: 핑크 배경 + 20pt 텍스트"""
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL, bold=False)


def add_hypothesis_row(slide, x, y, w, h_id, title, result, accent_color):
    """가설 행: GRAY_BG 배경 + 좌측 accent 바 + H# 배지 + 제목 + 결과 + ✓ 채택"""
    row_h = 0.48

    # 카드 배경
    add_rect(slide, x, y, w, row_h, GRAY_BG, radius=0.04)

    # 좌측 accent 바
    add_rect(slide, x, y + 0.05, 0.035, row_h - 0.10, accent_color)

    # H# 배지
    add_rect(slide, x + 0.15, y + 0.11, 0.45, 0.26,
             accent_color, radius=0.12)
    add_text(slide, x + 0.15, y + 0.12, 0.45, 0.26, h_id,
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 가설 제목
    add_text(slide, x + 0.7, y + 0.09, w - 2.6, 0.3,
             title, size=14, color=CHARCOAL, bold=True)

    # 결과 텍스트
    add_text(slide, x + w - 1.85, y + 0.09, 1.1, 0.3,
             result, size=13, color=accent_color, bold=True,
             align=PP_ALIGN.RIGHT)

    # ✓ 채택
    add_text(slide, x + w - 0.7, y + 0.09, 0.65, 0.3,
             "✓ 채택", size=12, color=CORAL, bold=True,
             align=PP_ALIGN.RIGHT)

    return row_h


def add_summary_card(slide, left, top, width, height,
                     value, label, value_color, bar_color):
    """하단 요약 카드: 큰 숫자 + 설명 + 좌측 바"""
    add_rect(slide, left, top, width, height, GRAY_BG, radius=0.04)
    add_rect(slide, left, top + 0.06, 0.035, height - 0.12, bar_color)
    add_text(slide, left + 0.2, top + 0.08, width - 0.4, 0.45,
             value, size=36, color=value_color, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.2, top + 0.55, width - 0.4, 0.3,
             label, size=13, color=GRAY_500, bold=False,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 14: 채택 7개 집중
# ══════════════════════════════════════════════
def build_slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # 1) 제목 + coral 밑줄
    add_header(slide, "7개 가설 채택 — 호스트 변수가 RevPAR을 결정한다")

    # 2) 인사이트 박스
    add_insight_box(
        slide,
        "💡 인사이트:  비모수 검정(p < 0.001) 기반, "
        "호스트 통제 변수(사진·최소숙박·슈퍼호스트)가 RevPAR에 유의미한 영향"
    )

    # 3) 2열 콘텐츠
    LEFT_X = 0.7
    LEFT_W = 5.45
    RIGHT_X = 6.5
    RIGHT_W = 5.8
    CONTENT_Y = 2.05

    # 좌측 컬럼 헤더
    add_text(slide, LEFT_X, CONTENT_Y, 4, 0.25,
             "호스트 관점 (H1–H5)", size=15, color=CHARCOAL, bold=True)
    add_text(slide, LEFT_X + 2.5, CONTENT_Y + 0.03, 3, 0.22,
             "호스트가 통제 가능한 변수", size=11, color=GRAY_500)

    # 우측 컬럼 헤더
    add_text(slide, RIGHT_X, CONTENT_Y, 4, 0.25,
             "자치구 관점 (H6–H7)", size=15, color=CHARCOAL, bold=True)
    add_text(slide, RIGHT_X + 2.5, CONTENT_Y + 0.03, 3, 0.22,
             "시장 구조·입지 변수", size=11, color=GRAY_500)

    # 호스트 가설 5개
    host_hypotheses = [
        ("H1", "슈퍼호스트 → 수익 프리미엄",  "+83.1%"),
        ("H2", "숙소 유형 → 수익 구조 차이",   "전체숙소 2.7배"),
        ("H3", "사진 수 → 수익 (비선형)",      "21~35장 최적"),
        ("H4", "최소숙박일 → 수익 영향",       "2~3박 최적"),
        ("H5", "추가요금 정책 → 수익",         "지역별 차이"),
    ]

    row_gap = 0.54
    rows_start = CONTENT_Y + 0.38

    for i, (h_id, title, result) in enumerate(host_hypotheses):
        ry = rows_start + i * row_gap
        add_hypothesis_row(slide, LEFT_X, ry, LEFT_W,
                           h_id, title, result, CORAL)

    # 자치구 가설 2개
    district_hypotheses = [
        ("H6", "공급 집중 → 수익 하락",   "마포구 리스크"),
        ("H7", "관광지 유형 → 수익",      "여행코스 최고"),
    ]

    for i, (h_id, title, result) in enumerate(district_hypotheses):
        ry = rows_start + i * row_gap
        add_hypothesis_row(slide, RIGHT_X, ry, RIGHT_W,
                           h_id, title, result, CORAL)

    # 우측 빈 공간: 어펜딕스 안내 카드
    note_y = rows_start + 2 * row_gap + 0.1
    add_rect(slide, RIGHT_X, note_y, RIGHT_W, 0.42, GRAY_BG, radius=0.04)
    add_rect(slide, RIGHT_X, note_y + 0.06, 0.035, 0.30, GRAY_400)
    add_text(slide, RIGHT_X + 0.2, note_y + 0.08, RIGHT_W - 0.4, 0.28,
             "H8·H9 부분 채택 → Appendix 참조 (자치구 소표본 n=25)",
             size=12, color=GRAY_500)

    # 4) 하단 요약 카드 2개
    CARD_Y = 5.35
    CARD_H = 0.9
    CARD_W = 5.55
    GAP = 0.3

    add_summary_card(slide, LEFT_X, CARD_Y, CARD_W, CARD_H,
                     "7", "채택", CORAL, CORAL)

    add_summary_card(slide, LEFT_X + CARD_W + GAP, CARD_Y, CARD_W, CARD_H,
                     "p < 0.001", "모든 결과 유의미", TEAL, TEAL)

    # 5) 페이지 번호
    add_text(slide, 12.5, 7.0, 0.5, 0.3, "14",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    return slide


# ══════════════════════════════════════════════
# 어펜딕스: 부분 채택 가설 (H8, H9)
# ══════════════════════════════════════════════
def build_appendix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # 1) 제목
    add_header(slide, "Appendix — 부분 채택 가설 (H8, H9)")

    # 2) 인사이트 박스: 소표본 한계 설명
    add_insight_box(
        slide,
        "💡 자치구 수준(n=25) 소표본 분석으로 통계적 검정력이 제한적이며, "
        "보수적으로 부분 채택 판정",
        height=0.6
    )

    # 3) H8 카드
    card_x = 0.7
    card_w = 11.8
    card_h = 1.6

    h8_y = 2.1
    add_rect(slide, card_x, h8_y, card_w, card_h, GRAY_BG, radius=0.04)
    add_rect(slide, card_x, h8_y + 0.1, 0.035, card_h - 0.2, ORANGE)

    # H8 배지
    add_rect(slide, card_x + 0.2, h8_y + 0.2, 0.5, 0.3, ORANGE, radius=0.12)
    add_text(slide, card_x + 0.2, h8_y + 0.21, 0.5, 0.3, "H8",
             size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, card_x + 0.85, h8_y + 0.18, 4, 0.35,
             "인구 규모 → 수익 관계", size=20, color=CHARCOAL, bold=True)

    add_text(slide, card_x + 0.85, h8_y + 0.55, card_w - 1.5, 0.9,
             "결과:  약한 음의 상관 — 인구가 많은 자치구일수록 RevPAR이 소폭 낮음\n"
             "해석:  인구 규모는 숙소 공급량과 연동되어 경쟁 심화 효과가 나타나지만,\n"
             "         상관계수가 낮아 단독 예측 변수로는 부족 → 부분 채택",
             size=14, color=GRAY_500)

    # △ 부분채택 배지
    add_text(slide, card_x + card_w - 1.5, h8_y + 0.2, 1.3, 0.3,
             "△ 부분채택", size=14, color=ORANGE, bold=True,
             align=PP_ALIGN.RIGHT)

    # 4) H9 카드
    h9_y = 4.0
    add_rect(slide, card_x, h9_y, card_w, card_h, GRAY_BG, radius=0.04)
    add_rect(slide, card_x, h9_y + 0.1, 0.035, card_h - 0.2, ORANGE)

    # H9 배지
    add_rect(slide, card_x + 0.2, h9_y + 0.2, 0.5, 0.3, ORANGE, radius=0.12)
    add_text(slide, card_x + 0.2, h9_y + 0.21, 0.5, 0.3, "H9",
             size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, card_x + 0.85, h9_y + 0.18, 4, 0.35,
             "Dormant 비율 → 시장 건전성", size=20, color=CHARCOAL, bold=True)

    add_text(slide, card_x + 0.85, h9_y + 0.55, card_w - 1.5, 0.9,
             "결과:  음의 상관 — Dormant 비율이 높은 자치구일수록 RevPAR 하락\n"
             "해석:  비활성 리스팅은 시장 포화·경쟁 실패의 신호이나,\n"
             "         n=25 소표본에서 인과 추론은 제한적 → 부분 채택",
             size=14, color=GRAY_500)

    # △ 부분채택 배지
    add_text(slide, card_x + card_w - 1.5, h9_y + 0.2, 1.3, 0.3,
             "△ 부분채택", size=14, color=ORANGE, bold=True,
             align=PP_ALIGN.RIGHT)

    # 5) 페이지 번호
    add_text(slide, 12.5, 7.0, 0.5, 0.3, "A-1",
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)

    return slide


# ── 메인 ──
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_14(prs)
    build_appendix(prs)

    out = "presentation/slide_14_hypothesis.pptx"
    prs.save(out)
    print(f"✓ 저장 완료: {out} (2장: 본 슬라이드 + 어펜딕스)")
