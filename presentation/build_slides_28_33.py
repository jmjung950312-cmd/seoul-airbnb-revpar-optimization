"""
슬라이드 28~33 — 섹션 6 후반 + 클로징 (6장)

디자인 시스템: build_eda_slides.py 패턴 통일
인지 과학 원칙:
  - Miller's Law: 슬라이드당 3~5 청크
  - Peak-End Rule: 31(액션)=Peak, 32(임팩트)=End
  - Progressive Disclosure: 구현 상세 3장→1장 압축
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 디자인 토큰 (build_eda_slides.py 기준) ──
CORAL       = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL    = RGBColor(0x48, 0x48, 0x48)
GRAY_BG     = RGBColor(0xF7, 0xF7, 0xF7)
GRAY_400    = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500    = RGBColor(0x76, 0x76, 0x76)
GRAY_200    = RGBColor(0xE8, 0xE8, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT  = RGBColor(0xE0, 0xF7, 0xF5)
ORANGE      = RGBColor(0xFC, 0x64, 0x2D)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xED)
BLUE        = RGBColor(0x4A, 0x90, 0xD9)
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)

FONT = "Pretendard"
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


# ═══════════════════════════════════════════
# 공통 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title):
    add_text(slide, 0.7, 0.35, 10.5, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL)


def add_page_num(slide, num):
    add_text(slide, 12.5, 7.0, 0.5, 0.3, str(num),
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_kpi_card(slide, left, top, width, height, value, label,
                 value_color=CHARCOAL, bar_color=None, accent=False,
                 value_size=36, label_size=13):
    """KPI 카드: 큰 숫자 + 설명 + 좌측 바"""
    bg = CORAL_LIGHT if accent else GRAY_BG
    add_rect(slide, left, top, width, height, bg, radius=0.04)
    if bar_color:
        add_rect(slide, left, top + 0.06, 0.035, height - 0.12, bar_color)
    add_text(slide, left + 0.2, top + 0.1, width - 0.4, 0.5,
             value, size=value_size, color=value_color, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.2, top + height - 0.4, width - 0.4, 0.3,
             label, size=label_size, color=GRAY_500, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# 슬라이드 28: 요금 시뮬레이터
# ═══════════════════════════════════════════
def build_slide_28(prs):
    s = new_slide(prs)
    add_header(s, "가격을 바꾸면 수익이 어떻게 되나요?")
    add_insight_box(s,
        "💡 인사이트:  군집별 가격탄성도가 다르므로 강남과 관악의 최적 가격이 다릅니다. "
        "시뮬레이터로 즉시 확인 가능")

    # ── 좌측: 가격탄성도 4카드 ──
    LX = 0.7
    LW = 5.45

    add_text(s, LX, 2.1, 3, 0.25,
             "군집별 가격탄성도", size=15, color=CHARCOAL, bold=True)

    elasticities = [
        ("핫플 수익형",       "-0.7", "가격 내성 강함",   CORAL),
        ("프리미엄 비즈니스", "-0.8", "비즈니스 수요 안정", TEAL),
        ("로컬 주거형",       "-0.9", "가격 민감",        ORANGE),
        ("가성비 신흥형",     "-1.1", "매우 가격 민감",   BLUE),
    ]

    for i, (name, val, desc, color) in enumerate(elasticities):
        y = 2.5 + i * 0.72
        add_rect(s, LX, y, LW, 0.62, GRAY_BG, radius=0.04)
        add_rect(s, LX, y + 0.06, 0.035, 0.50, color)
        add_text(s, LX + 0.2, y + 0.1, 2.5, 0.25,
                 name, size=14, color=CHARCOAL, bold=True)
        add_text(s, LX + 0.2, y + 0.34, 2.5, 0.22,
                 desc, size=12, color=GRAY_500)
        add_text(s, LX + LW - 1.2, y + 0.1, 1.0, 0.4,
                 val, size=28, color=color, bold=True, align=PP_ALIGN.RIGHT)

    # ── 우측: 월 손익 + 추천 ──
    RX = 6.5
    RW = 5.8

    add_text(s, RX, 2.1, 3, 0.25,
             "월 손익 계산서", size=15, color=CHARCOAL, bold=True)

    # 손익 구조 카드
    add_rect(s, RX, 2.5, RW, 1.8, GRAY_BG, radius=0.04)
    add_rect(s, RX, 2.56, 0.035, 1.68, CORAL)
    pnl_lines = [
        ("  월 매출", "ADR × 점유율 × 30일"),
        ("- 수수료", "플랫폼 3%"),
        ("- 운영비", "호스트 직접 입력"),
        ("= 월 순이익", "실시간 계산"),
    ]
    for j, (label, note) in enumerate(pnl_lines):
        py = 2.62 + j * 0.42
        is_result = j == 3
        add_text(s, RX + 0.2, py, 2.0, 0.3,
                 label, size=15,
                 color=CORAL if is_result else CHARCOAL,
                 bold=is_result)
        add_text(s, RX + 2.5, py, 3.0, 0.3,
                 note, size=13, color=GRAY_500)

    # 추천 3단계
    add_text(s, RX, 4.55, 3, 0.25,
             "적정 요금 추천", size=15, color=CHARCOAL, bold=True)

    tiers = [
        ("신규", "리뷰 < 10",    "시장 하위 25%",  ORANGE),
        ("안정", "리뷰 10~50",   "시장 평균",      TEAL),
        ("프리미엄", "50+, 슈퍼", "시장 상위 25%", CORAL),
    ]
    for j, (level, cond, rec, color) in enumerate(tiers):
        ty = 4.9 + j * 0.52
        add_rect(s, RX, ty, RW, 0.44, GRAY_BG, radius=0.04)
        add_rect(s, RX, ty + 0.05, 0.035, 0.34, color)
        # 배지
        add_rect(s, RX + 0.2, ty + 0.09, 0.9, 0.26, color, radius=0.12)
        add_text(s, RX + 0.2, ty + 0.10, 0.9, 0.26,
                 level, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, RX + 1.25, ty + 0.08, 1.5, 0.28,
                 cond, size=12, color=GRAY_500)
        add_text(s, RX + RW - 2.0, ty + 0.08, 1.8, 0.28,
                 f"→ {rec}", size=13, color=color, bold=True,
                 align=PP_ALIGN.RIGHT)

    add_page_num(s, 28)
    return s


# ═══════════════════════════════════════════
# 슬라이드 29: 위험 감지 + 자동화 (29+30 병합)
# ═══════════════════════════════════════════
def build_slide_29(prs):
    s = new_slide(prs)
    add_header(s, "24/7 위험 감지 — 5개 규칙이 쉬지 않습니다")
    add_insight_box(s,
        "💡 인사이트:  데이터 변경 시 자동 실행 → 리스크 분류 → "
        "HIGH 발견 시 즉시 이메일 알림. 호스트 개입 없이 24시간 동작",
        height=0.6)

    # ── R1~R5 미니카드 5개 (가로 1행) ──
    rules = [
        ("R1", "유령수익",       "리뷰 0 + 수익 > 1천만", "HIGH",   CORAL),
        ("R2", "점유율-수익\n불일치", "점유율 0% & 수익 > 0", "HIGH", CORAL),
        ("R3", "가격 이상치",    "ADR > 구평균+3σ",       "MEDIUM", ORANGE),
        ("R4", "유령 액티브",    "Active & 수익 < 5백만",  "MEDIUM", ORANGE),
        ("R5", "평점 없는\n고수익", "수익 > 1천만 & 평점 무", "HIGH", CORAL),
    ]

    card_w = 2.2
    card_gap = 0.15
    card_start_x = 0.7
    card_y = 1.95
    card_h = 1.55

    for i, (rid, name, cond, level, color) in enumerate(rules):
        cx = card_start_x + i * (card_w + card_gap)

        # 카드 배경
        add_rect(s, cx, card_y, card_w, card_h, GRAY_BG, radius=0.04)
        # 상단 accent 바
        add_rect(s, cx + 0.1, card_y, card_w - 0.2, 0.035, color)

        # R# 배지 + LEVEL 배지
        add_rect(s, cx + 0.15, card_y + 0.15, 0.42, 0.24, color, radius=0.10)
        add_text(s, cx + 0.15, card_y + 0.16, 0.42, 0.24,
                 rid, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        level_bg = CORAL_LIGHT if level == "HIGH" else ORANGE_LIGHT
        add_rect(s, cx + card_w - 0.85, card_y + 0.15, 0.7, 0.24,
                 level_bg, radius=0.10)
        add_text(s, cx + card_w - 0.85, card_y + 0.16, 0.7, 0.24,
                 level, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)

        # 규칙명
        add_text(s, cx + 0.15, card_y + 0.5, card_w - 0.3, 0.45,
                 name, size=14, color=CHARCOAL, bold=True)

        # 조건 (회색 박스)
        add_rect(s, cx + 0.1, card_y + 1.05, card_w - 0.2, 0.38,
                 RGBColor(0xEE, 0xEE, 0xEE), radius=0.03)
        add_text(s, cx + 0.2, card_y + 1.1, card_w - 0.4, 0.3,
                 cond, size=11, color=GRAY_500)

    # ── 자동화 플로우 (간소화) ──
    flow_y = 3.7
    add_text(s, 0.7, flow_y, 3, 0.25,
             "자동 실행 플로우", size=15, color=CHARCOAL, bold=True)

    steps = [
        ("데이터 변경", "hooks.py", GRAY_BG, CHARCOAL),
        ("R1~R5 검사", "detector.py", GRAY_BG, CHARCOAL),
        ("IQR / Z-score", "scorer.py", GRAY_BG, CHARCOAL),
        ("리스크 분류", "HIGH / MEDIUM", GRAY_BG, CHARCOAL),
        ("이메일 알림", "자동 발송", CORAL_LIGHT, CORAL),
    ]

    step_w = 2.1
    step_gap = 0.22
    step_y = 4.05

    for i, (title, sub, bg, tc) in enumerate(steps):
        sx = 0.7 + i * (step_w + step_gap)
        add_rect(s, sx, step_y, step_w, 0.75, bg, radius=0.04)
        if bg == CORAL_LIGHT:
            add_rect(s, sx, step_y + 0.06, 0.035, 0.63, CORAL)
        add_text(s, sx + 0.15, step_y + 0.12, step_w - 0.3, 0.3,
                 title, size=13, color=tc, bold=True)
        add_text(s, sx + 0.15, step_y + 0.42, step_w - 0.3, 0.25,
                 sub, size=11, color=GRAY_500)

        # 화살표
        if i < len(steps) - 1:
            ax = sx + step_w + 0.02
            add_text(s, ax, step_y + 0.2, 0.2, 0.3,
                     "→", size=18, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 하단 요약 ──
    add_rect(s, 0.7, 5.15, 11.8, 0.55, TEAL_LIGHT, radius=0.04)
    add_rect(s, 0.7, 5.21, 0.035, 0.43, TEAL)
    add_text(s, 1.0, 5.25, 11.0, 0.4,
             "Claude Code Hooks 기반 — 데이터 파일 수정 시 자동 트리거, 중복 알림 7일 필터, "
             "HIGH 리스크(2개+ 규칙) 시에만 이메일 발송",
             size=14, color=TEAL, bold=True)

    # 탐지 실적 카드 2개
    add_kpi_card(s, 0.7, 5.9, 5.7, 0.9,
                 "14건", "전체 플래그 탐지 (R1: 14건, R3: 3건)",
                 CORAL, CORAL)
    add_kpi_card(s, 6.6, 5.9, 5.9, 0.9,
                 "HIGH 1건", "listing_idx=5996 · 송파구 · ₩20.8M",
                 CORAL, CORAL, accent=True)

    add_page_num(s, 29)
    return s


# ═══════════════════════════════════════════
# 슬라이드 30: Multi-Agent AI (간소화)
# ═══════════════════════════════════════════
def build_slide_30(prs):
    s = new_slide(prs)
    add_header(s, "Multi-Agent AI — 분석 자동화 엔진")
    add_insight_box(s,
        "💡 인사이트:  Orchestrator가 6개 에이전트를 순차 호출, "
        "도메인 리서치부터 모델 검증까지 End-to-End 자동화. Enterprise 상품의 핵심",
        height=0.6)

    # ── Orchestrator 바 ──
    orch_y = 1.95
    add_rect(s, 0.7, orch_y, 11.8, 0.55, CORAL, radius=0.04)
    add_text(s, 0.7, orch_y + 0.05, 11.8, 0.2,
             "Central", size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0.7, orch_y + 0.2, 11.8, 0.3,
             "Orchestrator — SOP 기반 순차 호출 · JSON 산출물 관리",
             size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 화살표
    add_text(s, 0.7, 2.55, 11.8, 0.3,
             "↓", size=22, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 6개 에이전트 (3×2 그리드) ──
    agents = [
        ("Domain Research", "도메인 지식 수집",     CORAL),
        ("Hypothesis",      "H1~H9 자동 생성",     ORANGE),
        ("EDA ×3",          "38개 차트 + 통계 검증", TEAL),
        ("Feature Eng.",    "14개 파생변수 설계",    BLUE),
        ("Modeling",        "LightGBM + SHAP",     PURPLE),
        ("Model Review",    "누수 검출 · 품질 평가", RGBColor(0x2C, 0x3E, 0x50)),
    ]

    agent_w = 3.7
    agent_h = 0.8
    gap_x = 0.35
    gap_y = 0.2
    start_x = 0.7
    start_y = 2.85

    for i, (name, desc, color) in enumerate(agents):
        col = i % 3
        row = i // 3
        ax = start_x + col * (agent_w + gap_x)
        ay = start_y + row * (agent_h + gap_y)

        add_rect(s, ax, ay, agent_w, agent_h, GRAY_BG, radius=0.04)
        # 상단 accent 바
        add_rect(s, ax + 0.1, ay, agent_w - 0.2, 0.03, color)
        add_text(s, ax + 0.2, ay + 0.15, agent_w - 0.4, 0.3,
                 name, size=15, color=CHARCOAL, bold=True)
        add_text(s, ax + 0.2, ay + 0.45, agent_w - 0.4, 0.25,
                 desc, size=12, color=GRAY_500)

    # 화살표
    add_text(s, 0.7, 4.72, 11.8, 0.3,
             "↓", size=22, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 산출물 4개 KPI ──
    results = [
        ("9",      "가설 검증",  CORAL),
        ("14",     "파생변수",   TEAL),
        ("38",     "차트",       BLUE),
        ("R² 0.85", "모델 성능", PURPLE),
    ]

    result_w = 2.75
    result_gap = 0.2
    result_y = 5.05

    for i, (val, label, color) in enumerate(results):
        rx = 0.7 + i * (result_w + result_gap)
        add_rect(s, rx, result_y, result_w, 0.95, GRAY_BG, radius=0.04)
        add_rect(s, rx, result_y + 0.06, 0.035, 0.83, color)
        add_text(s, rx + 0.15, result_y + 0.1, result_w - 0.3, 0.5,
                 val, size=32, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, rx + 0.15, result_y + 0.6, result_w - 0.3, 0.25,
                 label, size=13, color=GRAY_500, align=PP_ALIGN.CENTER)

    add_page_num(s, 30)
    return s


# ═══════════════════════════════════════════
# 슬라이드 31: 호스트 액션 Top 3 (PEAK 슬라이드)
# ═══════════════════════════════════════════
def build_slide_31(prs):
    s = new_slide(prs)
    add_header(s, "호스트가 지금 바로 실행할 액션 Top 3")
    add_insight_box(s,
        "💡 인사이트:  SHAP 분석 + 가설 검증 결과를 종합한 우선순위. "
        "대시보드에서 호스트별 맞춤 추천 자동 생성")

    # ── 3개 큰 카드 (가로 배치) ──
    actions = [
        {
            "rank": "1", "medal": "🥇",
            "action": "리뷰 요청\n자동화",
            "evidence": "SHAP 1위 (중요도 0.51)",
            "effect": "+5~10%",
            "effect_sub": "리뷰 10건당 RevPAR",
            "difficulty": "낮음",
            "color": CORAL,
        },
        {
            "rank": "2", "medal": "🥈",
            "action": "슈퍼호스트\n달성",
            "evidence": "H1 검증 완료",
            "effect": "+83.1%",
            "effect_sub": "RevPAR 프리미엄",
            "difficulty": "중간",
            "color": TEAL,
        },
        {
            "rank": "3", "medal": "🥉",
            "action": "최소숙박\n2~3박",
            "evidence": "H4 검증 완료",
            "effect": "최적점",
            "effect_sub": "RevPAR 수익 극대화",
            "difficulty": "낮음",
            "color": BLUE,
        },
    ]

    card_w = 3.7
    card_gap = 0.3
    card_start = 0.7
    card_y = 2.1
    card_h = 3.5

    for i, a in enumerate(actions):
        cx = card_start + i * (card_w + card_gap)
        color = a["color"]

        # 카드 배경
        add_rect(s, cx, card_y, card_w, card_h, GRAY_BG, radius=0.04)
        # 좌측 accent 바
        add_rect(s, cx, card_y + 0.1, 0.04, card_h - 0.2, color)

        # 순위 배지
        add_rect(s, cx + 0.25, card_y + 0.25, 0.55, 0.40, color, radius=0.10)
        add_text(s, cx + 0.25, card_y + 0.28, 0.55, 0.40,
                 a["rank"], size=20, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)

        # 액션명
        add_text(s, cx + 0.95, card_y + 0.25, card_w - 1.2, 0.6,
                 a["action"], size=20, color=CHARCOAL, bold=True)

        # 구분선
        add_rect(s, cx + 0.25, card_y + 0.95, card_w - 0.5, 0.015, GRAY_200)

        # 근거
        add_text(s, cx + 0.25, card_y + 1.1, 1.0, 0.2,
                 "근거", size=11, color=GRAY_500, bold=True)
        add_text(s, cx + 0.25, card_y + 1.35, card_w - 0.5, 0.25,
                 a["evidence"], size=14, color=CHARCOAL)

        # 기대효과 — 큰 숫자
        add_text(s, cx + 0.25, card_y + 1.75, 1.0, 0.2,
                 "기대효과", size=11, color=GRAY_500, bold=True)
        add_text(s, cx + 0.25, card_y + 2.0, card_w - 0.5, 0.55,
                 a["effect"], size=40, color=color, bold=True)
        add_text(s, cx + 0.25, card_y + 2.6, card_w - 0.5, 0.25,
                 a["effect_sub"], size=12, color=GRAY_500)

        # 난이도 배지
        diff_bg = TEAL_LIGHT if a["difficulty"] == "낮음" else ORANGE_LIGHT
        diff_color = TEAL if a["difficulty"] == "낮음" else ORANGE
        add_rect(s, cx + 0.25, card_y + 3.0, 1.2, 0.3, diff_bg, radius=0.12)
        add_text(s, cx + 0.25, card_y + 3.02, 1.2, 0.3,
                 f"난이도: {a['difficulty']}", size=11, color=diff_color,
                 bold=True, align=PP_ALIGN.CENTER)

    # ── 하단: 추가 액션 안내 ──
    add_rect(s, 0.7, 5.85, 11.8, 0.45, GRAY_BG, radius=0.04)
    add_rect(s, 0.7, 5.91, 0.035, 0.33, GRAY_400)
    add_text(s, 1.0, 5.93, 11.0, 0.35,
             "+4개 추가 액션 (사진 21~35장 · 평점 4.85~4.95 · 추가요금 제거 · Entire Home 전환) "
             "→ 대시보드에서 호스트별 맞춤 제공",
             size=13, color=GRAY_500)

    add_page_num(s, 31)
    return s


# ═══════════════════════════════════════════
# 슬라이드 32: 프로젝트 임팩트 (END 슬라이드)
# ═══════════════════════════════════════════
def build_slide_32(prs):
    s = new_slide(prs)
    add_header(s, "데이터로 호스트의 수익을 설계하다")

    # 인사이트 박스 대신 슬로건 서브타이틀
    add_text(s, 0.7, 1.15, 11.8, 0.4,
             "분석 → 모델 → 대시보드 → 자동화 → 액션",
             size=22, color=GRAY_500, align=PP_ALIGN.CENTER)

    # ── 4개 임팩트 KPI (2×2 그리드) ──
    impacts = [
        ("R² = 0.85",     "Dual Model로\nRevPAR의 85% 설명",        CORAL),
        ("5컴포넌트",      "건강점수로\n숙소 상태 종합 진단 (A~F)",      TEAL),
        ("5규칙 24/7",     "위험 감지 엔진이\n쉬지 않고 자동 감시",      BLUE),
        ("End-to-End",    "분석 → 대시보드 →\n자동화 → AI 시스템",    PURPLE),
    ]

    kpi_w = 5.7
    kpi_h = 2.0
    gap_x = 0.4
    gap_y = 0.35
    start_x = 0.7
    start_y = 1.85

    for i, (val, desc, color) in enumerate(impacts):
        col = i % 2
        row = i // 2
        kx = start_x + col * (kpi_w + gap_x)
        ky = start_y + row * (kpi_h + gap_y)

        add_rect(s, kx, ky, kpi_w, kpi_h, GRAY_BG, radius=0.04)
        add_rect(s, kx, ky + 0.1, 0.04, kpi_h - 0.2, color)

        # 큰 숫자
        add_text(s, kx + 0.3, ky + 0.3, kpi_w - 0.6, 0.7,
                 val, size=42, color=color, bold=True)
        # 설명
        add_text(s, kx + 0.3, ky + 1.1, kpi_w - 0.6, 0.7,
                 desc, size=16, color=GRAY_500)

    # ── 하단 마무리 슬로건 ──
    add_rect(s, 0.7, 6.3, 11.8, 0.55, CORAL, radius=0.04)
    add_text(s, 0.7, 6.38, 11.8, 0.4,
             "이 프로젝트는 분석에서 끝나지 않습니다. "
             "호스트의 행동을 바꿉니다.",
             size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_page_num(s, 32)
    return s


# ═══════════════════════════════════════════
# 슬라이드 33: Q&A
# ═══════════════════════════════════════════
def build_slide_33(prs):
    s = new_slide(prs)

    # 중앙 coral 카드
    card_w = 7.0
    card_h = 3.5
    card_x = (13.333 - card_w) / 2
    card_y = (7.5 - card_h) / 2 - 0.3

    add_rect(s, card_x, card_y, card_w, card_h, CORAL, radius=0.15)

    # Q & A 텍스트
    add_text(s, card_x, card_y + 0.6, card_w, 1.0,
             "Q & A", size=56, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, card_x, card_y + 1.8, card_w, 0.5,
             "질문과 의견을 환영합니다", size=20, color=WHITE,
             align=PP_ALIGN.CENTER)

    # 하단 팀 정보
    add_text(s, 0.7, 6.5, 11.8, 0.3,
             "4조 · 데이터분석 부트캠프 · 2026",
             size=14, color=GRAY_400, align=PP_ALIGN.CENTER)

    add_page_num(s, 33)
    return s


# ═══════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_28(prs)
    build_slide_29(prs)
    build_slide_30(prs)
    build_slide_31(prs)
    build_slide_32(prs)
    build_slide_33(prs)

    out = "presentation/slides_28_33.pptx"
    prs.save(out)
    print(f"✓ 저장 완료: {out} (6장: 슬라이드 28~33)")
