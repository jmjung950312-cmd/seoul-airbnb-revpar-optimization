"""
슬라이드 10 — 피처 엔지니어링 (네이티브 python-pptx, 편집 가능)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 색상 ──
CORAL = RGBColor(0xFF, 0x5A, 0x5F)
CORAL_LIGHT_BG = RGBColor(0xFF, 0xF0, 0xF3)
TEAL = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT_BG = RGBColor(0xE0, 0xF7, 0xF5)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
GRAY_100 = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_200 = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xE0, 0x43, 0x47)
DANGER_BG = RGBColor(0xFF, 0xEB, 0xEB)

FONT = "Pretendard"
MONO = "SF Mono"


def add_textbox(slide, left, top, width, height, text,
                font_size=14, color=CHARCOAL, bold=False,
                font_name=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """텍스트 박스 헬퍼"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tb


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """둥근 모서리 사각형"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # 모서리 반경
    shape.adjustments[0] = 0.05
    return shape


def add_left_bar(slide, left, top, height, color, width_pt=4):
    """좌측 컬러 바 (얇은 사각형)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Pt(width_pt), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, bg_color, text_color, width=1.2):
    """배지 (pill shape)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.32)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5  # 완전 둥근 모서리
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_feature_card(slide, left, top, width, height, var_name, desc, note):
    """모델 피처 카드: 배경 + 좌측 바 + 변수명 + 설명"""
    # 카드 배경
    add_rounded_rect(slide, left, top, width, height, GRAY_100)
    # 좌측 coral 바
    add_left_bar(slide, left, top + 0.05, height - 0.10, CORAL)

    # 변수명 (코드 스타일 배지)
    code_bg = add_rounded_rect(slide, left + 0.22, top + 0.18, len(var_name) * 0.095 + 0.28, 0.28, WHITE, CORAL)
    tf = code_bg.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = var_name
    p.font.size = Pt(11)
    p.font.color.rgb = CORAL
    p.font.bold = True
    p.font.name = MONO
    p.alignment = PP_ALIGN.CENTER

    # 설명
    add_textbox(slide, left + 0.22, top + 0.55, width - 0.4, 0.25,
                desc, 13, CHARCOAL, bold=True)
    # 보조 노트
    add_textbox(slide, left + 0.22, top + 0.82, width - 0.4, 0.22,
                note, 11, GRAY_500)


def add_ref_row(slide, top, var_name, desc, note, role, role_color, role_bg,
                strikethrough=False):
    """타겟/참조 변수 행"""
    row_left = 0.8
    row_width = 11.7

    # 행 배경
    add_rounded_rect(slide, row_left, top, row_width, 0.38, GRAY_100)

    # 변수명
    code_w = len(var_name) * 0.088 + 0.3
    code_bg = add_rounded_rect(slide, row_left + 0.2, top + 0.06, code_w, 0.26, WHITE)
    tf = code_bg.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = var_name
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.font.bold = True
    p.font.name = MONO
    p.alignment = PP_ALIGN.CENTER

    # 설명
    add_textbox(slide, 3.2, top + 0.06, 3.0, 0.28, desc, 12, CHARCOAL,
                anchor=MSO_ANCHOR.MIDDLE)
    # 노트
    add_textbox(slide, 6.8, top + 0.06, 3.5, 0.28, note, 11, GRAY_500,
                anchor=MSO_ANCHOR.MIDDLE)

    # 역할 배지
    badge_w = 0.8
    badge = add_badge(slide, row_left + row_width - badge_w - 0.15, top + 0.04,
                      role, role_bg, role_color, width=badge_w)
    if strikethrough:
        for p in badge.text_frame.paragraphs:
            p.font.strikethrough = True


def add_bottom_card(slide, left, top, width, number, label_line1, label_line2,
                    num_color, bar_color):
    """하단 인사이트 카드"""
    # 배경
    bg = add_rounded_rect(slide, left, top, width, 0.75, GRAY_100)
    # 좌측 바
    add_left_bar(slide, left, top + 0.05, 0.65, bar_color)

    # 큰 숫자
    add_textbox(slide, left + 0.25, top + 0.08, 1.2, 0.55,
                number, 32, num_color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # 라벨
    tb = slide.shapes.add_textbox(
        Inches(left + 1.5), Inches(top + 0.12), Inches(width - 1.8), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = label_line1
    p1.font.size = Pt(12)
    p1.font.color.rgb = CHARCOAL
    p1.font.name = FONT

    p2 = tf.add_paragraph()
    p2.text = label_line2
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_500
    p2.font.bold = True
    p2.font.name = FONT


def build_slide_10():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ── 헤더 ──
    add_textbox(slide, 0.8, 0.4, 2, 0.25, "전처리", 12, GRAY_500, bold=True)
    add_textbox(slide, 0.8, 0.65, 10, 0.45,
                "피처 엔지니어링 — 원본 데이터를 분석 가능한 형태로", 26, CHARCOAL, bold=True)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.12), Inches(3.6), Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background()

    # ── 배지 행 ──
    badge_y = 1.3
    add_badge(slide, 0.8, badge_y, "구간화 4개", CORAL_LIGHT_BG, CORAL, 1.05)
    add_badge(slide, 1.95, badge_y, "비율·지표 3개", TEAL_LIGHT_BG, TEAL, 1.2)
    add_badge(slide, 3.25, badge_y, "타겟 변환 1개", RGBColor(0xF0, 0xF0, 0xF0), GRAY_500, 1.15)
    add_textbox(slide, 4.55, badge_y + 0.04, 4, 0.25,
                "원본 42열 → 14개 파생변수 추가 생성", 11, GRAY_500)

    # ── 섹션 1: 구간화 ──
    s1y = 1.72
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(s1y + 0.05), Pt(9), Pt(9))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CORAL
    dot.line.fill.background()
    add_textbox(slide, 1.0, s1y - 0.02, 8, 0.28,
                "구간화 — 연속형 변수를 의미 있는 구간으로 변환", 14, CHARCOAL, bold=True)

    group_cards = [
        ("photos_tier", "사진 수 → 구간 분류", "0~10 / 11~20 / 21~35 / 36+"),
        ("bedrooms_group", "침실 수 → 구간 분류", "0 / 1 / 2 / 3+"),
        ("baths_group", "욕실 수 → 구간 분류", "1 / 1.5 / 2 / 2.5+"),
        ("guests_group", "수용 인원 → 구간 분류", "1~2 / 3~4 / 5~6 / 7+"),
    ]
    cw, ch = 2.78, 1.35
    gx = 0.12
    sx = 0.8
    sy = 2.05
    for i, (var, desc, note) in enumerate(group_cards):
        x = sx + i * (cw + gx)
        add_feature_card(slide, x, sy, cw, ch, var, desc, note)

    # ── 섹션 2: 비율·지표 ──
    s2y = 3.6
    dot2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(s2y + 0.05), Pt(9), Pt(9))
    dot2.fill.solid()
    dot2.fill.fore_color.rgb = TEAL
    dot2.line.fill.background()
    add_textbox(slide, 1.0, s2y - 0.02, 8, 0.28,
                "비율·지표 — 기존 컬럼을 조합하여 새로운 의미 부여", 14, CHARCOAL, bold=True)

    ratio_cards = [
        ("poi_dist_category", "POI 거리 카테고리화", "도보권 / 근거리 / 원거리 등급"),
        ("extra_guest_fee_policy", "추가요금 유무 이진화", "추가인원 요금 있음/없음"),
        ("revpar_trend", "최근 90일 vs 연평균", "수익 방향성 (성장/하락)"),
    ]
    r_sy = s2y + 0.32
    for i, (var, desc, note) in enumerate(ratio_cards):
        x = sx + i * (cw + gx)
        add_feature_card(slide, x, r_sy, cw, ch, var, desc, note)

    # 타겟 변환 카드 (teal)
    tx = sx + 3 * (cw + gx)
    add_rounded_rect(slide, tx, r_sy, cw, ch, TEAL_LIGHT_BG)
    add_left_bar(slide, tx, r_sy + 0.05, ch - 0.10, TEAL)
    code_bg = add_rounded_rect(slide, tx + 0.22, r_sy + 0.22, 1.65, 0.28, WHITE, TEAL)
    tf = code_bg.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "log_ttm_revpar"
    p.font.size = Pt(12)
    p.font.color.rgb = TEAL
    p.font.bold = True
    p.font.name = MONO
    p.alignment = PP_ALIGN.CENTER
    add_textbox(slide, tx + 0.22, r_sy + 0.6, cw - 0.4, 0.25,
                "RevPAR → log1p 변환", 14, CHARCOAL, bold=True)
    add_textbox(slide, tx + 0.22, r_sy + 0.88, cw - 0.4, 0.22,
                "skewness 3.76 → 정규분포화", 12, GRAY_500)

    # ── 인사이트 박스 ──
    iy = 5.25
    add_rounded_rect(slide, 0.8, iy, 11.7, 0.82, CORAL_LIGHT_BG)
    add_left_bar(slide, 0.8, iy + 0.05, 0.72, CORAL)
    add_textbox(slide, 1.15, iy + 0.08, 11, 0.3,
                "왜 이런 변수를 만들었나?", 15, CHARCOAL, bold=True)
    points = [
        "→  사진 21장과 22장의 차이는 무의미 — 의미 있는 구간으로 묶어 패턴 탐지",
        "→  기존 컬럼의 비율·조합으로 숨은 패턴을 드러냄",
        "→  RevPAR 분포가 극단적으로 치우쳐 log 변환으로 분석 안정성 확보",
    ]
    for i, pt in enumerate(points):
        add_textbox(slide, 1.15 + i * 3.8, iy + 0.42, 3.6, 0.35,
                    pt, 12, CHARCOAL)

    # ── 하단 요약 바 ──
    by = 6.25
    add_bottom_card(slide, 0.8, by, 3.7, "14", "파생변수 생성",
                    "원본 42열에서 추가 파생", TEAL, TEAL)
    add_bottom_card(slide, 4.65, by, 3.7, "4+3+1", "구간화 · 비율 · 타겟",
                    "세 가지 유형으로 설계", CORAL, CORAL)
    add_bottom_card(slide, 8.5, by, 3.95, "log1p", "타겟 변환",
                    "skewness 3.76 → 정규분포화", CHARCOAL, GRAY_400)

    # ── 페이지 번호 ──
    add_textbox(slide, 12.3, 7.0, 0.5, 0.3, "10", 12, GRAY_400,
                align=PP_ALIGN.RIGHT)

    return prs


if __name__ == "__main__":
    prs = build_slide_10()
    out = "slide_10_only.pptx"
    prs.save(out)
    print(f"저장 완료: {out}")
