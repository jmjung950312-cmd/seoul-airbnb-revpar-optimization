"""
하이브리드 PPT 빌드 스크립트
- HTML 슬라이드: Playwright로 렌더링 → PNG → 전체 이미지 삽입
- 네이티브 슬라이드: python-pptx 직접 생성
- 차트 슬라이드: 기존 reports/ PNG 삽입
"""

import os
import sys
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 프로젝트 경로
PROJECT_DIR = Path(__file__).parent.parent
PRES_DIR = PROJECT_DIR / "presentation"
RENDERED_DIR = PRES_DIR / "rendered"
HTML_DIR = PRES_DIR / "html_slides"
REPORTS_DIR = PROJECT_DIR / "reports"

# 색상
CORAL = RGBColor(0xFF, 0x5A, 0x5F)
TEAL = RGBColor(0x00, 0xA6, 0x99)
CHARCOAL = RGBColor(0x48, 0x48, 0x48)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_400 = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500 = RGBColor(0x76, 0x76, 0x76)


# ── STEP 1: HTML → PNG 렌더링 ─────────────────────────────

def render_html_slides():
  """모든 HTML 슬라이드를 PNG로 렌더링"""
  from playwright.sync_api import sync_playwright
  import html_templates as tpl

  # HTML 생성 목록: (파일명, HTML 생성 함수)
  html_slides = [
    ("slide_01", tpl.slide_01_cover),
    ("slide_02", tpl.slide_02_toc),
    ("slide_03", lambda: tpl.slide_section("01", "문제정의", "Why · What · How · Result", 3,
      purpose="호스트가 직면한 수익 딜레마를 정의하고 분석 방향을 설정합니다.",
      items=[
        {"text": "시장 현황 진단 (32,061 리스팅)", "color": "var(--coral)"},
        {"text": "RevPAR = ADR × 점유율 구조 분석", "color": "var(--coral)"},
        {"text": "6단계 분석 프레임워크 수립", "color": "var(--teal)"},
        {"text": "기대 결과물 정의", "color": "var(--teal)"},
      ])),
    ("slide_04", tpl.slide_04_why),
    ("slide_05", tpl.slide_05_what),
    ("slide_06", tpl.slide_06_how),
    ("slide_07", tpl.slide_07_result),
    ("slide_08", lambda: tpl.slide_section("02", "데이터 개요·전처리", "", 8,
      purpose="32,061개 리스팅 데이터의 품질 이슈를 해결하고 분석 준비를 완료합니다.",
      items=[
        {"text": "데이터 소개 — 42열, TTM 12개월", "color": "var(--coral)"},
        {"text": "품질 이슈 대응 3건", "color": "var(--coral)"},
        {"text": "피처 엔지니어링 14개 파생변수", "color": "var(--teal)"},
      ])),
    ("slide_09", tpl.slide_09_data_overview),
    ("slide_10", tpl.slide_10_feature_eng),
    ("slide_11", lambda: tpl.slide_section("03", "탐색적 데이터 분석", "호스트 관점 가설 검증", 11,
      purpose="호스트가 통제 가능한 변수 중심으로 9개 가설을 데이터로 검증합니다.",
      items=[
        {"text": "핵심 발견 Top 5", "color": "var(--coral)"},
        {"text": "차트 증거 (2×2 그리드)", "color": "var(--teal)"},
      ])),
    ("slide_12", tpl.slide_12_eda_top5),
    # slide_13: 차트 슬라이드 (기존 PNG 사용) — 아래 별도 처리
    ("slide_14", lambda: tpl.slide_section("04", "통계검증", "", 14,
      purpose="9개 가설에 대한 비모수 통계 검정 결과를 요약합니다.",
      items=[
        {"text": "가설 검증 결과 요약 (7/9 채택)", "color": "var(--coral)"},
        {"text": "p-value < 0.001 수준", "color": "var(--teal)"},
      ])),
    # slide_15: 네이티브 테이블
    ("slide_16", lambda: tpl.slide_section("05", "모델링", "A. Dual Model  ·  B. 건강점수  ·  C. 자치구 군집", 16,
      purpose="호스트 수익을 예측하고 숙소 상태를 진단하는 3개 모델을 구축합니다.",
      items=[
        {"text": "A. Dual Model — ADR × Occupancy 분리 예측", "color": "var(--coral)"},
        {"text": "B. 건강점수 — 5컴포넌트 × A~F 등급", "color": "var(--teal)"},
        {"text": "C. 자치구 군집 — K-Means (k=4)", "color": "#3498DB"},
      ])),
    ("slide_17", tpl.slide_17_dual_model),
    # slide_18: 네이티브 (모델 선택 근거)
    ("slide_19", tpl.slide_19_performance),
    ("slide_20", tpl.slide_20_shap),
    ("slide_21", tpl.slide_21_health_score),
    # slide_22: 네이티브 (군집 모델 설명)
    ("slide_23", tpl.slide_23_cluster_result),
    ("slide_24", lambda: tpl.slide_section("06", "대시보드·자동화", "★ 분석이 끝나는 곳에서, 비즈니스가 시작된다", 24,
      purpose="분석 결과를 대시보드와 자동화 시스템으로 비즈니스에 연결합니다.",
      items=[
        {"text": "비즈니스 모델 (Free → Pro)", "color": "var(--coral)"},
        {"text": "위험 감지 엔진 R1~R5", "color": "var(--coral)"},
        {"text": "자동화 플로우 & Multi-Agent AI", "color": "var(--teal)"},
        {"text": "호스트 액션 가이드", "color": "var(--teal)"},
      ])),
    ("slide_25", tpl.slide_25_biz_model),
    # slide_26~28: 네이티브/스크린샷
    ("slide_29", tpl.slide_29_risk_rules),
    ("slide_30", tpl.slide_30_automation_flow),
    ("slide_31", tpl.slide_31_multi_agent),
    ("slide_32", tpl.slide_32_action_guide),
    ("slide_33", tpl.slide_33_impact),
    ("slide_34", tpl.slide_34_thankyou),
  ]

  # HTML 파일 생성
  for name, gen_func in html_slides:
    html_path = HTML_DIR / f"{name}.html"
    html_path.write_text(gen_func(), encoding="utf-8")
    print(f"  ✓ HTML 생성: {name}.html")

  # Playwright로 PNG 렌더링
  print("\n📸 HTML → PNG 렌더링 시작...")
  with sync_playwright() as p:
    browser = p.chromium.launch()
    page = browser.new_page(viewport={"width": 1920, "height": 1080})

    for name, _ in html_slides:
      html_path = HTML_DIR / f"{name}.html"
      png_path = RENDERED_DIR / f"{name}.png"
      page.goto(f"file://{html_path.absolute()}")
      page.wait_for_load_state("networkidle")
      # 폰트 로딩 대기
      page.wait_for_timeout(1500)
      page.screenshot(path=str(png_path), full_page=False)
      print(f"  ✓ 렌더링: {name}.png")

    browser.close()

  print(f"\n✅ HTML 슬라이드 {len(html_slides)}장 렌더링 완료")
  return [name for name, _ in html_slides]


# ── STEP 2: PPTX 빌드 ────────────────────────────────────

def add_image_slide(prs, img_path):
  """전체 이미지 슬라이드 추가 (16:9 전체 채움)"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
  # 16:9 = 13.333" × 7.5"
  slide.shapes.add_picture(
    str(img_path),
    left=Emu(0), top=Emu(0),
    width=Emu(12192000),  # 13.333 inches
    height=Emu(6858000),  # 7.5 inches
  )
  return slide


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 font_color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Pretendard"):
  """텍스트 박스 추가 헬퍼"""
  txBox = slide.shapes.add_textbox(
    Inches(left), Inches(top), Inches(width), Inches(height)
  )
  tf = txBox.text_frame
  tf.word_wrap = True
  p = tf.paragraphs[0]
  p.text = text
  p.font.size = Pt(font_size)
  p.font.color.rgb = font_color
  p.font.bold = bold
  p.font.name = font_name
  p.alignment = alignment
  return txBox


def add_native_slide_15(prs):
  """슬라이드 15: 가설 검증 결과 요약 (네이티브 테이블)"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  # 배경 라이트
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  # 제목
  add_text_box(slide, 0.8, 0.5, 4, 0.3, "통계검증", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.8, 8, 0.5, "가설 검증 결과 요약", 32, CHARCOAL, True)

  # 테이블
  rows, cols = 8, 5
  tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(8.5), Inches(4.5)).table

  headers = ["가설", "내용", "검증 방법", "결과", ""]
  data = [
    ["H1", "슈퍼호스트 → RevPAR ↑", "Mann-Whitney U", "✅ +83.1%", ""],
    ["H2", "객실유형 → RevPAR 구조", "Kruskal-Wallis", "✅ Entire 2.7배", ""],
    ["H3", "사진 → RevPAR (비선형)", "Spearman + 구간 분석", "✅ 21~35장", ""],
    ["H4", "최소숙박 → RevPAR", "구간별 중위값 비교", "✅ 2~3박", ""],
    ["H6", "공급 집중 → RevPAR ↓", "Pearson + 군집", "✅ 마포 리스크", ""],
    ["H7", "POI 유형 → RevPAR", "Kruskal-Wallis", "✅ 여행코스 최고", ""],
    ["", "", "", "", ""],
  ]

  # 헤더
  for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
      p.font.size = Pt(13)
      p.font.bold = True
      p.font.color.rgb = GRAY_500
      p.font.name = "Noto Sans KR"

  # 데이터
  for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
      cell = tbl.cell(r + 1, c)
      cell.text = val
      for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.name = "Noto Sans KR"
        p.font.color.rgb = CHARCOAL

  # 우측 핵심 메시지 박스
  add_text_box(slide, 10.0, 1.8, 2.8, 0.4, "핵심 메시지", 13, GRAY_500, True)
  add_text_box(slide, 10.0, 2.3, 2.8, 0.5, "9개 가설 중 7개 채택\n2개 부분 채택", 16, CHARCOAL, True)
  add_text_box(slide, 10.0, 3.2, 2.8, 0.8, "모든 검증은 비모수 검정 기반\n(RevPAR 분포 비정규)\np-value < 0.001 수준", 13, GRAY_500)

  # 페이지 번호
  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "15", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def add_native_slide_18(prs):
  """슬라이드 18: 모델 선택 근거 & 누수 방지"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  add_text_box(slide, 0.8, 0.5, 4, 0.3, "모델링 A", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.8, 10, 0.5, "호스트 모델 — 선택 근거 & 누수 방지", 32, CHARCOAL, True)

  # 좌측: 모델 선택 근거
  add_text_box(slide, 0.8, 1.7, 5, 0.3, "모델 선택 근거", 16, CHARCOAL, True)
  models_text = (
    "• LightGBM: 비선형 관계 포착, 범주형 네이티브, SHAP 해석 가능\n"
    "• Random Forest: 앙상블 비교군\n"
    "• Ridge: 선형 베이스라인\n"
    "• DummyRegressor: 중위값 기준선 (모델 개선 증명용)\n\n"
    "→ 단순 성능이 아닌 해석 가능성 + 배포 용이성 기준"
  )
  add_text_box(slide, 0.8, 2.2, 5.5, 3.0, models_text, 14, CHARCOAL)

  # 우측: 누수 방지
  add_text_box(slide, 7.0, 1.7, 5, 0.3, "데이터 누수 방지", 16, CORAL, True)
  leakage_text = (
    "CRITICAL: revpar_vs_district_median — 타겟 직접 파생\n"
    "HIGH: revpar_trend — 타겟 구성요소\n"
    "MEDIUM: price_efficiency — ADR/점유율 비율\n"
    "+ TTM 수익 관련 6개 변수 전체 제외\n\n"
    "→ Model Review 에이전트가 자동 누수 검출"
  )
  add_text_box(slide, 7.0, 2.2, 5.5, 3.0, leakage_text, 14, CHARCOAL)

  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "18", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def add_native_slide_22(prs):
  """슬라이드 22: 자치구 군집 모델"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  add_text_box(slide, 0.8, 0.5, 4, 0.3, "모델링 C", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.8, 10, 0.5, "자치구 군집 모델 — K-Means (k=4)", 32, CHARCOAL, True)

  # 좌측
  add_text_box(slide, 0.8, 1.7, 5.5, 0.3, "모델링 접근", 16, CHARCOAL, True)
  approach_text = (
    "왜 군집 모델인가?\n"
    "• 자치구 단위 (n=25)는 소표본\n"
    "• 지도학습 대신 비지도학습 선택\n\n"
    "방법: K-Means (k=4), 표준화 후 군집화\n"
    "보조: Ridge + LeaveOneOut CV\n"
    "입력: RevPAR 중위값, Dormant 비율,\n"
    "슈퍼호스트 비율, 공급량, POI 분포"
  )
  add_text_box(slide, 0.8, 2.2, 5.5, 3.5, approach_text, 14, CHARCOAL)

  # 우측
  add_text_box(slide, 7.0, 1.7, 5.5, 0.3, "모델 선택 근거", 16, TEAL, True)
  reason_text = (
    "• Elbow Method → k=4에서 변곡점\n"
    "• Silhouette Score 확인\n\n"
    "호스트 예측 모델 (개별) +\n"
    "자치구 군집 모델 (시장) = 상호 보완\n\n"
    "→ 대시보드 '시장 유형 & 전략' 탭 연동\n"
    "→ 가격탄성도 클러스터별 차등 적용"
  )
  add_text_box(slide, 7.0, 2.2, 5.5, 3.5, reason_text, 14, CHARCOAL)

  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "22", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def add_chart_slide_13(prs):
  """슬라이드 13: 2×2 차트 그리드 (기존 PNG 활용)"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  add_text_box(slide, 0.8, 0.3, 4, 0.3, "EDA", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.6, 10, 0.5, "호스트 관점 EDA — 차트 증거", 28, CHARCOAL, True)

  # 차트 이미지 배치 (2×2)
  charts = [
    ("fig_01b_A_adr_occ_quadrant.png", "ADR-점유율 4분면"),
    ("fig_01d_XE_reviews_superhost_revpar.png", "슈퍼호스트 RevPAR"),
    ("fig_01b_B_reviews_revpar.png", "사진/리뷰 vs RevPAR"),
    ("fig_01d_XA_rating_photos_heatmap.png", "평점-사진 히트맵"),
  ]

  positions = [
    (0.5, 1.3, 6.0, 3.0),   # 좌상
    (6.8, 1.3, 6.0, 3.0),   # 우상
    (0.5, 4.4, 6.0, 3.0),   # 좌하
    (6.8, 4.4, 6.0, 3.0),   # 우하
  ]

  for (fname, label), (l, t, w, h) in zip(charts, positions):
    img_path = REPORTS_DIR / fname
    if img_path.exists():
      slide.shapes.add_picture(
        str(img_path),
        Inches(l), Inches(t), Inches(w), Inches(h)
      )
    else:
      add_text_box(slide, l, t, w, h, f"[{label}]", 16, GRAY_400)

  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "13", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def add_native_slide_26(prs):
  """슬라이드 26: 무료 대시보드 (이메일 + 이상신호)"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  add_text_box(slide, 0.8, 0.5, 6, 0.3, "대시보드 · 자동화", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.8, 10, 0.5, "무료 버전 — 이메일 자동화 + 이상 신호 3종", 28, CHARCOAL, True)

  # 분석 플로우
  add_text_box(slide, 0.8, 1.6, 6, 0.3, "분석 플로우", 14, GRAY_500, True)
  flow = "호스트 URL 입력 → 자동 분석 → 이상 신호 3종 감지 → 리포트 생성 → 이메일 발송"
  add_text_box(slide, 0.8, 2.0, 7.5, 0.5, flow, 14, CHARCOAL)

  # 이상 신호 3종
  add_text_box(slide, 0.8, 2.7, 6, 0.3, "이상 신호 3종", 14, GRAY_500, True)
  signals = (
    "① Z-score 이상치: ADR/RevPAR |Z| > 2.0\n"
    "② 추세 이탈: L90D vs TTM/4 하락률 > -20%\n"
    "③ STL 잔차 이상: 시계열 분해 잔차 > 2σ"
  )
  add_text_box(slide, 0.8, 3.1, 7.5, 1.5, signals, 14, CHARCOAL)

  # 시장 포지셔닝
  add_text_box(slide, 0.8, 4.7, 6, 0.3, "시장 포지셔닝 (4사분면)", 14, GRAY_500, True)
  quad = "물량형 (고점유·저가) / 고수익형 (고점유·고가) / 침체형 (저점유·저가) / 고가위험형 (저점유·고가)"
  add_text_box(slide, 0.8, 5.1, 7.5, 0.7, quad, 13, CHARCOAL)

  # 우측: 이메일 리포트 설명
  add_text_box(slide, 9.0, 1.6, 4, 0.3, "이메일 리포트 구조", 14, GRAY_500, True)
  email_desc = (
    "상단 (공개):\n"
    "  • 시장 포지션 카드\n"
    "  • 이상 신호 진단\n\n"
    "하단 (blur 처리):\n"
    "  • 건강점수 상세\n"
    "  • 최적화 가이드\n\n"
    "CTA: \"유료 대시보드에서\n"
    "       확인하세요\""
  )
  add_text_box(slide, 9.0, 2.0, 3.5, 4.0, email_desc, 14, CHARCOAL)

  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "26", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def add_native_slide_27(prs):
  """슬라이드 27: 유료 핵심요약 탭"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  add_text_box(slide, 0.8, 0.5, 6, 0.3, "대시보드 · 자동화", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.8, 10, 0.5, "유료 핵심요약 탭 — 사분면 + 건강점수 + ML 진단", 26, CHARCOAL, True)

  # 4개 기능 설명
  features = [
    ("① 사분면 포지셔닝 차트", "ADR 백분위 vs 점유율 백분위\n물량형/고수익형/침체형/고가위험형\n호스트 현재 위치 빨간 점 표시"),
    ("② 건강점수 게이지", "0~100점 게이지 + A~F 등급 뱃지\n5컴포넌트 레이더 차트"),
    ("③ KPI 3종 카드", "예측 RevPAR\n시장 대비 위치 (백분위)\n개선 잠재력 (%)"),
    ("④ ML 시장 진단", "Dual Model 결과\n예측 ADR vs 현재 ADR\nprice_gap 시각화, 점유율 예측"),
  ]

  for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.6 + row * 2.8
    add_text_box(slide, x, y, 5.5, 0.3, title, 16, CORAL, True)
    add_text_box(slide, x, y + 0.4, 5.5, 2.0, desc, 14, CHARCOAL)

  add_text_box(slide, 0.8, 6.8, 12, 0.5, "모델링의 Dual Model + 건강점수 + 군집이 이 화면 하나에 통합", 15, GRAY_500, alignment=PP_ALIGN.CENTER)
  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "27", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def add_native_slide_28(prs):
  """슬라이드 28: 요금 시뮬레이션 & 월 손익"""
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  bg = slide.background
  fill = bg.fill
  fill.solid()
  fill.fore_color.rgb = LIGHT_BG

  add_text_box(slide, 0.8, 0.5, 6, 0.3, "대시보드 · 자동화", 14, CORAL, True)
  add_text_box(slide, 0.8, 0.8, 10, 0.5, "요금 시뮬레이션 & 월 손익 계산서", 28, CHARCOAL, True)

  # 좌측: 시뮬레이터
  add_text_box(slide, 0.8, 1.6, 6, 0.3, "요금 변경 시뮬레이터", 16, CHARCOAL, True)
  sim_text = (
    "슬라이더: 현재 요금 대비 -30% ~ +50%\n\n"
    "가격탄성도 (클러스터별 차등):\n"
    "  • 핫플 수익형: -0.7 (가격 내성 강함)\n"
    "  • 프리미엄 비즈니스: -0.8\n"
    "  • 로컬 주거형: -0.9 (가격 민감)\n"
    "  • 가성비 신흥형: -1.1 (매우 민감)\n\n"
    "실시간 결과: 예상 점유율 / RevPAR / 월 순이익"
  )
  add_text_box(slide, 0.8, 2.0, 5.5, 3.5, sim_text, 14, CHARCOAL)

  # 우측: 월 손익 + 추천
  add_text_box(slide, 7.0, 1.6, 6, 0.3, "월 손익 계산서", 16, CHARCOAL, True)
  pnl_text = (
    "월 매출 (ADR × 점유율 × 30일)\n"
    "- 플랫폼 수수료 (3%)\n"
    "- 월 운영비 (호스트 입력)\n"
    "= 월 순이익\n\n"
    "적정 요금 추천 3단계:\n"
    "  • 신규 (리뷰 <10): 시장 하위 25%\n"
    "  • 안정 (10~50): 시장 평균\n"
    "  • 프리미엄 (50+, 슈퍼): 상위 25%"
  )
  add_text_box(slide, 7.0, 2.0, 5.5, 3.5, pnl_text, 14, CHARCOAL)

  add_text_box(slide, 12.5, 7.0, 0.5, 0.3, "28", 14, GRAY_400, alignment=PP_ALIGN.RIGHT)
  return slide


def build_pptx(rendered_slides):
  """최종 PPTX 빌드"""
  prs = Presentation()
  prs.slide_width = Emu(12192000)   # 13.333"
  prs.slide_height = Emu(6858000)   # 7.5"

  # 슬라이드 순서 정의
  # HTML 렌더링 슬라이드: 이미지로 삽입
  # 네이티브 슬라이드: python-pptx로 직접 생성
  # 차트 슬라이드: 기존 PNG 삽입

  slide_order = [
    # 오프닝
    ("html", "slide_01"),   # 1. 표지
    ("html", "slide_02"),   # 2. 목차
    # ① 문제정의
    ("html", "slide_03"),   # 3. 섹션
    ("html", "slide_04"),   # 4. WHY
    ("html", "slide_05"),   # 5. WHAT
    ("html", "slide_06"),   # 6. HOW
    ("html", "slide_07"),   # 7. RESULT
    # ② 데이터
    ("html", "slide_08"),   # 8. 섹션
    ("html", "slide_09"),   # 9. 데이터 소개
    ("html", "slide_10"),   # 10. 피처 엔지니어링
    # ③ EDA
    ("html", "slide_11"),   # 11. 섹션
    ("html", "slide_12"),   # 12. 핵심 발견 Top 5
    ("chart", "slide_13"),  # 13. 차트 증거
    # ④ 통계검증
    ("html", "slide_14"),   # 14. 섹션
    ("native", "slide_15"), # 15. 가설 검증 요약
    # ⑤ 모델링
    ("html", "slide_16"),   # 16. 섹션
    ("html", "slide_17"),   # 17. Dual Model
    ("native", "slide_18"), # 18. 모델 선택 근거
    ("html", "slide_19"),   # 19. 성능 비교
    ("html", "slide_20"),   # 20. SHAP
    ("html", "slide_21"),   # 21. 건강점수
    ("native", "slide_22"), # 22. 군집 모델
    ("html", "slide_23"),   # 23. 군집 결과
    # ⑥ 대시보드·자동화
    ("html", "slide_24"),   # 24. 섹션
    ("html", "slide_25"),   # 25. 비즈니스 모델
    ("native", "slide_26"), # 26. 무료 대시보드
    ("native", "slide_27"), # 27. 유료 핵심요약
    ("native", "slide_28"), # 28. 요금 시뮬레이션
    ("html", "slide_29"),   # 29. 위험 감지
    ("html", "slide_30"),   # 30. 자동화 플로우
    ("html", "slide_31"),   # 31. Multi-Agent
    ("html", "slide_32"),   # 32. 액션 가이드
    # 클로징
    ("html", "slide_33"),   # 33. 임팩트 요약
    ("html", "slide_34"),   # 34. Q&A
  ]

  for i, (stype, name) in enumerate(slide_order):
    slide_num = i + 1
    print(f"  [{slide_num:02d}/34] ", end="")

    if stype == "html":
      img_path = RENDERED_DIR / f"{name}.png"
      if img_path.exists():
        add_image_slide(prs, img_path)
        print(f"✓ HTML이미지 — {name}")
      else:
        print(f"✗ 이미지 없음: {name}.png")
        prs.slides.add_slide(prs.slide_layouts[6])

    elif stype == "native":
      if name == "slide_15":
        add_native_slide_15(prs)
      elif name == "slide_18":
        add_native_slide_18(prs)
      elif name == "slide_22":
        add_native_slide_22(prs)
      elif name == "slide_26":
        add_native_slide_26(prs)
      elif name == "slide_27":
        add_native_slide_27(prs)
      elif name == "slide_28":
        add_native_slide_28(prs)
      print(f"✓ 네이티브 — {name}")

    elif stype == "chart":
      add_chart_slide_13(prs)
      print(f"✓ 차트 — {name}")

  # 저장
  output_path = PRES_DIR / "서울_에어비앤비_RevPAR최적화.pptx"
  prs.save(str(output_path))
  print(f"\n📄 PPTX 저장: {output_path}")
  return output_path


# ── MAIN ──────────────────────────────────────────────────

def main():
  print("=" * 60)
  print("🎨 하이브리드 PPT 빌드 시작")
  print("=" * 60)

  # 디렉토리 준비
  RENDERED_DIR.mkdir(parents=True, exist_ok=True)
  HTML_DIR.mkdir(parents=True, exist_ok=True)

  # STEP 1: HTML 렌더링
  print("\n📝 STEP 1: HTML 슬라이드 생성 & 렌더링")
  print("-" * 40)
  rendered = render_html_slides()

  # STEP 2: PPTX 빌드
  print("\n🔨 STEP 2: PPTX 빌드")
  print("-" * 40)
  output = build_pptx(rendered)

  print("\n" + "=" * 60)
  print(f"✅ 완료! 34장 하이브리드 PPT 생성")
  print(f"   📄 {output}")
  print(f"   HTML 렌더링: {len(rendered)}장")
  print(f"   네이티브: 6장")
  print(f"   차트: 1장")
  print("=" * 60)


if __name__ == "__main__":
  main()
