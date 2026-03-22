"""
슬라이드 28~37 — 섹션 6 + 마무리 (10장, 기존 7장 + 신규 3장)

개선된 흐름:
  p28: 비즈니스 모델
  p29: [NEW] 자동화 개요 — 3가지 자동화
  p30: [NEW] 무료 자동화 — 이메일 진단
  p31: 위험 감지 R1~R5 (기존 p29)
  p32: [NEW] 대시보드 전환 — 브릿지
  p33: 대시보드 핵심 기능 (기존 p30)
  p34: 자동화 + AI 에이전트 (기존 p31)
  p35: 임팩트 요약 (기존 p32)
  p36: 한계점 + 로드맵 (기존 p33)
  p37: Q&A (기존 p34)
"""
from pathlib import Path
import importlib.util

# build_slides_28_34.py 의 함수들을 임포트
spec_old = importlib.util.spec_from_file_location(
    "old", Path(__file__).parent / "build_slides_28_34.py")
old = importlib.util.module_from_spec(spec_old)
spec_old.loader.exec_module(old)

# build_new_transition_slides.py 의 함수들을 임포트
spec_new = importlib.util.spec_from_file_location(
    "new", Path(__file__).parent / "build_new_transition_slides.py")
new = importlib.util.module_from_spec(spec_new)
spec_new.loader.exec_module(new)

from pptx import Presentation
from pptx.util import Emu


if __name__ == "__main__":
    print("=" * 55)
    print("슬라이드 28~37 통합 빌드 (기존 7장 + 신규 3장 = 10장)")
    print("=" * 55)

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    # p28: 비즈니스 모델
    print("\n[01/10] p28: 비즈니스 모델")
    old.build_slide_28(prs)

    # p29: [NEW] 자동화 개요
    print("[02/10] p29: [NEW] 자동화 개요")
    new.build_slide_new_a(prs, page_num=29)

    # p30: [NEW] 무료 자동화 이메일
    print("[03/10] p30: [NEW] 무료 자동화 이메일")
    new.build_slide_new_b(prs, page_num=30)

    # p31: 위험 감지 R1~R5 (기존 p29, 번호 변경)
    print("[04/10] p31: 위험 감지 R1~R5")
    s = old.build_slide_29(prs)
    # 페이지 번호 패치: 29 → 31
    for shape in s.shapes:
        if hasattr(shape, 'text') and shape.text.strip() == '29':
            shape.text_frame.paragraphs[0].runs[0].text = '31'
            break

    # p32: [NEW] 대시보드 전환
    print("[05/10] p32: [NEW] 대시보드 전환")
    new.build_slide_new_c(prs, page_num=32)

    # p33: 대시보드 핵심 기능 (기존 p30, 번호 변경)
    print("[06/10] p33: 대시보드 핵심 기능")
    s = old.build_slide_30(prs)
    for shape in s.shapes:
        if hasattr(shape, 'text') and shape.text.strip() == '30':
            shape.text_frame.paragraphs[0].runs[0].text = '33'
            break

    # p34: 자동화 + AI 에이전트 (기존 p31, 번호 변경)
    print("[07/10] p34: 자동화 + AI 에이전트")
    s = old.build_slide_31(prs)
    for shape in s.shapes:
        if hasattr(shape, 'text') and shape.text.strip() == '31':
            shape.text_frame.paragraphs[0].runs[0].text = '34'
            break

    # p35: 임팩트 요약 (기존 p32, 번호 변경)
    print("[08/10] p35: 임팩트 요약")
    s = old.build_slide_32(prs)
    for shape in s.shapes:
        if hasattr(shape, 'text') and shape.text.strip() == '32':
            shape.text_frame.paragraphs[0].runs[0].text = '35'
            break

    # p36: 한계점 + 로드맵 (기존 p33, 번호 변경)
    print("[09/10] p36: 한계점 + 로드맵")
    s = old.build_slide_33(prs)
    for shape in s.shapes:
        if hasattr(shape, 'text') and shape.text.strip() == '33':
            shape.text_frame.paragraphs[0].runs[0].text = '36'
            break

    # p37: Q&A (기존 p34)
    print("[10/10] p37: Q&A")
    old.build_slide_34(prs)

    out = Path(__file__).parent / "slides_28_37.pptx"
    prs.save(str(out))
    print(f"\n✅ 저장 완료: {out}")
    print("   10장: 슬라이드 28~37 (기존 7장 + 신규 3장)")
    print("\n개선된 흐름:")
    print("  p28 비즈니스 모델")
    print("  p29 [NEW] 자동화 개요")
    print("  p30 [NEW] 무료 자동화 이메일")
    print("  p31 위험 감지 R1~R5")
    print("  p32 [NEW] 대시보드 전환")
    print("  p33 대시보드 핵심 기능")
    print("  p34 자동화 + AI 에이전트")
    print("  p35 임팩트 요약")
    print("  p36 한계점 + 로드맵")
    print("  p37 Q&A")
