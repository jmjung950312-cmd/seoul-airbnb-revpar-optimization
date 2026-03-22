"""
미니 디바이더(브릿지) 슬라이드 3장 생성
- 큰 섹션 디바이더와 차별화: 흰 배경 + 가벼운 텍스트
- 삽입 위치: 슬라이드 13→14 사이, (선택) 18→19 사이
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import json

# 디자인 시스템 로드
with open("design_system.json", "r") as f:
    ds = json.load(f)

CORAL = RGBColor(0xFF, 0x38, 0x5C)
TEAL = RGBColor(0x00, 0xA6, 0x99)
DARK = RGBColor(0x48, 0x48, 0x48)
GREY = RGBColor(0x76, 0x76, 0x76)
LIGHT_PINK = RGBColor(0xFF, 0xF0, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_mini_divider(prs, top_text, main_text, sub_text, accent_color=CORAL):
    """
    미니 브릿지 슬라이드: 중앙 정렬, 가벼운 톤
    - top_text: 작은 상단 레이블 (예: "DEEP DIVE")
    - main_text: 메인 브릿지 문구
    - sub_text: 부가 설명
    """
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # 상단 얇은 악센트 라인 (좌→우 전체 폭)
    line_h = Pt(3)
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(1.8),
        SLIDE_W, line_h
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    # 상단 레이블 (작은 글씨, 악센트 색상)
    tf_top = slide.shapes.add_textbox(
        Inches(0), Inches(2.2), SLIDE_W, Inches(0.5)
    ).text_frame
    tf_top.word_wrap = True
    p = tf_top.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = top_text
    run.font.size = Pt(14)
    run.font.color.rgb = accent_color
    run.font.bold = True

    # 메인 텍스트 (큰 글씨, 진한 색)
    tf_main = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.9), Inches(10.333), Inches(1.5)
    ).text_frame
    tf_main.word_wrap = True
    p = tf_main.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = main_text
    run.font.size = Pt(36)
    run.font.color.rgb = DARK
    run.font.bold = True

    # 하단 부가 설명 (중간 글씨, 회색)
    tf_sub = slide.shapes.add_textbox(
        Inches(2), Inches(4.6), Inches(9.333), Inches(0.8)
    ).text_frame
    tf_sub.word_wrap = True
    p = tf_sub.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sub_text
    run.font.size = Pt(16)
    run.font.color.rgb = GREY
    run.font.bold = False

    # 하단 얇은 악센트 라인
    line2 = slide.shapes.add_shape(
        1,
        Inches(5.5), Inches(5.6),
        Inches(2.333), Pt(2)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = accent_color
    line2.line.fill.background()

    return slide


# ─── 슬라이드 1: 가설 개요 → 상세 분석 전환 (13→14 사이) ───
add_mini_divider(
    prs,
    top_text="DEEP DIVE",
    main_text="각 가설을 데이터로 검증합니다",
    sub_text="5가지 핵심 요인을 하나씩 살펴봅니다  →",
    accent_color=CORAL,
)

# ─── 슬라이드 2: EDA → 통계 검증 전환 (18→19 사이) ───
add_mini_divider(
    prs,
    top_text="FROM PATTERN TO PROOF",
    main_text="눈에 보인 패턴, 정말 유의미한가?",
    sub_text="EDA에서 발견한 차이를 비모수 검정으로 검증합니다",
    accent_color=TEAL,
)

# ─── 슬라이드 3 (선택): Section 01 RESULT 전 브릿지 (6→7 사이) ───
add_mini_divider(
    prs,
    top_text="PREVIEW",
    main_text="이 분석이 만들어낸 결과물",
    sub_text="분석 과정은 이후 섹션에서 상세히 다룹니다",
    accent_color=CORAL,
)

out_path = "mini_dividers.pptx"
prs.save(out_path)
print(f"✅ 미니 디바이더 3장 저장: {out_path}")
