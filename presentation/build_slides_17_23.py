"""
슬라이드 17~23 — 섹션 5 모델링 (7장)
PDF 기준 페이지 번호 사용

디자인 시스템: build_slides_28_34.py 패턴 통일
레이아웃:
  p17: Layout E (플로우) — Dual Model 아키텍처
  p18: Layout F (좌우 분할) — 모델 선택 & 누수 방지
  p19: Layout G+C (테이블 + KPI) — 성능 비교
  p20: Layout F (좌우 분할) — SHAP → 비즈니스 액션
  p21: Layout F (좌우 분할) — 건강점수 5컴포넌트
  p22: Layout F (좌우 분할) — 자치구 군집 모델
  p23: Layout C (4카드 그리드) — 자치구 군집 결과
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 디자인 토큰 ──
CORAL       = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL    = RGBColor(0x48, 0x48, 0x48)
GRAY_BG     = RGBColor(0xF7, 0xF7, 0xF7)
GRAY_100    = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_200    = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_400    = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500    = RGBColor(0x76, 0x76, 0x76)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT  = RGBColor(0xE0, 0xF7, 0xF5)
ORANGE      = RGBColor(0xFC, 0x64, 0x2D)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xED)
BLUE        = RGBColor(0x4A, 0x90, 0xD9)
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)
PURPLE_LIGHT = RGBColor(0xF3, 0xE8, 0xFA)

FONT = "Pretendard"
SLIDE_W = Emu(12192000)  # 13.333"
SLIDE_H = Emu(6858000)   # 7.5"


# ═══════════════════════════════════════════
# 공통 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title):
    add_text(slide, 0.7, 0.35, 10.5, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL)


def add_page_num(slide, num):
    add_text(slide, 12.5, 7.0, 0.5, 0.3, str(num),
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_kpi_card(slide, left, top, width, height, value, label,
                 value_color=CHARCOAL, bar_color=None, accent=False,
                 value_size=36, label_size=13):
    bg = CORAL_LIGHT if accent else GRAY_BG
    add_rect(slide, left, top, width, height, bg, radius=0.04)
    if bar_color:
        add_rect(slide, left, top + 0.06, 0.035, height - 0.12, bar_color)
    add_text(slide, left + 0.2, top + 0.1, width - 0.4, 0.5,
             value, size=value_size, color=value_color, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.2, top + height - 0.4, width - 0.4, 0.3,
             label, size=label_size, color=GRAY_500, align=PP_ALIGN.CENTER)


def add_circle(slide, left, top, diameter, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_shadow_card(slide, left, top, width, height, fill_color=GRAY_BG,
                    border_color=None, radius=0.04, shadow_offset=0.04):
    add_rect(slide, left + shadow_offset, top + shadow_offset,
             width, height, GRAY_200, radius=radius)
    return add_rect(slide, left, top, width, height, fill_color,
                    border_color=border_color, radius=radius)


# ═══════════════════════════════════════════
# p17: Dual Model 아키텍처 — Layout E (플로우)
# ═══════════════════════════════════════════
def build_slide_17(prs):
    s = new_slide(prs)
    add_header(s, "Dual Model 아키텍처 — ADR × Occupancy 분리 예측")
    add_insight_box(s,
        "보충 모델링: 메인 모델(단일 RevPAR 예측)에 이어, "
        "ADR과 점유율을 분리하여 호스트가 각각 조종 가능하도록 설계")

    # ── 4-step 가로 플로우 ──
    steps = [
        {
            "label": "Model A",
            "sub": "ADR 예측",
            "desc": "위치 · 규모 · POI",
            "bg": CHARCOAL,
            "accent": CORAL,
            "text_color": WHITE,
        },
        {
            "label": "Bridge",
            "sub": "price_gap",
            "desc": "현재 ADR\n− 예측 ADR",
            "bg": TEAL_LIGHT,
            "accent": TEAL,
            "text_color": CHARCOAL,
        },
        {
            "label": "Model B",
            "sub": "Occupancy 예측",
            "desc": "예약정책 · 리뷰\n평점 + price_gap",
            "bg": TEAL_LIGHT,
            "accent": TEAL,
            "text_color": CHARCOAL,
        },
        {
            "label": "Calibration",
            "sub": "Isotonic",
            "desc": "단조증가 보정\n→ 최종 RevPAR",
            "bg": GRAY_BG,
            "accent": GRAY_400,
            "text_color": CHARCOAL,
        },
    ]

    card_w = 2.55
    card_h = 3.2
    card_gap = 0.35
    total_w = 4 * card_w + 3 * card_gap
    start_x = (13.333 - total_w) / 2
    card_y = 2.2

    for i, step in enumerate(steps):
        cx = start_x + i * (card_w + card_gap)

        # 카드 배경
        add_rect(s, cx, card_y, card_w, card_h, step["bg"], radius=0.06)

        # 상단 accent 바
        add_rect(s, cx, card_y, card_w, 0.06, step["accent"], radius=0.06)

        # 스텝 번호 원형 배지
        badge_d = 0.4
        badge_x = cx + (card_w - badge_d) / 2
        add_circle(s, badge_x, card_y + 0.25, badge_d, step["accent"])
        add_text(s, badge_x, card_y + 0.27, badge_d, badge_d,
                 str(i + 1), size=18, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 레이블
        add_text(s, cx + 0.15, card_y + 0.8, card_w - 0.3, 0.35,
                 step["label"], size=20, color=step["text_color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 서브레이블
        add_text(s, cx + 0.15, card_y + 1.15, card_w - 0.3, 0.3,
                 step["sub"], size=15, color=step["accent"] if step["bg"] != CHARCOAL else CORAL,
                 bold=True, align=PP_ALIGN.CENTER)

        # 구분선
        line_color = GRAY_400 if step["bg"] == CHARCOAL else GRAY_200
        add_rect(s, cx + 0.3, card_y + 1.6, card_w - 0.6, 0.015, line_color)

        # 설명
        add_text(s, cx + 0.2, card_y + 1.8, card_w - 0.4, 1.0,
                 step["desc"], size=14, color=step["text_color"],
                 align=PP_ALIGN.CENTER)

        # 화살표 (마지막 제외)
        if i < 3:
            ax = cx + card_w + 0.02
            add_text(s, ax, card_y + card_h / 2 - 0.2, 0.3, 0.4,
                     "→", size=28, color=GRAY_400, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

    # ── 하단 핵심 메시지 바 ──
    bar_y = card_y + card_h + 0.4
    add_rect(s, 0.7, bar_y, 11.8, 0.6, CORAL_LIGHT, radius=0.04)
    add_rect(s, 0.7, bar_y, 0.04, 0.6, CORAL)
    add_text(s, 1.0, bar_y + 0.1, 11.2, 0.4,
             "메인 모델(p18~19)과 별도로, ADR·점유율을 분리 예측하여 "
             "호스트 액션을 구체화하는 보충 모델",
             size=16, color=CHARCOAL, bold=True)

    add_page_num(s, 17)


# ═══════════════════════════════════════════
# p18: 모델 선택 & 누수 방지 — Layout F (좌우 분할)
# ═══════════════════════════════════════════
def build_slide_18(prs):
    s = new_slide(prs)
    add_header(s, "호스트 모델 — 선택 근거 & 누수 방지")
    add_insight_box(s,
        "성능만 추구하지 않았습니다. "
        "TreeSHAP 정확 해석 · 비선형 포착 · 경량 배포 — 세 기준으로 선택했습니다")

    content_y = 2.1
    left_x = 0.7
    right_x = 6.95
    col_w = 5.85

    # ═══ 좌측: 모델 선택 ═══
    add_text(s, left_x, content_y, col_w, 0.35,
             "모델 선택 근거", size=20, color=CHARCOAL, bold=True)

    models = [
        {"name": "LightGBM", "desc": "TreeSHAP 정확 해석 · 비선형 포착 · 경량 배포",
         "badge": "→ 선택", "bg": CORAL_LIGHT, "accent": CORAL, "badge_bg": CORAL},
        {"name": "Random Forest", "desc": "TreeSHAP 가능하나 트리 수백 개 → 느림",
         "badge": None, "bg": GRAY_BG, "accent": GRAY_400, "badge_bg": None},
        {"name": "Ridge", "desc": "KernelSHAP만 가능 (근사치) · 선형 가정 한계",
         "badge": None, "bg": GRAY_BG, "accent": GRAY_400, "badge_bg": None},
        {"name": "DummyRegressor", "desc": "중위값 기준선",
         "badge": None, "bg": GRAY_BG, "accent": GRAY_400, "badge_bg": None},
    ]

    for i, m in enumerate(models):
        cy = content_y + 0.5 + i * 0.95
        add_rect(s, left_x, cy, col_w, 0.8, m["bg"], radius=0.04)
        add_rect(s, left_x, cy + 0.08, 0.04, 0.64, m["accent"])
        add_text(s, left_x + 0.25, cy + 0.1, 2.5, 0.3,
                 m["name"], size=16, color=CHARCOAL, bold=True)
        add_text(s, left_x + 0.25, cy + 0.42, 3.8, 0.3,
                 m["desc"], size=13, color=GRAY_500)
        if m["badge"]:
            add_rect(s, left_x + col_w - 1.3, cy + 0.2, 1.1, 0.4,
                     m["badge_bg"], radius=0.06)
            add_text(s, left_x + col_w - 1.3, cy + 0.22, 1.1, 0.36,
                     m["badge"], size=14, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 하단 강조
    add_text(s, left_x, content_y + 4.35, col_w, 0.3,
             "SHAP은 model-agnostic이나, TreeSHAP은 트리 모델 전용 → 정확 + 고속",
             size=12, color=GRAY_500, bold=True)

    # ═══ 우측: 누수 방지 ═══
    add_text(s, right_x, content_y, col_w, 0.35,
             "데이터 누수 방지", size=20, color=CHARCOAL, bold=True)

    # 왜 누수를 제거하는가 — 이유 설명
    why_y = content_y + 0.45
    add_rect(s, right_x, why_y, col_w, 0.6, GRAY_BG, radius=0.04)
    add_rect(s, right_x, why_y + 0.06, 0.04, 0.48, CHARCOAL)
    add_text(s, right_x + 0.2, why_y + 0.05, col_w - 0.4, 0.5,
             "실전 배포 시 쓸 수 없는 R² 0.95보다,\n"
             "실전에서 재현되는 R² 0.85가 신뢰할 수 있는 모델",
             size=12, color=CHARCOAL)

    leaks = [
        {"level": "CRITICAL", "var": "revpar_vs_district_median",
         "reason": "타겟 직접 파생",
         "bg": CORAL_LIGHT, "accent": CORAL, "label_bg": CORAL},
        {"level": "HIGH", "var": "revpar_trend, review_rate",
         "reason": "타겟 구성요소 포함",
         "bg": ORANGE_LIGHT, "accent": ORANGE, "label_bg": ORANGE},
        {"level": "MEDIUM", "var": "price_efficiency",
         "reason": "ADR / Occupancy 간접 누수",
         "bg": BLUE_LIGHT, "accent": BLUE, "label_bg": BLUE},
    ]

    for i, lk in enumerate(leaks):
        cy = content_y + 1.2 + i * 0.82
        add_rect(s, right_x, cy, col_w, 0.7, lk["bg"], radius=0.04)
        add_rect(s, right_x, cy + 0.06, 0.04, 0.58, lk["accent"])

        # 레벨 배지
        add_rect(s, right_x + 0.2, cy + 0.08, 1.2, 0.3,
                 lk["label_bg"], radius=0.06)
        add_text(s, right_x + 0.2, cy + 0.08, 1.2, 0.3,
                 lk["level"], size=12, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 변수명
        add_text(s, right_x + 1.55, cy + 0.06, col_w - 1.8, 0.3,
                 lk["var"], size=13, color=CHARCOAL, bold=True)

        # 누수 사유
        add_text(s, right_x + 1.55, cy + 0.38, col_w - 1.8, 0.25,
                 lk["reason"], size=11, color=GRAY_500)

    # TTM 관련 제외 변수
    ttm_y = content_y + 3.7
    add_rect(s, right_x, ttm_y, col_w, 0.75, GRAY_BG, radius=0.04)
    add_text(s, right_x + 0.2, ttm_y + 0.06, col_w - 0.4, 0.25,
             "제외 (TTM/L90D 직접 구성요소)", size=12, color=GRAY_500, bold=True)
    excluded = ("ttm_occupancy · ttm_avg_rate · ttm_revenue\n"
                "l90d_revpar · l90d_occupancy · l90d_avg_rate")
    add_text(s, right_x + 0.2, ttm_y + 0.32, col_w - 0.4, 0.4,
             excluded, size=11, color=CHARCOAL)

    # 자동 검출 배지
    add_rect(s, right_x, ttm_y + 0.9, col_w, 0.4, TEAL_LIGHT, radius=0.04)
    add_text(s, right_x + 0.2, ttm_y + 0.92, col_w - 0.4, 0.35,
             "Model Review 에이전트가 자동 검출",
             size=13, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, 18)


# ═══════════════════════════════════════════
# p19: 성능 비교 — Layout G+C (테이블 + KPI)
# ═══════════════════════════════════════════
def build_slide_19(prs):
    s = new_slide(prs)
    add_header(s, "호스트 모델 — 성능 비교")
    add_insight_box(s,
        "호스트 통제 변수만으로 RevPAR의 85%를 예측할 수 있습니다")

    # ── 테이블 (네이티브 shape 기반) ──
    table_x = 0.7
    table_y = 2.15
    col_widths = [3.0, 2.2, 2.2, 2.2, 2.2]  # 합계 11.8
    row_h = 0.55
    headers = ["모델", "CV R²", "Test R²", "Test MAE", "Test MAPE"]

    # 헤더 행
    cx = table_x
    for j, hdr in enumerate(headers):
        add_rect(s, cx, table_y, col_widths[j], row_h, CHARCOAL)
        add_text(s, cx + 0.1, table_y + 0.05, col_widths[j] - 0.2, row_h - 0.1,
                 hdr, size=14, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_widths[j]

    # 데이터 행
    rows = [
        {"model": "LightGBM", "cv": "0.847", "test": "0.85",
         "mae": "₩19,554", "mape": "32.6%", "highlight": True},
        {"model": "Random Forest", "cv": "0.744", "test": "0.75",
         "mae": "₩24,470", "mape": "45.6%", "highlight": False},
        {"model": "Ridge", "cv": "0.579", "test": "0.57",
         "mae": "₩27,595", "mape": "70.4%", "highlight": False},
        {"model": "Baseline (Dummy)", "cv": "-0.025", "test": "—",
         "mae": "—", "mape": "—", "highlight": False},
    ]

    for i, row in enumerate(rows):
        ry = table_y + row_h + i * row_h
        bg = CORAL_LIGHT if row["highlight"] else (GRAY_100 if i % 2 == 0 else WHITE)
        vals = [row["model"], row["cv"], row["test"], row["mae"], row["mape"]]

        cx = table_x
        for j, val in enumerate(vals):
            add_rect(s, cx, ry, col_widths[j], row_h, bg,
                     border_color=GRAY_200)
            txt_color = CORAL if row["highlight"] and j == 0 else CHARCOAL
            txt_bold = row["highlight"] and j in (0, 2)
            add_text(s, cx + 0.1, ry + 0.05, col_widths[j] - 0.2, row_h - 0.1,
                     val, size=14, color=txt_color, bold=txt_bold,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_widths[j]

    # ── 하단 KPI 3개 (shadow 카드) ──
    kpi_y = 5.15
    kpi_w = 3.73
    kpi_h = 1.5
    kpi_gap = 0.3
    kpi_start = 0.7

    kpis = [
        {"value": "85%", "label": "호스트 통제 변수만으로\nRevPAR 설명",
         "color": CORAL, "bar": CORAL},
        {"value": "+0.875", "label": "Baseline 대비 R² 개선\n→ 모델 가치 증명",
         "color": TEAL, "bar": TEAL},
        {"value": "₩19,554", "label": "평균 예측 오차\n→ 실무 활용 수준",
         "color": CHARCOAL, "bar": CHARCOAL},
    ]

    for i, k in enumerate(kpis):
        kx = kpi_start + i * (kpi_w + kpi_gap)
        add_shadow_card(s, kx, kpi_y, kpi_w, kpi_h, GRAY_BG, radius=0.04)
        add_rect(s, kx, kpi_y + 0.06, 0.04, kpi_h - 0.12, k["bar"])
        add_text(s, kx + 0.2, kpi_y + 0.15, kpi_w - 0.4, 0.6,
                 k["value"], size=36, color=k["color"], bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, kx + 0.2, kpi_y + 0.85, kpi_w - 0.4, 0.5,
                 k["label"], size=13, color=GRAY_500,
                 align=PP_ALIGN.CENTER)

    add_page_num(s, 19)


# ═══════════════════════════════════════════
# p20: SHAP → 비즈니스 액션 — Layout F (좌우 분할)
# ═══════════════════════════════════════════
def build_slide_20(prs):
    s = new_slide(prs)

    # ── 섹션 브레드크럼 ──
    add_circle(s, 0.7, 0.3, 0.22, CORAL)
    add_text(s, 1.0, 0.3, 5.0, 0.25,
             "SECTION 05 — 모델링 성능", size=12, color=GRAY_500)

    # ── 제목 (브레드크럼 아래) ──
    add_text(s, 0.7, 0.65, 10.5, 0.65,
             "SHAP 특성 중요도 분석", size=43, color=CHARCOAL, bold=True)

    # ── 인사이트 박스 (CORAL 배경 + 흰색 텍스트) ──
    insight_y = 1.45
    insight_h = 0.95
    add_rect(s, 0.7, insight_y, 11.8, insight_h, CORAL, radius=0.06)
    add_text(s, 1.0, insight_y + 0.12, 11.2, 0.3,
             "핵심 인사이트 — 실행 우선순위",
             size=16, color=WHITE, bold=True)
    add_text(s, 1.0, insight_y + 0.45, 11.2, 0.45,
             "SHAP 분석 결과, 리뷰 수(0.508)가 압도적 1순위입니다. "
             "호스트는 리뷰 유도 → 평점 4.8+ 유지 → 최소숙박 2~3박 → "
             "사진 21~35장 순서로 집중하여 RevPAR를 단계적으로 개선할 수 있습니다.",
             size=13, color=WHITE)

    # ═══════════════════════════════════════
    # 좌측: SHAP 중요도 바차트 (카드 컨테이너)
    # ═══════════════════════════════════════
    panel_x = 0.7
    panel_y = 2.65
    panel_w = 7.4
    panel_h = 4.3

    # 카드 컨테이너
    add_rect(s, panel_x, panel_y, panel_w, panel_h, WHITE,
             border_color=GRAY_200, radius=0.05)

    # 소제목 + TEAL accent 바
    add_text(s, panel_x + 0.35, panel_y + 0.2, panel_w - 0.7, 0.35,
             "SHAP 중요도 (호스트 통제 변수)",
             size=16, color=CORAL, bold=True)
    add_rect(s, panel_x + 0.35, panel_y + 0.55, panel_w - 0.7, 0.03, TEAL)

    # ── 바차트 ──
    features = [
        {"name": "리뷰 수",        "value": 0.508},
        {"name": "평점",           "value": 0.120},
        {"name": "최소숙박",        "value": 0.079},
        {"name": "사진 수",        "value": 0.071},
        {"name": "POI 거리",       "value": 0.054},
    ]

    # 그라데이션 핑크 (진 → 연)
    bar_colors = [
        CORAL,
        RGBColor(0xFF, 0x5A, 0x7A),
        RGBColor(0xFF, 0x7C, 0x98),
        RGBColor(0xFF, 0xA0, 0xB5),
        RGBColor(0xFF, 0xC0, 0xD0),
    ]

    label_x = panel_x + 0.3
    label_w = 1.8
    bar_x = panel_x + 2.2
    bar_max_w = 4.5
    bar_h = 0.42
    bar_gap = 0.22
    chart_y = panel_y + 0.8
    max_val = 0.60  # 축 최대값

    for i, f in enumerate(features):
        by = chart_y + i * (bar_h + bar_gap)
        bw = (f["value"] / max_val) * bar_max_w

        # 피처명 (우측 정렬)
        add_text(s, label_x, by, label_w, bar_h,
                 f["name"], size=14, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

        # 바
        add_rect(s, bar_x, by, bw, bar_h, bar_colors[i], radius=0.04)

    # ── 축 눈금 (0.00 ~ 0.60) ──
    axis_y = chart_y + 5 * (bar_h + bar_gap) - 0.1
    ticks = [0.00, 0.10, 0.20, 0.30, 0.40, 0.50, 0.60]
    for t in ticks:
        tx = bar_x + (t / max_val) * bar_max_w
        add_text(s, tx - 0.25, axis_y, 0.5, 0.25,
                 f"{t:.2f}", size=10, color=GRAY_400,
                 align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════
    # 우측: 호스트 액션 우선순위
    # ═══════════════════════════════════════
    right_x = 8.35
    right_w = 4.15
    right_y = panel_y

    # 카드 컨테이너
    add_rect(s, right_x, right_y, right_w, panel_h, WHITE,
             border_color=GRAY_200, radius=0.05)

    # 소제목 + TEAL accent 바
    add_text(s, right_x + 0.3, right_y + 0.2, right_w - 0.6, 0.35,
             "호스트 액션 우선순위",
             size=16, color=CORAL, bold=True)
    add_rect(s, right_x + 0.3, right_y + 0.55, right_w - 0.6, 0.03, TEAL)

    actions = [
        {"num": "1", "action": "리뷰 유도 메시지 발송",
         "detail": "중요도: 0.508 | 체크아웃 후 리뷰 리마인드",
         "hero": True},
        {"num": "2", "action": "평점 4.8 이상 유지",
         "detail": "중요도: 0.120 | 체크인 안내 · 문제 즉시 보상",
         "hero": False},
        {"num": "3", "action": "최소숙박 2~3박 설정",
         "detail": "중요도: 0.079 | RevPAR 최적 구간",
         "hero": False},
        {"num": "4", "action": "사진 21~35장 업로드",
         "detail": "중요도: 0.071 | 최적 구간 초과 시 한계효용 체감",
         "hero": False},
    ]

    card_start_y = right_y + 0.75
    card_gap = 0.15

    for i, a in enumerate(actions):
        # 1번: CORAL 배경 + 큰 카드 / 나머지: CORAL_LIGHT 배경 + 작은 카드
        if a["hero"]:
            ch = 0.85
            cy = card_start_y
            card_bg = CORAL
            num_bg = WHITE
            num_color = CORAL
            title_color = WHITE
            detail_color = WHITE
        else:
            ch = 0.7
            cy = card_start_y + 0.85 + card_gap + (i - 1) * (0.7 + card_gap)
            card_bg = CORAL_LIGHT
            num_bg = WHITE
            num_color = CORAL
            title_color = CHARCOAL
            detail_color = GRAY_500

        # 카드 배경
        add_rect(s, right_x + 0.2, cy, right_w - 0.4, ch,
                 card_bg, radius=0.05)

        # 넘버 원형
        num_d = 0.45
        add_circle(s, right_x + 0.35, cy + (ch - num_d) / 2,
                   num_d, num_bg)
        add_text(s, right_x + 0.35, cy + (ch - num_d) / 2,
                 num_d, num_d,
                 a["num"], size=20, color=num_color, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 액션명
        add_text(s, right_x + 0.95, cy + 0.08, right_w - 1.3, 0.3,
                 a["action"], size=14, color=title_color, bold=True)

        # 상세
        add_text(s, right_x + 0.95, cy + 0.38, right_w - 1.3, 0.3,
                 a["detail"], size=10, color=detail_color)

    add_page_num(s, 20)


# ═══════════════════════════════════════════
# p21: 건강점수 5컴포넌트 — Layout F (좌우 분할)
# ═══════════════════════════════════════════
def build_slide_21(prs):
    s = new_slide(prs)
    add_header(s, "숙소 건강점수 — 5컴포넌트 × A~F 등급")
    add_insight_box(s,
        "호스트가 '지금 어디가 부족한가'를 한눈에 파악하는 종합 진단 도구")

    content_y = 2.1
    left_x = 0.7
    right_x = 6.95
    col_w = 5.85

    # ═══ 좌측: 5컴포넌트 ═══
    add_text(s, left_x, content_y, col_w, 0.35,
             "건강점수 (0~100)", size=20, color=CHARCOAL, bold=True)

    components = [
        {"name": "리뷰 신호", "weight": "20%",
         "desc": "리뷰수 + 평점 백분위", "color": CORAL},
        {"name": "사진 품질", "weight": "20%",
         "desc": "23~35장 = 100점", "color": TEAL},
        {"name": "예약 정책", "weight": "20%",
         "desc": "즉시예약 + 최소박 + 추가요금", "color": ORANGE},
        {"name": "위치", "weight": "20%",
         "desc": "POI 근접도 역백분위", "color": BLUE},
        {"name": "숙소 구성", "weight": "20%",
         "desc": "침실 · 욕실 백분위", "color": PURPLE},
    ]

    for i, c in enumerate(components):
        cy = content_y + 0.55 + i * 0.78
        add_rect(s, left_x, cy, col_w, 0.65, GRAY_BG, radius=0.04)

        # 컬러 도트
        dot_d = 0.35
        add_circle(s, left_x + 0.15, cy + 0.15, dot_d, c["color"])
        add_text(s, left_x + 0.15, cy + 0.15, dot_d, dot_d,
                 str(i + 1), size=14, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 이름 + 가중치
        add_text(s, left_x + 0.65, cy + 0.08, 2.5, 0.25,
                 f"{c['name']} ({c['weight']})",
                 size=15, color=CHARCOAL, bold=True)

        # 설명
        add_text(s, left_x + 0.65, cy + 0.35, col_w - 1.0, 0.25,
                 c["desc"], size=12, color=GRAY_500)

    # ═══ 우측: 등급 매트릭스 ═══
    add_text(s, right_x, content_y, col_w, 0.35,
             "등급 매트릭스", size=20, color=CHARCOAL, bold=True)

    grades = [
        {"grade": "A", "range": "≥ 80", "label": "상위 20%",
         "color": TEAL, "bg": TEAL_LIGHT},
        {"grade": "B", "range": "≥ 60", "label": "양호",
         "color": BLUE, "bg": BLUE_LIGHT},
        {"grade": "C", "range": "≥ 40", "label": "보통",
         "color": ORANGE, "bg": ORANGE_LIGHT},
        {"grade": "D", "range": "≥ 20", "label": "취약",
         "color": CORAL, "bg": CORAL_LIGHT},
        {"grade": "F", "range": "< 20", "label": "위험",
         "color": CHARCOAL, "bg": GRAY_BG},
    ]

    for i, g in enumerate(grades):
        gy = content_y + 0.55 + i * 0.68
        add_rect(s, right_x, gy, col_w, 0.55, g["bg"], radius=0.04)

        # 등급 배지
        badge_d = 0.4
        add_circle(s, right_x + 0.15, gy + 0.07, badge_d, g["color"])
        add_text(s, right_x + 0.15, gy + 0.07, badge_d, badge_d,
                 g["grade"], size=18, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 범위 + 라벨
        add_text(s, right_x + 0.7, gy + 0.08, 1.0, 0.4,
                 g["range"], size=16, color=CHARCOAL, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, right_x + 1.8, gy + 0.08, 3.5, 0.4,
                 g["label"], size=14, color=GRAY_500,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ── 하단 예시 박스 ──
    example_y = content_y + 4.2
    add_rect(s, 0.7, example_y, 11.8, 0.65, GRAY_BG, radius=0.04)
    add_rect(s, 0.7, example_y, 0.04, 0.65, TEAL)
    add_text(s, 1.0, example_y + 0.1, 11.2, 0.45,
             "예시: 마포구 원룸, 리뷰 12건, 평점 4.6, 즉시예약 미활성 → "
             "건강점수 C (41점) — 예약정책 · 숙소구성 취약",
             size=15, color=CHARCOAL)

    add_page_num(s, 21)


# ═══════════════════════════════════════════
# p22: 자치구 군집 모델 — Layout F (좌우 분할)
# ═══════════════════════════════════════════
def build_slide_22(prs):
    s = new_slide(prs)
    add_header(s, "자치구 군집 모델 — K-Means (k=4)")
    add_insight_box(s,
        "개별 호스트 예측은 LightGBM, 시장 포지셔닝은 K-Means. "
        "두 모델이 서로 다른 의사결정 지원")

    content_y = 2.1
    left_x = 0.7
    right_x = 6.95
    col_w = 5.85

    # ═══ 좌측: 왜 군집 모델인가? ═══
    add_text(s, left_x, content_y, col_w, 0.35,
             "왜 군집 모델인가?", size=20, color=CHARCOAL, bold=True)

    reasons = [
        {"num": "1", "text": "자치구 단위 (n=25) → 소표본에 적합한 군집화", "color": CORAL},
        {"num": "2", "text": "시장 유형 분류 → 호스트 모델과 상호 보완", "color": TEAL},
    ]

    for i, r in enumerate(reasons):
        ry = content_y + 0.55 + i * 0.78
        add_rect(s, left_x, ry, col_w, 0.65, GRAY_BG, radius=0.04)
        badge_d = 0.4
        add_circle(s, left_x + 0.15, ry + 0.12, badge_d, r["color"])
        add_text(s, left_x + 0.15, ry + 0.12, badge_d, badge_d,
                 r["num"], size=16, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left_x + 0.7, ry + 0.08, col_w - 1.0, 0.5,
                 r["text"], size=15, color=CHARCOAL,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 방법 배지
    badge_y = content_y + 2.95
    add_rect(s, left_x, badge_y, 2.6, 0.45, CHARCOAL, radius=0.06)
    add_text(s, left_x, badge_y + 0.02, 2.6, 0.4,
             "K-Means (k=4) · 표준화", size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, left_x + 2.8, badge_y, 3.0, 0.45, GRAY_BG, radius=0.06)
    add_text(s, left_x + 2.8, badge_y + 0.02, 3.0, 0.4,
             "StandardScaler 정규화", size=14, color=GRAY_500, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ═══ 우측: 모델 선택 근거 ═══
    add_text(s, right_x, content_y, col_w, 0.35,
             "모델 선택 근거", size=20, color=CHARCOAL, bold=True)

    # Elbow Method 카드
    ey = content_y + 0.55
    add_rect(s, right_x, ey, col_w, 0.85, GRAY_BG, radius=0.04)
    add_rect(s, right_x, ey + 0.08, 0.04, 0.69, TEAL)
    add_text(s, right_x + 0.25, ey + 0.1, col_w - 0.5, 0.3,
             "Elbow Method", size=16, color=CHARCOAL, bold=True)
    add_text(s, right_x + 0.25, ey + 0.42, col_w - 0.5, 0.35,
             "k=2~8 범위 탐색 → k=4에서 뚜렷한 elbow point",
             size=13, color=GRAY_500)

    # Silhouette Score 카드
    sy = ey + 1.0
    add_rect(s, right_x, sy, col_w, 0.85, GRAY_BG, radius=0.04)
    add_rect(s, right_x, sy + 0.08, 0.04, 0.69, ORANGE)
    add_text(s, right_x + 0.25, sy + 0.1, col_w - 0.5, 0.3,
             "Silhouette Score 확인", size=16, color=CHARCOAL, bold=True)
    add_text(s, right_x + 0.25, sy + 0.42, col_w - 0.5, 0.35,
             "군집 간 분리도 검증 → 안정적 군집 확인",
             size=13, color=GRAY_500)

    # 입력 피처 카드
    fy = sy + 1.0
    add_rect(s, right_x, fy, col_w, 1.1, GRAY_BG, radius=0.04)
    add_rect(s, right_x, fy + 0.08, 0.04, 0.94, BLUE)
    add_text(s, right_x + 0.25, fy + 0.1, col_w - 0.5, 0.3,
             "입력 피처", size=16, color=CHARCOAL, bold=True)
    features_text = ("median_revpar_ao · dormant_ratio\n"
                     "superhost_rate · total_listings")
    add_text(s, right_x + 0.25, fy + 0.42, col_w - 0.5, 0.6,
             features_text, size=13, color=GRAY_500)

    # ── 하단 연결 배지 ──
    link_y = content_y + 4.3
    add_rect(s, 0.7, link_y, 11.8, 0.5, TEAL_LIGHT, radius=0.04)
    add_text(s, 1.0, link_y + 0.05, 11.2, 0.4,
             "→ 대시보드 '시장 유형 & 전략' 탭으로 연결",
             size=16, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, 22)


# ═══════════════════════════════════════════
# p23: 자치구 군집 결과 — Layout C (4카드 그리드)
# ═══════════════════════════════════════════
def build_slide_23(prs):
    s = new_slide(prs)
    add_header(s, "자치구 군집 결과 — 4대 시장 전략")
    add_insight_box(s,
        "K-Means (k=4) · 25개 자치구 군집화")

    # ── 4카드 가로 배치 ──
    clusters = [
        {
            "id": "C1", "name": "핫플 수익형",
            "revpar": "₩63,366", "districts": "마포구",
            "strategy": "품질 밀도 강화\n공급 과잉 경계",
            "dormant": "Dormant 50.8%",
            "accent": CORAL, "bg": CORAL_LIGHT,
        },
        {
            "id": "C2", "name": "프리미엄 비즈니스",
            "revpar": "₩53,052", "districts": "강남 · 종로 · 용산",
            "strategy": "슈퍼호스트\n전환 인센티브",
            "dormant": "Dormant 52.1%",
            "accent": TEAL, "bg": TEAL_LIGHT,
        },
        {
            "id": "C3", "name": "로컬 주거형",
            "revpar": "₩35,881", "districts": "13개 자치구",
            "strategy": "사진 · 평점\n표준 가이드",
            "dormant": "Dormant 55.4%",
            "accent": BLUE, "bg": BLUE_LIGHT,
        },
        {
            "id": "C4", "name": "가성비 신흥형",
            "revpar": "₩26,413", "districts": "동대문 · 관악 등",
            "strategy": "비활성 정리\nDormant 집중 관리",
            "dormant": "Dormant 63.3%",
            "accent": ORANGE, "bg": ORANGE_LIGHT,
        },
    ]

    card_w = 2.78
    card_h = 4.3
    card_gap = 0.2
    total_w = 4 * card_w + 3 * card_gap
    start_x = (13.333 - total_w) / 2
    card_y = 2.1

    for i, cl in enumerate(clusters):
        cx = start_x + i * (card_w + card_gap)

        # 카드 배경
        add_shadow_card(s, cx, card_y, card_w, card_h, WHITE,
                        border_color=GRAY_200, radius=0.05)

        # 상단 accent 바
        add_rect(s, cx, card_y, card_w, 0.07, cl["accent"], radius=0.05)

        # 클러스터 배지
        badge_w = 0.6
        badge_h = 0.35
        add_rect(s, cx + 0.15, card_y + 0.25, badge_w, badge_h,
                 cl["accent"], radius=0.06)
        add_text(s, cx + 0.15, card_y + 0.25, badge_w, badge_h,
                 cl["id"], size=14, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 클러스터명
        add_text(s, cx + 0.85, card_y + 0.25, card_w - 1.0, 0.35,
                 cl["name"], size=16, color=CHARCOAL, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

        # RevPAR 히어로 숫자
        add_text(s, cx + 0.1, card_y + 0.8, card_w - 0.2, 0.55,
                 cl["revpar"], size=30, color=cl["accent"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 대표 자치구
        add_rect(s, cx + 0.15, card_y + 1.45, card_w - 0.3, 0.35,
                 cl["bg"], radius=0.04)
        add_text(s, cx + 0.15, card_y + 1.47, card_w - 0.3, 0.3,
                 cl["districts"], size=12, color=CHARCOAL,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 구분선
        add_rect(s, cx + 0.3, card_y + 2.0, card_w - 0.6, 0.015, GRAY_200)

        # 전략 키워드
        add_text(s, cx + 0.15, card_y + 2.15, card_w - 0.3, 0.8,
                 cl["strategy"], size=14, color=CHARCOAL,
                 align=PP_ALIGN.CENTER)

        # Dormant %
        dormant_y = card_y + 3.15
        add_rect(s, cx + 0.15, dormant_y, card_w - 0.3, 0.35,
                 GRAY_BG, radius=0.04)
        add_text(s, cx + 0.15, dormant_y + 0.02, card_w - 0.3, 0.3,
                 cl["dormant"], size=12, color=GRAY_500,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 하단 연결 메시지 ──
    footer_y = card_y + card_h + 0.25
    add_rect(s, 0.7, footer_y, 11.8, 0.5, GRAY_BG, radius=0.04)
    add_rect(s, 0.7, footer_y, 0.04, 0.5, TEAL)
    add_text(s, 1.0, footer_y + 0.05, 11.2, 0.4,
             "이 군집 결과가 대시보드의 '자치구 벤치마킹' 기능으로 이어집니다",
             size=15, color=CHARCOAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, 23)


# ═══════════════════════════════════════════
# 메인: 빌드 & 저장
# ═══════════════════════════════════════════
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_17(prs)
    build_slide_18(prs)
    build_slide_19(prs)
    build_slide_20(prs)
    build_slide_21(prs)
    build_slide_22(prs)
    build_slide_23(prs)

    out = "presentation/slides_17_23.pptx"
    prs.save(out)
    print(f"✓ {out} 저장 완료 — {len(prs.slides)}장")

    # ── 검증 ──
    errors = []
    max_w = 13.333
    max_h = 7.5
    for idx, slide in enumerate(prs.slides):
        page = 17 + idx
        for shape in slide.shapes:
            r = shape.left / 914400 + shape.width / 914400
            b = shape.top / 914400 + shape.height / 914400
            if r > max_w + 0.05:
                errors.append(f"p{page}: 가로 오버플로우 ({r:.2f}\" > {max_w}\")")
            if b > max_h + 0.05:
                errors.append(f"p{page}: 세로 오버플로우 ({b:.2f}\" > {max_h}\")")

    if errors:
        print(f"\n⚠ 오버플로우 {len(errors)}건:")
        for e in errors:
            print(f"  · {e}")
    else:
        print("✓ 오버플로우 없음 — 모든 shape이 13.333\" × 7.5\" 내")

    print(f"✓ 슬라이드 수: {len(prs.slides)}장")
