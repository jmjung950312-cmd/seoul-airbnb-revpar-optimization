"""
신규 슬라이드 3장 — 흐름 개선용 전환 슬라이드 (v2 디자인 리뉴얼)

디자인 개선 포인트:
  ✗ 이모지 아이콘 → ✓ 컬러 원형 배지 + 텍스트 심볼
  ✗ 연한 배경색 → ✓ GRAY_BG + 강한 accent 바 + 그림자 효과
  ✗ 세로 플로우 → ✓ 가로 플로우 (기존 PDF 패턴 통일)
  ✗ 약한 시각 계층 → ✓ 히어로 숫자 + 카드 + 바 3단 구조
  ✗ 심플 전환 슬라이드 → ✓ 정식 섹션 디바이더 패턴 (Layout D)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 디자인 토큰 (build_slides_28_34.py 동일) ──
CORAL       = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL    = RGBColor(0x48, 0x48, 0x48)
GRAY_BG     = RGBColor(0xF7, 0xF7, 0xF7)
GRAY_400    = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500    = RGBColor(0x76, 0x76, 0x76)
GRAY_200    = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_100    = RGBColor(0xF5, 0xF5, 0xF5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT  = RGBColor(0xE0, 0xF7, 0xF5)
ORANGE      = RGBColor(0xFC, 0x64, 0x2D)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xED)
BLUE        = RGBColor(0x4A, 0x90, 0xD9)
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)

FONT = "Pretendard"
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# 이메일 스크린샷 경로
EMAIL_SCREENSHOT = Path(__file__).parent.parent / "slide_preview" / "dashboard_free_3x.png"


# ═══════════════════════════════════════════
# 공통 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return tb


def add_rich_text(slide, left, top, width, height, runs,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 14))
        run.font.color.rgb = r.get("color", CHARCOAL)
        run.font.bold = r.get("bold", False)
        run.font.name = r.get("font", FONT)
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, diameter, fill_color):
    """컬러 원형 배지"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_header(slide, title):
    add_text(slide, 0.7, 0.35, 10.5, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 2.5, 0.04, CORAL)


def add_insight_box(slide, runs_list, left=0.7, top=1.15, width=11.8, height=0.55):
    """인사이트 박스 (rich text 지원)"""
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_rich_text(slide, left + 0.25, top + 0.1, width - 0.5, height - 0.2,
                  runs_list, anchor=MSO_ANCHOR.MIDDLE)


def add_page_num(slide, num):
    add_text(slide, 12.5, 7.0, 0.5, 0.3, str(num),
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_shadow_card(slide, left, top, width, height, fill_color=GRAY_BG,
                    border_color=None, radius=0.04, shadow_offset=0.04):
    """그림자 효과가 있는 카드"""
    # 그림자
    add_rect(slide, left + shadow_offset, top + shadow_offset,
             width, height, GRAY_200, radius=radius)
    # 카드 본체
    return add_rect(slide, left, top, width, height, fill_color,
                    border_color=border_color, radius=radius)


# ═══════════════════════════════════════════
# NEW A: 자동화 개요 — 3가지 자동화 시스템
# 기존 PDF p24(비즈니스 모델) 패턴 참조: 큰 카드 3열 + 화살표
# ═══════════════════════════════════════════
def build_slide_new_a(prs, page_num=29):
    s = new_slide(prs)
    add_header(s, "분석이 끝이 아니다 — 자동화로 연결한다")
    add_insight_box(s, [
        {"text": "인사이트 :  ", "size": 15, "color": CORAL, "bold": True},
        {"text": "'분석 → 대시보드 → 자동화'까지 이어지는 ", "size": 15},
        {"text": "End-to-End 체계", "size": 15, "bold": True},
        {"text": "를 구축했습니다", "size": 15},
    ])

    # ═══ 3개 카드 (비즈니스 모델 p24 패턴: 큰 카드 + 좌측 accent 바) ═══
    cards = [
        {
            "icon_letter": "M",      # Mail
            "icon_color": TEAL,
            "title": "이메일 진단 리포트",
            "price": "무료",
            "price_color": TEAL,
            "features": [
                "숙소 사전 진단",
                "시장 위치 분석",
                "개선 포인트 제안",
            ],
            "bg": WHITE,
            "accent": False,
        },
        {
            "icon_letter": "R",      # Risk
            "icon_color": CORAL,
            "title": "위험 감지 자동화",
            "price": "무료",
            "price_color": CORAL,
            "features": [
                "5가지 규칙 엔진",
                "통계 보조 검사",
                "이메일 알림 발송",
            ],
            "bg": WHITE,
            "accent": False,
        },
        {
            "icon_letter": "AI",     # AI Agent
            "icon_color": PURPLE,
            "title": "AI 에이전트 분석",
            "price": "Enterprise",
            "price_color": WHITE,
            "features": [
                "6개 에이전트 자동 호출",
                "가설 검증 자동화",
                "자동 리포팅",
            ],
            "bg": CORAL,
            "accent": True,
        },
    ]

    card_w = 3.6
    card_gap = 0.3
    card_start = 0.85
    card_y = 1.95
    card_h = 4.9

    for i, c in enumerate(cards):
        cx = card_start + i * (card_w + card_gap)
        is_accent = c["accent"]
        bg = c["bg"]
        text_color = WHITE if is_accent else CHARCOAL
        sub_color = WHITE if is_accent else GRAY_500
        line_color = WHITE if is_accent else GRAY_200

        # 카드 배경 (그림자)
        if not is_accent:
            add_rect(s, cx + 0.03, card_y + 0.03,
                     card_w, card_h,
                     GRAY_200, radius=0.06)
        add_rect(s, cx, card_y, card_w, card_h, bg,
                 border_color=GRAY_200 if not is_accent else None,
                 radius=0.06)

        # 아이콘 원형 배지 (상단 센터)
        icon_d = 0.7
        icon_x = cx + (card_w - icon_d) / 2
        icon_y = card_y + 0.5
        add_circle(s, icon_x, icon_y, icon_d, c["icon_color"])
        add_text(s, icon_x, icon_y + 0.05, icon_d, icon_d - 0.1,
                 c["icon_letter"], size=22, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 타이틀
        add_text(s, cx, card_y + 1.4, card_w, 0.4,
                 c["title"], size=22, color=text_color, bold=True,
                 align=PP_ALIGN.CENTER)

        # 가격
        price_bg = c["price_color"] if not is_accent else WHITE
        price_text = c["price_color"] if not is_accent else CORAL
        badge_w = 1.4
        add_rect(s, cx + (card_w - badge_w) / 2, card_y + 1.95,
                 badge_w, 0.35, price_bg, radius=0.5)
        add_text(s, cx + (card_w - badge_w) / 2, card_y + 1.97,
                 badge_w, 0.35,
                 c["price"], size=14, color=price_text if is_accent else WHITE,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 구분선
        add_rect(s, cx + 0.5, card_y + 2.55, card_w - 1.0, 0.015, line_color)

        # 기능 리스트
        for j, feat in enumerate(c["features"]):
            fy = card_y + 2.85 + j * 0.55
            # 불릿
            bullet_color = WHITE if is_accent else c["icon_color"]
            dot_size = 0.08
            add_rect(s, cx + 0.5, fy + 0.07, dot_size, dot_size,
                     bullet_color, radius=1.0)
            add_text(s, cx + 0.75, fy, card_w - 1.2, 0.35,
                     feat, size=15, color=text_color)

        # 화살표 (마지막 제외)
        if i < 2:
            ax = cx + card_w + 0.02
            add_text(s, ax, card_y + card_h / 2 - 0.2, 0.25, 0.4,
                     "→", size=24, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 하단 CTA 바 ──
    add_rect(s, 0.7, 7.05, 11.8, 0.03, GRAY_200)

    add_page_num(s, page_num)
    return s


# ═══════════════════════════════════════════
# NEW B: 무료 자동화 — 호스트 이메일 진단
# 기존 PDF p25 패턴 참조: 좌측 내용 + 우측 스크린샷
# ═══════════════════════════════════════════
def build_slide_new_b(prs, page_num=30):
    s = new_slide(prs)
    add_header(s, "무료: 숙소 진단 이메일 — 자동 발송")
    add_insight_box(s, [
        {"text": "인사이트 :  ", "size": 15, "color": CORAL, "bold": True},
        {"text": "호스트가 숙소 정보만 입력하면 AI가 자동으로 시장 분석 → ", "size": 15},
        {"text": "맞춤형 진단 리포트", "size": 15, "bold": True},
        {"text": "를 이메일로 발송합니다", "size": 15},
    ])

    # ═══ 좌측: 3단계 가로 플로우 (기존 PDF p9 전처리 패턴) ═══
    LX = 0.7
    LW = 5.5

    add_text(s, LX, 1.95, 4, 0.3,
             "자동 발송 플로우", size=18, color=CHARCOAL, bold=True)

    # 3단계 가로 플로우
    flow_steps = [
        ("1", "정보 입력", "자치구 · 유형\n가격 입력", TEAL),
        ("2", "AI 분석", "시장 위치 ·\n수익 예측", CORAL),
        ("3", "이메일", "HTML 리포트\n자동 발송", BLUE),
    ]

    step_w = 1.55
    step_h = 1.3
    step_gap = 0.15
    step_y = 2.4

    for i, (num, title, desc, color) in enumerate(flow_steps):
        sx = LX + i * (step_w + step_gap + 0.15)  # 화살표 공간

        # 카드
        add_rect(s, sx, step_y, step_w, step_h, GRAY_BG,
                 border_color=GRAY_200, radius=0.04)
        # 상단 accent 바
        add_rect(s, sx, step_y, step_w, 0.04, color)

        # 번호 원형
        circle_d = 0.4
        circle_x = sx + (step_w - circle_d) / 2
        add_circle(s, circle_x, step_y + 0.2, circle_d, color)
        add_text(s, circle_x, step_y + 0.22, circle_d, circle_d - 0.04,
                 num, size=16, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 타이틀
        add_text(s, sx + 0.1, step_y + 0.7, step_w - 0.2, 0.2,
                 title, size=13, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)

        # 설명
        add_text(s, sx + 0.1, step_y + 0.95, step_w - 0.2, 0.35,
                 desc, size=10, color=GRAY_500, align=PP_ALIGN.CENTER)

        # 화살표
        if i < 2:
            ax = sx + step_w + 0.02
            add_text(s, ax, step_y + step_h / 2 - 0.15, 0.12, 0.3,
                     "→", size=18, color=GRAY_400, align=PP_ALIGN.CENTER)

    # ── 좌측: 이메일 포함 내용 (리스트) ──
    list_y = 3.95
    add_text(s, LX, list_y, 4, 0.3,
             "이메일에 포함되는 내용", size=18, color=CHARCOAL, bold=True)

    email_items = [
        ("시장 위치", "서울 전체 / 자치구 내 수익 백분위 순위", CORAL),
        ("수익 진단", "구 평균 대비 · 최근 3개월 추세 · 유형별 비교", TEAL),
        ("경쟁력 진단", "같은 유형 숙소 대비 포지셔닝 (4분면)", BLUE),
        ("개선 포인트", "사진 · 요금 · 최소숙박일 등 구체적 행동 제안", PURPLE),
    ]

    for j, (title, desc, color) in enumerate(email_items):
        iy = list_y + 0.4 + j * 0.7

        # 카드 배경
        add_rect(s, LX, iy, LW, 0.58, GRAY_BG, radius=0.04)
        # 좌측 accent 바
        add_rect(s, LX, iy + 0.06, 0.04, 0.46, color)

        # 아이콘 배지
        badge_w = len(title) * 0.17 + 0.2
        add_rect(s, LX + 0.2, iy + 0.14, badge_w, 0.28, color, radius=0.08)
        add_text(s, LX + 0.2, iy + 0.15, badge_w, 0.28,
                 title, size=11, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 설명 텍스트
        add_text(s, LX + 0.3 + badge_w + 0.1, iy + 0.15, LW - badge_w - 0.8, 0.3,
                 desc, size=12, color=GRAY_500, anchor=MSO_ANCHOR.MIDDLE)

    # ═══ 우측: 이메일 미리보기 (브라우저 프레임) ═══
    RX = 6.5
    RW = 5.9
    RY = 1.95

    add_text(s, RX, RY, 4, 0.3,
             "이메일 미리보기", size=18, color=CHARCOAL, bold=True)

    frame_y = RY + 0.4
    frame_h = 4.85
    BAR_H = 0.32

    # 그림자
    add_rect(s, RX + 0.04, frame_y + 0.04, RW, frame_h, GRAY_200, radius=0.04)
    # 프레임 배경
    add_rect(s, RX, frame_y, RW, frame_h, WHITE,
             border_color=GRAY_200, radius=0.04)

    # 브라우저 바
    add_rect(s, RX, frame_y, RW, BAR_H,
             RGBColor(0x3C, 0x3C, 0x3C), radius=0.04)
    add_rect(s, RX, frame_y + BAR_H - 0.06, RW, 0.06,
             RGBColor(0x3C, 0x3C, 0x3C))

    # 트래픽 라이트
    for bx, bc in [(RX + 0.15, RGBColor(0xFF, 0x5F, 0x57)),
                   (RX + 0.3, RGBColor(0xFF, 0xBD, 0x2E)),
                   (RX + 0.45, RGBColor(0x28, 0xCA, 0x41))]:
        add_circle(s, bx, frame_y + 0.1, 0.1, bc)

    # URL 바
    add_rect(s, RX + 0.7, frame_y + 0.06, 3.2, 0.2,
             RGBColor(0x55, 0x55, 0x55), radius=0.5)
    add_text(s, RX + 0.85, frame_y + 0.065, 3.0, 0.2,
             "https://mail.google.com", size=8, color=GRAY_400,
             anchor=MSO_ANCHOR.MIDDLE)

    # LIVE 배지 (기존 PDF p7/p26 패턴)
    add_rect(s, RX + RW - 0.7, frame_y + 0.06, 0.5, 0.2,
             TEAL, radius=0.5)
    add_text(s, RX + RW - 0.7, frame_y + 0.065, 0.5, 0.2,
             "● FREE", size=7, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 이미지 삽입
    img_x = RX + 0.05
    img_y = frame_y + BAR_H + 0.03
    img_w = RW - 0.1
    img_h = frame_h - BAR_H - 0.08

    if EMAIL_SCREENSHOT.exists():
        try:
            from PIL import Image
            with Image.open(str(EMAIL_SCREENSHOT)) as img:
                orig_w, orig_h = img.size
                img_ratio = orig_w / orig_h
                frame_ratio = img_w / img_h
                if img_ratio > frame_ratio:
                    final_w = img_w
                    final_h = img_w / img_ratio
                else:
                    final_h = img_h
                    final_w = img_h * img_ratio
                offset_x = (img_w - final_w) / 2
                s.shapes.add_picture(
                    str(EMAIL_SCREENSHOT),
                    Inches(img_x + offset_x), Inches(img_y),
                    Inches(final_w), Inches(final_h))
                print(f"  ✓ 이메일 스크린샷 삽입: {final_w:.1f}×{final_h:.1f}in")
        except Exception as e:
            print(f"  ⚠ 이미지 삽입 실패: {e}")
            _add_placeholder(s, img_x, img_y, img_w, img_h)
    else:
        print(f"  ⚠ 스크린샷 미발견: {EMAIL_SCREENSHOT}")
        _add_placeholder(s, img_x, img_y, img_w, img_h)

    add_page_num(s, page_num)
    return s


def _add_placeholder(slide, x, y, w, h):
    """이미지 없을 때 플레이스홀더"""
    add_rect(slide, x, y, w, h, GRAY_100, radius=0.02)
    items = [
        ("내 숙소, 서울에서\n어디쯤 있을까요?", CORAL),
        ("내 숙소 시장 위치", CHARCOAL),
        ("수익 진단 결과", CHARCOAL),
    ]
    add_text(slide, x + 0.3, y + 0.4, w - 0.6, 0.5,
             items[0][0], size=20, color=items[0][1],
             bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, x + 0.3, y + 1.0, w - 0.6, 0.3,
             "지금 바로 수익 진단 결과를 확인해 보세요.",
             size=12, color=GRAY_500)

    # 3개 미니 KPI
    kpi_data = [("상위 99.1%", "서울 전체에서"),
                ("237번째 수익", "마포구 내"),
                ("가성비 신흥형", "시장 유형")]
    kpi_w = (w - 0.8) / 3
    for j, (val, label) in enumerate(kpi_data):
        kx = x + 0.3 + j * (kpi_w + 0.05)
        add_rect(slide, kx, y + 1.5, kpi_w, 0.8, WHITE, radius=0.03)
        add_text(slide, kx + 0.1, y + 1.55, kpi_w - 0.2, 0.15,
                 label, size=8, color=GRAY_500, align=PP_ALIGN.CENTER)
        add_text(slide, kx + 0.1, y + 1.8, kpi_w - 0.2, 0.35,
                 val, size=13, color=CORAL, bold=True,
                 align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# NEW C: 대시보드 전환 — Layout D 섹션 디바이더 정식 패턴
# 기존 PDF p23 (SECTION 06) 패턴 100% 참조
# ═══════════════════════════════════════════
def build_slide_new_c(prs, page_num=32):
    s = new_slide(prs)

    # ═══ 좌측 coral 패널 (기존 섹션 디바이더 정확히 재현) ═══
    panel_w = 4.8
    add_rect(s, 0, 0, panel_w, 7.5, CORAL)

    # "DASHBOARD" 대형 텍스트 (기존 "SECTION" 패턴)
    add_text(s, 0.5, 1.8, panel_w - 1.0, 0.6,
             "DASHBOARD", size=32, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT, font="Pretendard")

    # 구분선
    add_rect(s, 0.5, 3.3, 0.6, 0.05, WHITE)

    # 한글 제목
    add_text(s, 0.5, 3.6, panel_w - 1.0, 0.5,
             "호스트 컨설팅\n대시보드", size=24, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT)

    # ═══ 우측 내용 (기존 섹션 디바이더 패턴) ═══
    RX = 5.5
    RW = 7.0

    # 목적
    add_circle(s, RX, 1.3, 0.35, TEAL)
    add_text(s, RX + 0.05, 1.33, 0.35, 0.3,
             "◎", size=16, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + 0.5, 1.3, 3, 0.35,
             "목적", size=20, color=CHARCOAL, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, RX + 0.5, 1.75, RW - 1.0, 0.5,
             "분석 결과를 호스트가 직접 탐색하고\n즉시 행동으로 연결하는 의사결정 도구",
             size=15, color=GRAY_500)

    # 핵심 기능 제목
    add_circle(s, RX, 2.65, 0.35, CORAL)
    add_text(s, RX + 0.05, 2.68, 0.35, 0.3,
             "✦", size=14, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + 0.5, 2.65, 3, 0.35,
             "핵심 기능", size=20, color=CHARCOAL, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    # 4개 기능 카드 (기존 섹션 디바이더의 리스트 아이템 패턴)
    features = [
        {
            "icon": "◉",
            "icon_bg": CORAL_LIGHT, "icon_color": CORAL,
            "title": "시장 위치 4분면",
            "desc": "내 숙소의 경쟁력을 한눈에 확인",
        },
        {
            "icon": "♥",
            "icon_bg": TEAL_LIGHT, "icon_color": TEAL,
            "title": "건강점수 게이지",
            "desc": "A~F 등급으로 숙소 상태 종합 진단",
        },
        {
            "icon": "▲",
            "icon_bg": BLUE_LIGHT, "icon_color": BLUE,
            "title": "수익 예측 시뮬레이터",
            "desc": "조건 변경 시 수익 변화를 즉시 확인",
        },
        {
            "icon": "★",
            "icon_bg": ORANGE_LIGHT, "icon_color": ORANGE,
            "title": "맞춤 행동 체크리스트",
            "desc": "우선순위별 개선 제안",
        },
    ]

    feat_y_start = 3.05
    feat_h = 0.65
    feat_gap = 0.12

    for j, f in enumerate(features):
        fy = feat_y_start + j * (feat_h + feat_gap)

        # 카드 배경
        add_rect(s, RX, fy, RW - 0.5, feat_h, GRAY_BG, radius=0.04)

        # 좌측 accent 바
        add_rect(s, RX, fy + 0.05, 0.04, feat_h - 0.1, f["icon_color"])

        # 아이콘 원형
        add_circle(s, RX + 0.2, fy + 0.1, 0.42, f["icon_bg"])
        add_text(s, RX + 0.2, fy + 0.12, 0.42, 0.38,
                 f["icon"], size=16, color=f["icon_color"],
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 타이틀
        add_text(s, RX + 0.8, fy + 0.08, RW - 1.8, 0.28,
                 f["title"], size=15, color=CHARCOAL, bold=True)

        # 설명
        add_text(s, RX + 0.8, fy + 0.35, RW - 1.8, 0.25,
                 f["desc"], size=12, color=GRAY_500)

    # ── 하단: 사용자 정보 ──
    bottom_y = feat_y_start + 4 * (feat_h + feat_gap) + 0.15

    add_circle(s, RX, bottom_y, 0.35, BLUE)
    add_text(s, RX + 0.05, bottom_y + 0.03, 0.35, 0.3,
             "◆", size=14, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + 0.5, bottom_y, 2, 0.35,
             "사용자", size=20, color=CHARCOAL, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + 1.6, bottom_y + 0.05, 3, 0.3,
             "호스트 · 운영자", size=14, color=GRAY_500,
             anchor=MSO_ANCHOR.MIDDLE)

    # Streamlit 배지
    add_rect(s, RX, bottom_y + 0.55, 3.0, 0.35, GRAY_BG, radius=0.5)
    add_rect(s, RX, bottom_y + 0.61, 0.04, 0.23, TEAL)
    add_text(s, RX + 0.15, bottom_y + 0.58, 2.7, 0.3,
             "Streamlit 기반 · 실시간 · 모바일 지원",
             size=11, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, page_num)
    return s


# ═══════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 55)
    print("전환 슬라이드 3장 빌드 (v2 디자인 리뉴얼)")
    print("=" * 55)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("\n[1/3] NEW A: 자동화 개요...")
    build_slide_new_a(prs, page_num=29)

    print("\n[2/3] NEW B: 무료 자동화 이메일...")
    build_slide_new_b(prs, page_num=30)

    print("\n[3/3] NEW C: 대시보드 전환...")
    build_slide_new_c(prs, page_num=32)

    out = Path(__file__).parent / "new_transition_slides.pptx"
    prs.save(str(out))
    print(f"\n✅ 저장 완료: {out} (3장)")
