"""
마지막 슬라이드 — 5가지 변형
인사이트: "단순한 인사이트를 넘어, 플랫폼 생태계의 선순환을 만드는 데이터 분석을 진행하였습니다"

Variant A: 선순환 사이클 (4카드 순환 플로우)
Variant B: 3-Value 임팩트 카드 (호스트·게스트·플랫폼)
Variant C: 내러티브 텍스트 + 인용문 스타일
Variant D: Before → After (좌우 대비)
Variant E: 히어로 메시지 + 3 핵심 성과 카드
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 디자인 토큰 ──
CORAL       = RGBColor(0xFF, 0x38, 0x5C)
CORAL_LIGHT = RGBColor(0xFF, 0xF0, 0xF3)
CHARCOAL    = RGBColor(0x48, 0x48, 0x48)
GRAY_BG     = RGBColor(0xF7, 0xF7, 0xF7)
GRAY_100    = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_200    = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_400    = RGBColor(0xB0, 0xB0, 0xB0)
GRAY_500    = RGBColor(0x76, 0x76, 0x76)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xA6, 0x99)
TEAL_LIGHT  = RGBColor(0xE0, 0xF7, 0xF5)
ORANGE      = RGBColor(0xFC, 0x64, 0x2D)
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xED)
BLUE        = RGBColor(0x4A, 0x90, 0xD9)
BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)
PURPLE_LIGHT = RGBColor(0xF3, 0xE8, 0xFA)

FONT = "Pretendard"
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


# ═══════════════════════════════════════════
# 공통 헬퍼
# ═══════════════════════════════════════════
def add_text(slide, left, top, width, height, text,
             size=14, color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, Inches(0))
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
    return tb


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if radius:
        shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title):
    add_text(slide, 0.7, 0.35, 10.5, 0.6, title, size=43, color=CHARCOAL, bold=True)
    add_rect(slide, 0.7, 0.95, 11.8, 0.03, CORAL)


def add_insight_box(slide, text, left=0.7, top=1.15, width=11.8, height=0.7):
    add_rect(slide, left, top, width, height, CORAL_LIGHT, radius=0.04)
    add_text(slide, left + 0.25, top + 0.12, width - 0.5, height - 0.2,
             text, size=20, color=CHARCOAL)


def add_page_num(slide, num):
    add_text(slide, 12.5, 7.0, 0.5, 0.3, str(num),
             size=14, color=GRAY_400, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_circle(slide, left, top, diameter, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_shadow_card(slide, left, top, width, height, fill_color=GRAY_BG,
                    border_color=None, radius=0.04, shadow_offset=0.04):
    add_rect(slide, left + shadow_offset, top + shadow_offset,
             width, height, GRAY_200, radius=radius)
    return add_rect(slide, left, top, width, height, fill_color,
                    border_color=border_color, radius=radius)


INSIGHT_TEXT = ("단순한 인사이트를 넘어, "
                "플랫폼 생태계의 선순환을 만드는 데이터 분석을 진행하였습니다")


# ═══════════════════════════════════════════
# Variant A: 선순환 사이클 (4카드 순환 플로우)
# ═══════════════════════════════════════════
def build_variant_a(prs):
    """가로 4-step 리니어 플로우 — 선순환 단계를 깔끔하게 표현"""
    s = new_slide(prs)
    add_header(s, "플랫폼 생태계 선순환")
    add_insight_box(s, INSIGHT_TEXT)

    # ── 4-step 가로 플로우 (좌→우, 마지막에서 처음으로 피드백) ──
    steps = [
        {
            "num": "1", "label": "호스트 최적화",
            "line1": "SHAP 기반 액션 추천",
            "line2": "건강점수 진단 → 개선",
            "color": CORAL, "bg": CORAL_LIGHT,
        },
        {
            "num": "2", "label": "게스트 경험 향상",
            "line1": "높은 평점 · 적정 가격",
            "line2": "최적 예약정책 → 만족도",
            "color": TEAL, "bg": TEAL_LIGHT,
        },
        {
            "num": "3", "label": "플랫폼 데이터 축적",
            "line1": "리뷰 증가 · 예약 패턴",
            "line2": "시장 인텔리전스 강화",
            "color": BLUE, "bg": BLUE_LIGHT,
        },
        {
            "num": "4", "label": "더 나은 인사이트",
            "line1": "모델 정밀도 향상",
            "line2": "맞춤 전략 고도화",
            "color": PURPLE, "bg": PURPLE_LIGHT,
        },
    ]

    card_w = 2.55
    card_h = 3.0
    card_gap = 0.35
    total_w = 4 * card_w + 3 * card_gap
    start_x = (13.333 - total_w) / 2
    card_y = 2.2

    for i, step in enumerate(steps):
        cx = start_x + i * (card_w + card_gap)

        # 카드 배경
        add_shadow_card(s, cx, card_y, card_w, card_h,
                        step["bg"], radius=0.05)

        # 상단 accent 바
        add_rect(s, cx, card_y, card_w, 0.06, step["color"], radius=0.05)

        # 넘버 배지
        badge_d = 0.5
        badge_x = cx + (card_w - badge_d) / 2
        add_circle(s, badge_x, card_y + 0.3, badge_d, step["color"])
        add_text(s, badge_x, card_y + 0.3, badge_d, badge_d,
                 step["num"], size=20, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 레이블
        add_text(s, cx + 0.15, card_y + 1.0, card_w - 0.3, 0.35,
                 step["label"], size=17, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)

        # 구분선
        add_rect(s, cx + 0.3, card_y + 1.5, card_w - 0.6, 0.015, GRAY_200)

        # 설명
        add_text(s, cx + 0.2, card_y + 1.7, card_w - 0.4, 0.4,
                 step["line1"], size=13, color=GRAY_500,
                 align=PP_ALIGN.CENTER)
        add_text(s, cx + 0.2, card_y + 2.1, card_w - 0.4, 0.4,
                 step["line2"], size=13, color=GRAY_500,
                 align=PP_ALIGN.CENTER)

        # 화살표 (마지막 카드 제외)
        if i < 3:
            ax = cx + card_w + 0.02
            add_text(s, ax, card_y + card_h / 2 - 0.2, 0.3, 0.4,
                     "→", size=28, color=GRAY_400,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 하단 피드백 루프 바 ──
    loop_y = card_y + card_h + 0.35
    add_rect(s, start_x, loop_y, total_w, 0.55, GRAY_BG, radius=0.04)
    add_rect(s, start_x, loop_y, 0.04, 0.55, CORAL)
    add_text(s, start_x + 0.25, loop_y + 0.05, total_w - 0.5, 0.45,
             "4 → 1  더 나은 인사이트가 다시 호스트 최적화로 — 자생적 개선 루프",
             size=15, color=CHARCOAL, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, "A")


# ═══════════════════════════════════════════
# Variant B: 3-Value 임팩트 카드 (호스트·게스트·플랫폼)
# ═══════════════════════════════════════════
def build_variant_b(prs):
    """3개 이해관계자별 가치 카드 — 생태계 3자 관점"""
    s = new_slide(prs)
    add_header(s, "플랫폼 생태계 선순환")
    add_insight_box(s, INSIGHT_TEXT)

    # ── 3-column 가치 카드 ──
    values = [
        {
            "icon": "H", "icon_bg": CORAL,
            "title": "호스트 가치",
            "hero": "RevPAR +85%",
            "hero_color": CORAL,
            "items": [
                "SHAP 기반 액션 우선순위",
                "건강점수로 취약점 진단",
                "자치구 벤치마킹 포지셔닝",
                "리스크 사전 감지 알림",
            ],
            "bottom": "수익 극대화의 과학적 근거",
            "accent": CORAL,
        },
        {
            "icon": "G", "icon_bg": TEAL,
            "title": "게스트 가치",
            "hero": "품질 표준화",
            "hero_color": TEAL,
            "items": [
                "사진 21~35장 가이드라인",
                "평점 4.85+ 유지 인센티브",
                "최소숙박 2~3박 최적화",
                "즉시예약 활성화 촉진",
            ],
            "bottom": "예측 가능한 숙박 경험",
            "accent": TEAL,
        },
        {
            "icon": "P", "icon_bg": BLUE,
            "title": "플랫폼 가치",
            "hero": "Dormant 54%→↓",
            "hero_color": BLUE,
            "items": [
                "비활성 리스팅 조기 식별",
                "시장 유형별 맞춤 정책",
                "공급 과잉 자치구 경보",
                "데이터 기반 정책 의사결정",
            ],
            "bottom": "건강한 마켓플레이스 생태계",
            "accent": BLUE,
        },
    ]

    card_w = 3.73
    card_h = 4.7
    card_gap = 0.3
    start_x = 0.7
    card_y = 2.1

    for i, v in enumerate(values):
        cx = start_x + i * (card_w + card_gap)

        # 카드 배경
        add_shadow_card(s, cx, card_y, card_w, card_h, WHITE,
                        border_color=GRAY_200, radius=0.05)

        # 상단 accent 바
        add_rect(s, cx, card_y, card_w, 0.06, v["accent"], radius=0.05)

        # 아이콘
        icon_d = 0.6
        icon_x = cx + (card_w - icon_d) / 2
        add_circle(s, icon_x, card_y + 0.25, icon_d, v["icon_bg"])
        add_text(s, icon_x, card_y + 0.25, icon_d, icon_d,
                 v["icon"], size=22, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 타이틀
        add_text(s, cx + 0.2, card_y + 1.0, card_w - 0.4, 0.35,
                 v["title"], size=18, color=CHARCOAL, bold=True,
                 align=PP_ALIGN.CENTER)

        # 히어로 숫자
        add_text(s, cx + 0.2, card_y + 1.35, card_w - 0.4, 0.4,
                 v["hero"], size=24, color=v["hero_color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 구분선
        add_rect(s, cx + 0.4, card_y + 1.85, card_w - 0.8, 0.015, GRAY_200)

        # 항목 리스트
        for j, item in enumerate(v["items"]):
            iy = card_y + 2.05 + j * 0.42
            add_circle(s, cx + 0.3, iy + 0.08, 0.12, v["accent"])
            add_text(s, cx + 0.55, iy, card_w - 0.8, 0.35,
                     item, size=13, color=CHARCOAL)

        # 하단 강조
        bottom_y = card_y + card_h - 0.55
        add_rect(s, cx + 0.15, bottom_y, card_w - 0.3, 0.4,
                 GRAY_BG, radius=0.04)
        add_text(s, cx + 0.15, bottom_y + 0.02, card_w - 0.3, 0.36,
                 v["bottom"], size=12, color=GRAY_500, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, "B")


# ═══════════════════════════════════════════
# Variant C: 내러티브 텍스트 + 인용문 스타일
# ═══════════════════════════════════════════
def build_variant_c(prs):
    """텍스트 중심 — 큰 인용문 + 핵심 문장 나열"""
    s = new_slide(prs)
    add_header(s, "플랫폼 생태계 선순환")
    add_insight_box(s, INSIGHT_TEXT)

    # ── 좌측 CORAL 인용 패널 ──
    panel_x = 0.7
    panel_y = 2.15
    panel_w = 5.5
    panel_h = 4.6

    add_rect(s, panel_x, panel_y, panel_w, panel_h, CORAL, radius=0.06)

    # 좌측 두꺼운 accent 바 (폰트 호환 문제 방지)
    add_rect(s, panel_x + 0.35, panel_y + 0.5, 0.06, 3.2, WHITE)

    # 인용문 3줄
    quote_lines = [
        "데이터는 숫자가 아니라",
        "호스트의 내일을",
        "바꾸는 도구입니다.",
    ]
    for qi, qline in enumerate(quote_lines):
        add_text(s, panel_x + 0.7, panel_y + 0.6 + qi * 0.95,
                 panel_w - 1.2, 0.7,
                 qline, size=28, color=WHITE, bold=True)

    # 구분선
    add_rect(s, panel_x + 0.7, panel_y + 3.5, 2.5, 0.025, WHITE)

    # 하단 부연
    add_text(s, panel_x + 0.7, panel_y + 3.7, panel_w - 1.2, 0.5,
             "서울 에어비앤비 RevPAR 최적화 프로젝트",
             size=14, color=RGBColor(0xFF, 0xCC, 0xD5))

    # ── 우측 핵심 문장 리스트 ──
    right_x = 6.6
    right_w = 5.9
    ry = 2.15

    statements = [
        {
            "num": "01",
            "title": "예측이 아닌, 행동을 설계합니다",
            "desc": "SHAP → 액션 전환 파이프라인으로\n"
                    "호스트가 '무엇을 해야 하는지' 직접 알려줍니다",
            "color": CORAL,
        },
        {
            "num": "02",
            "title": "개별 숙소와 시장, 두 눈으로 봅니다",
            "desc": "LightGBM(호스트 단위) + K-Means(자치구 단위)\n"
                    "미시와 거시를 동시에 분석합니다",
            "color": TEAL,
        },
        {
            "num": "03",
            "title": "위험을 예측하고, 기회를 제안합니다",
            "desc": "Phantom Revenue, Rate Outlier 등 5가지 리스크 규칙과\n"
                    "건강점수 A~F 등급으로 선제적 관리",
            "color": BLUE,
        },
        {
            "num": "04",
            "title": "분석이 끝이 아닌, 시작입니다",
            "desc": "Streamlit 대시보드 + 자동 이메일 알림으로\n"
                    "분석 결과가 호스트의 일상에 스며듭니다",
            "color": PURPLE,
        },
    ]

    for i, st in enumerate(statements):
        sy = ry + i * 1.15
        # 카드
        add_rect(s, right_x, sy, right_w, 1.0, GRAY_BG, radius=0.04)
        add_rect(s, right_x, sy + 0.08, 0.04, 0.84, st["color"])

        # 넘버
        add_text(s, right_x + 0.2, sy + 0.08, 0.5, 0.3,
                 st["num"], size=20, color=st["color"], bold=True)

        # 타이틀
        add_text(s, right_x + 0.7, sy + 0.08, right_w - 1.0, 0.3,
                 st["title"], size=15, color=CHARCOAL, bold=True)

        # 설명
        add_text(s, right_x + 0.7, sy + 0.42, right_w - 1.0, 0.5,
                 st["desc"], size=11, color=GRAY_500)

    add_page_num(s, "C")


# ═══════════════════════════════════════════
# Variant D: Before → After (좌우 대비)
# ═══════════════════════════════════════════
def build_variant_d(prs):
    """좌측 Before (회색) / 우측 After (컬러) 대비"""
    s = new_slide(prs)
    add_header(s, "플랫폼 생태계 선순환")
    add_insight_box(s, INSIGHT_TEXT)

    content_y = 2.15
    col_w = 5.85

    # ═══ 좌측: BEFORE ═══
    left_x = 0.7
    add_rect(s, left_x, content_y, col_w, 0.5, GRAY_200, radius=0.04)
    add_text(s, left_x + 0.2, content_y + 0.05, col_w - 0.4, 0.4,
             "BEFORE — 기존 방식", size=20, color=GRAY_500, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    befores = [
        {"text": "감에 의존한 가격 설정", "icon": "?"},
        {"text": "왜 예약이 안 되는지 모름", "icon": "?"},
        {"text": "비활성 리스팅 방치", "icon": "?"},
        {"text": "자치구 시장 상황 파악 불가", "icon": "?"},
        {"text": "리스크를 사후에 발견", "icon": "?"},
    ]

    for i, b in enumerate(befores):
        by = content_y + 0.7 + i * 0.72
        add_rect(s, left_x, by, col_w, 0.6, GRAY_100, radius=0.04)
        # X 표시
        add_circle(s, left_x + 0.15, by + 0.12, 0.36, GRAY_400)
        add_text(s, left_x + 0.15, by + 0.12, 0.36, 0.36,
                 "✕", size=16, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, left_x + 0.65, by + 0.08, col_w - 1.0, 0.44,
                 b["text"], size=15, color=GRAY_500,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ── 중앙 화살표 ──
    add_text(s, 6.3, 3.8, 0.6, 0.6,
             "→", size=36, color=CORAL, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ═══ 우측: AFTER ═══
    right_x = 6.95
    add_rect(s, right_x, content_y, col_w, 0.5, CORAL, radius=0.04)
    add_text(s, right_x + 0.2, content_y + 0.05, col_w - 0.4, 0.4,
             "AFTER — 데이터 기반", size=20, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

    afters = [
        {"text": "ADR 예측 모델로 적정 가격 제시", "color": CORAL, "bg": CORAL_LIGHT},
        {"text": "SHAP Top 5 액션으로 원인 진단", "color": TEAL, "bg": TEAL_LIGHT},
        {"text": "건강점수 F등급 조기 경보", "color": ORANGE, "bg": ORANGE_LIGHT},
        {"text": "K-Means 군집으로 시장 포지션 파악", "color": BLUE, "bg": BLUE_LIGHT},
        {"text": "5가지 규칙 엔진으로 리스크 사전 감지", "color": PURPLE, "bg": PURPLE_LIGHT},
    ]

    for i, a in enumerate(afters):
        ay = content_y + 0.7 + i * 0.72
        add_rect(s, right_x, ay, col_w, 0.6, a["bg"], radius=0.04)
        add_rect(s, right_x, ay + 0.06, 0.04, 0.48, a["color"])
        # 체크 표시
        add_circle(s, right_x + 0.15, ay + 0.12, 0.36, a["color"])
        add_text(s, right_x + 0.15, ay + 0.12, 0.36, 0.36,
                 "✓", size=16, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, right_x + 0.65, ay + 0.08, col_w - 1.0, 0.44,
                 a["text"], size=15, color=CHARCOAL,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, "D")


# ═══════════════════════════════════════════
# Variant E: 히어로 메시지 + 3 핵심 성과 카드
# ═══════════════════════════════════════════
def build_variant_e(prs):
    """큰 메시지 중심 + 하단 3개 성과 카드"""
    s = new_slide(prs)
    add_header(s, "플랫폼 생태계 선순환")
    add_insight_box(s, INSIGHT_TEXT)

    # ── 히어로 영역 (전체 너비) ──
    hero_y = 2.2
    hero_h = 2.5

    # 배경 카드
    add_rect(s, 0.7, hero_y, 11.8, hero_h, GRAY_BG, radius=0.06)

    # 큰 메시지 3줄
    lines = [
        {"text": "호스트의 성장이 → 게스트 만족으로,", "color": CORAL},
        {"text": "게스트 만족이 → 플랫폼 성장으로,", "color": TEAL},
        {"text": "플랫폼 성장이 → 더 나은 인사이트로.", "color": BLUE},
    ]

    for i, line in enumerate(lines):
        ly = hero_y + 0.3 + i * 0.7
        # 좌측 컬러 도트
        add_circle(s, 1.2, ly + 0.08, 0.25, line["color"])
        # 텍스트
        add_text(s, 1.7, ly, 10.0, 0.5,
                 line["text"], size=24, color=CHARCOAL, bold=True)

    # 하단 보충
    add_text(s, 1.7, hero_y + hero_h - 0.5, 10.0, 0.35,
             "이 프로젝트는 단일 분석이 아닌, 자생적 개선 루프를 설계합니다",
             size=14, color=GRAY_500)

    # ── 하단 3 성과 카드 ──
    kpi_y = 5.0
    kpi_w = 3.73
    kpi_h = 1.85
    kpi_gap = 0.3
    kpi_start = 0.7

    kpis = [
        {
            "value": "R² = 0.85",
            "label": "호스트 통제 변수만으로\nRevPAR 85% 설명",
            "sub_label": "LightGBM 5-Fold CV",
            "color": CORAL, "bar": CORAL,
        },
        {
            "value": "k = 4 군집",
            "label": "25개 자치구를\n4대 시장 유형으로 분류",
            "sub_label": "K-Means + Elbow Method",
            "color": TEAL, "bar": TEAL,
        },
        {
            "value": "5 규칙 엔진",
            "label": "Phantom Revenue 등\n리스크 자동 탐지",
            "sub_label": "IQR + Z-score 기반",
            "color": BLUE, "bar": BLUE,
        },
    ]

    for i, k in enumerate(kpis):
        kx = kpi_start + i * (kpi_w + kpi_gap)
        add_shadow_card(s, kx, kpi_y, kpi_w, kpi_h, WHITE,
                        border_color=GRAY_200, radius=0.05)
        # 상단 accent 바
        add_rect(s, kx, kpi_y, kpi_w, 0.05, k["bar"], radius=0.05)

        # 값
        add_text(s, kx + 0.2, kpi_y + 0.2, kpi_w - 0.4, 0.5,
                 k["value"], size=28, color=k["color"], bold=True,
                 align=PP_ALIGN.CENTER)

        # 설명
        add_text(s, kx + 0.2, kpi_y + 0.75, kpi_w - 0.4, 0.55,
                 k["label"], size=13, color=CHARCOAL,
                 align=PP_ALIGN.CENTER)

        # 방법론 태그
        tag_y = kpi_y + kpi_h - 0.45
        add_rect(s, kx + 0.3, tag_y, kpi_w - 0.6, 0.3,
                 GRAY_BG, radius=0.04)
        add_text(s, kx + 0.3, tag_y + 0.02, kpi_w - 0.6, 0.26,
                 k["sub_label"], size=11, color=GRAY_500,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(s, "E")


# ═══════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_variant_a(prs)  # 선순환 사이클
    build_variant_b(prs)  # 3-Value 임팩트
    build_variant_c(prs)  # 내러티브 + 인용문
    build_variant_d(prs)  # Before → After
    build_variant_e(prs)  # 히어로 + 성과 카드

    out = "presentation/slide_final_variants.pptx"
    prs.save(out)
    print(f"✓ {out} 저장 완료 — {len(prs.slides)}장 변형")

    # ── 검증 ──
    labels = ["A 선순환", "B 3-Value", "C 내러티브", "D Before/After", "E 히어로"]
    errors = []
    max_w, max_h = 13.333, 7.5
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            r = shape.left / 914400 + shape.width / 914400
            b = shape.top / 914400 + shape.height / 914400
            if r > max_w + 0.05:
                errors.append(f"Variant {labels[idx]}: 가로 오버플로우 ({r:.2f}\")")
            if b > max_h + 0.05:
                errors.append(f"Variant {labels[idx]}: 세로 오버플로우 ({b:.2f}\")")

    if errors:
        print(f"\n⚠ 오버플로우 {len(errors)}건:")
        for e in errors:
            print(f"  · {e}")
    else:
        print("✓ 오버플로우 없음")

    for idx, label in enumerate(labels):
        cnt = len(prs.slides[idx].shapes)
        print(f"  Variant {label}: {cnt} shapes")
