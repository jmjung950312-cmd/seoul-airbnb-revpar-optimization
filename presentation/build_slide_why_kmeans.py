"""
슬라이드: "왜 K-Means 자치구 분류인가?" — 단일 슬라이드 빌드
K-Means 군집 분석의 3가지 필요성을 시각적으로 전달
"""
import os
import json
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 경로 설정 ─────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(BASE_DIR)
DATA_PATH = os.path.join(PROJECT_DIR, "data/processed/district_clustered.csv")
DS_PATH = os.path.join(BASE_DIR, "design_system.json")
VIZ_DIR = os.path.join(BASE_DIR, "figures")
OUTPUT_PATH = os.path.join(BASE_DIR, "slide_why_kmeans.pptx")
os.makedirs(VIZ_DIR, exist_ok=True)

# ─── 디자인 시스템 로드 ──────────────────────────
with open(DS_PATH) as f:
    ds = json.load(f)

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

COLORS = {k: rgb(v) for k, v in ds['colors'].items()}
FONT = ds['meta']['font_primary']
FONT_FB = ds['meta']['font_fallback']

# ─── STEP 3-2: 차트 이미지 생성 ──────────────────
print("📊 차트 생성 중...")
df = pd.read_csv(DATA_PATH)

plt.rcParams['font.family'] = 'AppleGothic'
plt.rcParams['axes.unicode_minus'] = False

CLUSTER_COLORS = {
    "핫플 수익형": "#FF5A5F",
    "프리미엄 비즈니스": "#00A699",
    "로컬 주거형": "#767676",
    "가성비 신흥형": "#484848"
}
CLUSTER_ORDER = ["핫플 수익형", "프리미엄 비즈니스", "로컬 주거형", "가성비 신흥형"]

# --- 차트: 공급량 vs RevPAR 산점도 (클러스터 색상) ---
fig, ax = plt.subplots(figsize=(12, 7), dpi=150)
fig.patch.set_facecolor('white')
ax.set_facecolor('white')

for name in CLUSTER_ORDER:
    group = df[df['cluster_name'] == name]
    color = CLUSTER_COLORS[name]
    ax.scatter(
        group['total_listings'], group['median_revpar_ao'],
        s=200, label=name, color=color, alpha=0.85,
        edgecolors='white', linewidths=1.5, zorder=5
    )
    for _, row in group.iterrows():
        # 자치구명 한글 매핑 (주요 자치구만)
        district_kr = {
            'Mapo-gu': '마포', 'Gangnam-gu': '강남', 'Jongno-gu': '종로',
            'Jung-gu': '중구', 'Yongsan-gu': '용산', 'Seocho-gu': '서초',
            'Gwanak-gu': '관악', 'Seodaemun-gu': '서대문', 'Songpa-gu': '송파',
            'Gwangjin-gu': '광진', 'Seongdong-gu': '성동',
            'Dongdaemun-gu': '동대문', 'Dongjak-gu': '동작',
            'Geumcheon-gu': '금천', 'Guro-gu': '구로',
            'Yeongdeungpo-gu': '영등포', 'Gangseo-gu': '강서',
            'Gangdong-gu': '강동', 'Gangbuk-gu': '강북',
            'Eunpyeong-gu': '은평', 'Seongbuk-gu': '성북',
            'Yangcheon-gu': '양천', 'Nowon-gu': '노원',
            'Dobong-gu': '도봉', 'Jungnang-gu': '중랑',
        }
        label = district_kr.get(row['district'], row['district'])
        ax.annotate(
            label,
            (row['total_listings'], row['median_revpar_ao']),
            fontsize=8, ha='center', va='bottom',
            textcoords="offset points", xytext=(0, 8),
            color='#484848', fontweight='medium'
        )

# 스타일링
ax.set_xlabel('총 리스팅 수 (공급량)', fontsize=12, color='#484848', labelpad=10)
ax.set_ylabel('RevPAR 중위값 (₩)', fontsize=12, color='#484848', labelpad=10)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['left'].set_color('#E8E8E8')
ax.spines['bottom'].set_color('#E8E8E8')
ax.tick_params(colors='#767676', labelsize=10)
ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'₩{x:,.0f}'))
ax.grid(axis='y', alpha=0.2, color='#E8E8E8')

# 범례 — 마포구(우상단)와 겹치지 않도록 좌하단 배치
legend = ax.legend(
    fontsize=10, loc='lower right',
    frameon=True, fancybox=True, shadow=False,
    edgecolor='#E8E8E8', facecolor='white'
)

plt.tight_layout()
chart_path = os.path.join(VIZ_DIR, "ppt_why_kmeans_scatter.png")
plt.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print(f"  ✅ 차트 저장: {chart_path}")

# ─── STEP 4: PPTX 빌드 ──────────────────────────
print("\n🔨 슬라이드 빌드 중...")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

def add_textbox(slide, text, x, y, w, h, *,
                size=14, bold=False, color=None, font=FONT,
                align=PP_ALIGN.LEFT, wrap=True):
    """텍스트박스 헬퍼"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, *, fg=None, border=None):
    """색상 직사각형 헬퍼"""
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fg:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fg
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

# --- 제목 ---
add_textbox(slide, "왜 K-Means 자치구 분류인가?", 0.8, 0.5, 11.7, 0.8,
            size=28, bold=True, color=COLORS['dominant'])

# 제목 아래 악센트 언더라인
add_rect(slide, 0.8, 1.2, 2.0, 0.04, fg=COLORS['accent'])

# --- 좌측: 3가지 이유 카드 (세로 배치) ---
card_x, card_w = 0.8, 5.6
reasons = [
    {
        "num": "1",
        "title": "자치구마다 시장 구조가 다르다",
        "body": "마포구(₩63,366) vs 금천구(₩26,423) — RevPAR 2.4배 차이.\n같은 서울이라도 공급량·수익·비활성 비율이 근본적으로 다름.",
        "color": "accent"
    },
    {
        "num": "2",
        "title": "유사한 시장끼리 자동 그룹화",
        "body": "4개 피처(RevPAR·공급량·비활성비율·슈퍼호스트율) 기반\n객관적 거리 측정으로 25개 자치구를 4개 유형으로 분류.",
        "color": "accent_secondary"
    },
    {
        "num": "3",
        "title": "시장 유형별 차별화된 전략 도출",
        "body": "25개 개별 전략이 아닌 4개 유형 전략으로 실행 가능성↑\n호스트에게 '내 동네는 어떤 시장인가?' 답변 제공.",
        "color": "orange"
    }
]

card_y_start = 1.6
card_h = 1.55
card_gap = 0.2

for i, r in enumerate(reasons):
    cy = card_y_start + i * (card_h + card_gap)
    color_key = r['color']
    accent_rgb = COLORS[color_key]

    # 카드 배경
    add_rect(slide, card_x, cy, card_w, card_h, fg=COLORS['card_bg'])
    # 좌측 컬러 바
    add_rect(slide, card_x, cy, 0.06, card_h, fg=accent_rgb)

    # 번호 원형
    oval = slide.shapes.add_shape(
        9,  # OVAL
        Inches(card_x + 0.2), Inches(cy + 0.18),
        Inches(0.42), Inches(0.42)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = accent_rgb
    oval.line.fill.background()
    # 원 안의 숫자
    tf = oval.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = r['num']
    run.font.name = FONT
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # 카드 제목
    add_textbox(slide, r['title'], card_x + 0.75, cy + 0.15, card_w - 1.0, 0.4,
                size=16, bold=True, color=COLORS['dominant'])
    # 카드 본문
    add_textbox(slide, r['body'], card_x + 0.75, cy + 0.6, card_w - 1.0, 0.85,
                size=12, color=COLORS['supporting'])

# --- 우측: 클러스터 산점도 차트 ---
chart_x = 6.7
chart_y = 1.6
chart_w = 5.83
chart_img_path = os.path.join(VIZ_DIR, "ppt_why_kmeans_scatter.png")

if os.path.exists(chart_img_path):
    slide.shapes.add_picture(
        chart_img_path,
        Inches(chart_x), Inches(chart_y),
        width=Inches(chart_w)
    )
else:
    add_rect(slide, chart_x, chart_y, chart_w, 4.5, fg=COLORS['card_bg'], border=COLORS['border'])
    add_textbox(slide, "[차트 이미지 없음]", chart_x, chart_y + 2.0, chart_w, 0.4,
                size=12, color=COLORS['supporting'], align=PP_ALIGN.CENTER)

# --- 하단 인사이트 박스 ---
insight_y = 6.65
insight_w = 11.733
add_rect(slide, 0.8, insight_y, insight_w, 0.5, fg=COLORS['accent_light'])
add_rect(slide, 0.8, insight_y, 0.05, 0.5, fg=COLORS['accent'])

# 인사이트 텍스트 (rich text)
tb = slide.shapes.add_textbox(
    Inches(1.0), Inches(insight_y + 0.08),
    Inches(insight_w - 0.4), Inches(0.35)
)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]

parts = [
    ("● ", False, 'dominant'),
    ("핵심: ", True, 'accent'),
    ("개별 자치구 25개 분석 → 파편화 / 서울 전체 1개 분석 → 과잉 일반화. ", False, 'dominant'),
    ("4개 유형 분류가 실행 가능한 최적 단위.", True, 'accent'),
]

for text, bold, color_key in parts:
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color_key]

# --- 출처 ---
add_textbox(slide, "데이터: 서울 에어비앤비 32,061개 리스팅 · Active+Operating 14,399개 · TTM 12개월 기준",
            0.8, 7.15, 8.0, 0.25,
            size=9, color=COLORS['supporting'])

# ─── STEP 5: 검증 ────────────────────────────────
print("\n🔍 검증 중...")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE = Inches(0.15)
warnings = []

for shape in slide.shapes:
    x, y, w, h = shape.left, shape.top, shape.width, shape.height
    if x < 0:
        warnings.append(f"  좌측 밖: {shape.name}")
    if y < 0:
        warnings.append(f"  상단 밖: {shape.name}")
    if x + w > SLIDE_W + SAFE:
        warnings.append(f"  우측 초과: {shape.name} (x+w={Emu(x+w).inches:.2f}\")")
    if y + h > SLIDE_H + SAFE:
        warnings.append(f"  하단 초과: {shape.name} (y+h={Emu(y+h).inches:.2f}\")")

if warnings:
    print("  ⚠️ WARNING:")
    for w in warnings:
        print(w)
else:
    print("  ✅ 안전 영역: PASS")

# 텍스트 요소 수 확인
text_count = sum(1 for s in slide.shapes if s.has_text_frame)
print(f"  📝 텍스트 요소: {text_count}개")

# 시각 요소 확인
has_visual = any(
    s.shape_type == 13  # picture
    for s in slide.shapes
)
print(f"  🖼️ 이미지 포함: {'✅' if has_visual else '❌'}")

# ─── 저장 ────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"\n✅ 저장 완료: {OUTPUT_PATH}")
print(f"   슬라이드: 1장")
print(f"   요소: {len(slide.shapes)}개")
